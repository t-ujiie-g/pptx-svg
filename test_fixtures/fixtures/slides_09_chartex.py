"""Slides 89-94: ChartEx (cx: namespace) — placeholder slides patched by postprocess.

Mechanically carved from the legacy gen_test_features.py. Runs its
slide-building code at import time against the shared '_ctx.prs'.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
import base64, io, os, re, struct, tempfile, zipfile, zlib

from ._ctx import prs, blank

# ── Slides 89-94: Office 2016+ ChartEx (cx: namespace) ──────────────────────
cx_chart_names = [
    "Waterfall Chart",
    "Treemap Chart",
    "Sunburst Chart",
    "Histogram Chart",
    "Box & Whisker Chart",
    "Funnel Chart",
]
cx_chart_slides = []
for i, name in enumerate(cx_chart_names):
    s = prs.slides.add_slide(blank)
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    t.text_frame.paragraphs[0].text = f"Slide {89 + i}: {name} (cx:chart)"
    t.text_frame.paragraphs[0].font.size = Pt(18)
    t.text_frame.paragraphs[0].font.bold = True
    cx_chart_slides.append(s)

# Save first, then patch the OMML into the slide XML
