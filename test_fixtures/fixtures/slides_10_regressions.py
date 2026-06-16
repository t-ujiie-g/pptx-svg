"""Slide 96: regression fixtures.

Slides for behaviors that don't fit neatly into a category module — typically
small reproductions for bugs that have been fixed, kept so regressions can be
detected automatically.

Currently covered:
  - empty custGeom (issue #39): a custGeom with no path commands must NOT
    render as a colored rectangle.
  - p:style/a:fontRef text color: a shape whose <p:style> includes
    <a:fontRef> should apply that color to text runs that have no explicit
    color, matching PowerPoint's behavior (accent-styled rectangles get white
    text via fontRef -> lt1).
"""
from pptx.util import Inches
from lxml import etree

from ._ctx import prs, blank

ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

# ── Slide 96: empty custGeom + fontRef text color ──────────────────────────
slide96 = prs.slides.add_slide(blank)

# Reference rect with <p:style>/<a:fontRef idx="minor"><a:schemeClr val="lt1"/></a:fontRef>.
# python-pptx's add_shape() emits exactly this style block, so the only thing
# we need to do is add a text run with NO explicit color. PowerPoint renders
# the text in white (via fontRef -> lt1 -> theme background light).
ref = slide96.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(2.5), Inches(0.6))
ref.text_frame.text = "Reference rect (prstGeom)"

# Empty custGeom (the regressed shape). Must NOT render as a colored block.
spTree95 = slide96.shapes._spTree
empty_custgeom_xml = f"""
<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:nvSpPr><p:cNvPr id="9501" name="EmptyCustGeom"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="3000000" y="1500000"/><a:ext cx="2000000" cy="2000000"/></a:xfrm>
    <a:custGeom>
      <a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
      <a:rect l="0" t="0" r="0" b="0"/>
      <a:pathLst><a:path w="508000" h="508000"/></a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="00B4D8"/></a:solidFill>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
</p:sp>
"""
spTree95.append(etree.fromstring(empty_custgeom_xml))

# Valid custGeom — used as a control so we can verify normal shapes still
# render after the empty-custGeom fix.
valid_custgeom_xml = f"""
<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:nvSpPr><p:cNvPr id="9502" name="ValidCustGeom"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="6000000" y="1500000"/><a:ext cx="2000000" cy="2000000"/></a:xfrm>
    <a:custGeom>
      <a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
      <a:rect l="0" t="0" r="0" b="0"/>
      <a:pathLst>
        <a:path w="100" h="100">
          <a:moveTo><a:pt x="50" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="100" y="100"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="100"/></a:lnTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="06A77D"/></a:solidFill>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
</p:sp>
"""
spTree95.append(etree.fromstring(valid_custgeom_xml))

# ── Slide 97: header/footer field placeholders (date / footer / slide number) ──
# Reproduces sample.pptx slide 20: a slide carrying date (<a:fld type="datetime1">),
# footer (ftr) and slide-number (<a:fld type="slidenum"> with a STALE cached "1")
# placeholders. The renderer must (a) show the ACTUAL slide number, not the cached
# "1", (b) fill the date field with the host-provided current date, and (c) keep
# the footer text.
slide97 = prs.slides.add_slide(blank)
spTree97 = slide97.shapes._spTree

footer_ph_xml = f"""
<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:nvSpPr><p:cNvPr id="9701" name="Footer Placeholder"/><p:cNvSpPr/>
    <p:nvPr><p:ph type="ftr" idx="11"/></p:nvPr></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="3124200" y="6356350"/><a:ext cx="2438400" cy="365125"/></a:xfrm></p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr/><a:t>moon-pptx footer</a:t></a:r></a:p></p:txBody>
</p:sp>
"""
date_ph_xml = f"""
<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:nvSpPr><p:cNvPr id="9702" name="Date Placeholder"/><p:cNvSpPr/>
    <p:nvPr><p:ph type="dt" idx="10"/></p:nvPr></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="457200" y="6356350"/><a:ext cx="2438400" cy="365125"/></a:xfrm></p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p>
    <a:fld id="{{1F3B9E54-0000-4000-8000-000000000099}}" type="datetime1"><a:rPr/></a:fld></a:p></p:txBody>
</p:sp>
"""
sldnum_ph_xml = f"""
<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:nvSpPr><p:cNvPr id="9703" name="Slide Number Placeholder"/><p:cNvSpPr/>
    <p:nvPr><p:ph type="sldNum" idx="12"/></p:nvPr></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="6553200" y="6356350"/><a:ext cx="2438400" cy="365125"/></a:xfrm></p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p>
    <a:fld id="{{1F3B9E54-0000-4000-8000-000000000099}}" type="slidenum"><a:rPr/><a:t>1</a:t></a:fld></a:p></p:txBody>
</p:sp>
"""
for xml in (footer_ph_xml, date_ph_xml, sldnum_ph_xml):
    spTree97.append(etree.fromstring(xml))

# ── Slide 98: underline styles + underline color (a:uFill) + double strike ────
# Underline color round-trip (R1): the underline value (u=) and a:uFill color are
# parsed, round-tripped (data-ooxml-underline-color) and re-serialized to <a:uFill>.
# (Decoration STYLE/COLOR are not drawn visually — browsers ignore
# text-decoration-style/-color on SVG text — but must survive export.)
slide98 = prs.slides.add_slide(blank)
spTree98 = slide98.shapes._spTree

# Each run: u="<variant>" with a red a:uFill underline color; one dblStrike run.
deco_runs = (
    '<a:r><a:rPr lang="en-US" sz="2000" u="dbl"><a:uFill><a:solidFill>'
    '<a:srgbClr val="FF0000"/></a:solidFill></a:uFill></a:rPr><a:t>double </a:t></a:r>'
    '<a:r><a:rPr lang="en-US" sz="2000" u="wavy"><a:uFill><a:solidFill>'
    '<a:srgbClr val="00B050"/></a:solidFill></a:uFill></a:rPr><a:t>wavy </a:t></a:r>'
    '<a:r><a:rPr lang="en-US" sz="2000" u="dotted"/><a:t>dotted </a:t></a:r>'
    '<a:r><a:rPr lang="en-US" sz="2000" strike="dblStrike"/><a:t>dblstrike</a:t></a:r>'
)
deco_sp_xml = f"""
<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">
  <p:nvSpPr><p:cNvPr id="9801" name="DecoText"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="8229600" cy="914400"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p>{deco_runs}</a:p></p:txBody>
</p:sp>
"""
spTree98.append(etree.fromstring(deco_sp_xml))
