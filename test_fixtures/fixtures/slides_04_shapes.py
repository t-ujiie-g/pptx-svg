"""Slides 34-40: groups, connectors, preset geometry, custom geometry, gears, text rects, connection points.

Mechanically carved from the legacy gen_test_features.py. Runs its
slide-building code at import time against the shared '_ctx.prs'.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
import base64, io, os, re, struct, tempfile, zipfile, zlib

from ._ctx import prs, blank

# ── Slide 34: Group shapes ────────────────────────────────────────────────
slide34 = prs.slides.add_slide(blank)

title34 = slide34.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title34.text_frame.paragraphs[0].text = "Slide 34: Group Shapes"
title34.text_frame.paragraphs[0].font.size = Pt(24)
title34.text_frame.paragraphs[0].font.bold = True

# Create a group shape with two rectangles using raw XML
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

spTree34 = slide34.shapes._spTree
grp_xml = f'''<p:grpSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvGrpSpPr>
    <p:cNvPr id="100" name="Group 1"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm>
      <a:off x="457200" y="1371600"/>
      <a:ext cx="3657600" cy="2743200"/>
      <a:chOff x="0" y="0"/>
      <a:chExt cx="3657600" cy="2743200"/>
    </a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="101" name="Rect1"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="0" y="0"/>
        <a:ext cx="1828800" cy="1371600"/>
      </a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="FF6B6B"/></a:solidFill>
      <a:ln w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:lstStyle/>
      <a:p><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Child 1 (red)</a:t></a:r></a:p>
    </p:txBody>
  </p:sp>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="102" name="Rect2"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="1828800" y="1371600"/>
        <a:ext cx="1828800" cy="1371600"/>
      </a:xfrm>
      <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="4ECDC4"/></a:solidFill>
      <a:ln w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:lstStyle/>
      <a:p><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Child 2 (teal)</a:t></a:r></a:p>
    </p:txBody>
  </p:sp>
</p:grpSp>'''
spTree34.append(etree.fromstring(grp_xml))

# Second group: nested group (group inside group)
grp_nested_xml = f'''<p:grpSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvGrpSpPr>
    <p:cNvPr id="200" name="Outer Group"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm>
      <a:off x="4572000" y="1371600"/>
      <a:ext cx="4572000" cy="2743200"/>
      <a:chOff x="0" y="0"/>
      <a:chExt cx="4572000" cy="2743200"/>
    </a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="201" name="Outer Rect"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="0" y="0"/>
        <a:ext cx="4572000" cy="2743200"/>
      </a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="E8E8E8"/></a:solidFill>
      <a:ln w="12700"><a:solidFill><a:srgbClr val="999999"/></a:solidFill></a:ln>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:lstStyle/>
      <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Outer bg</a:t></a:r></a:p>
    </p:txBody>
  </p:sp>
  <p:grpSp>
    <p:nvGrpSpPr>
      <p:cNvPr id="210" name="Inner Group"/>
      <p:cNvGrpSpPr/>
      <p:nvPr/>
    </p:nvGrpSpPr>
    <p:grpSpPr>
      <a:xfrm>
        <a:off x="457200" y="457200"/>
        <a:ext cx="3657600" cy="1828800"/>
        <a:chOff x="0" y="0"/>
        <a:chExt cx="3657600" cy="1828800"/>
      </a:xfrm>
    </p:grpSpPr>
    <p:sp>
      <p:nvSpPr>
        <p:cNvPr id="211" name="Inner Circle"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="1828800" cy="1828800"/>
        </a:xfrm>
        <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="FFD93D"/></a:solidFill>
      </p:spPr>
      <p:txBody>
        <a:bodyPr anchor="ctr"/>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Nested A</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>
    <p:sp>
      <p:nvSpPr>
        <p:cNvPr id="212" name="Inner Rect"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="1828800" y="0"/>
          <a:ext cx="1828800" cy="1828800"/>
        </a:xfrm>
        <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="6C5CE7"/></a:solidFill>
      </p:spPr>
      <p:txBody>
        <a:bodyPr anchor="ctr"/>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>Nested B</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>
  </p:grpSp>
</p:grpSp>'''
spTree34.append(etree.fromstring(grp_nested_xml))

# Label
lbl34 = slide34.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
lbl34.text_frame.paragraphs[0].text = "Left: simple group (red rect + teal ellipse). Right: nested group (outer bg + inner circle + rounded rect)."
lbl34.text_frame.paragraphs[0].font.size = Pt(12)
lbl34.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 35: Connectors ────────────────────────────────────────────────────
slide35 = prs.slides.add_slide(blank)

title35 = slide35.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title35.text_frame.paragraphs[0].text = "Slide 35: Connectors (p:cxnSp)"
title35.text_frame.paragraphs[0].font.size = Pt(24)
title35.text_frame.paragraphs[0].font.bold = True

# Straight connector
spTree35 = slide35.shapes._spTree
cxn1_xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvCxnSpPr>
    <p:cNvPr id="300" name="Straight Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="1371600"/>
      <a:ext cx="3657600" cy="0"/>
    </a:xfrm>
    <a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>
    <a:ln w="25400">
      <a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>
      <a:tailEnd type="triangle"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree35.append(etree.fromstring(cxn1_xml))

# Diagonal connector with arrows on both ends
cxn2_xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvCxnSpPr>
    <p:cNvPr id="301" name="Diagonal Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="2286000"/>
      <a:ext cx="3657600" cy="1828800"/>
    </a:xfrm>
    <a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>
    <a:ln w="19050">
      <a:solidFill><a:srgbClr val="0066CC"/></a:solidFill>
      <a:prstDash val="dash"/>
      <a:headEnd type="diamond"/>
      <a:tailEnd type="stealth"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree35.append(etree.fromstring(cxn2_xml))

# Bent connector (L-shape)
cxn3_xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvCxnSpPr>
    <p:cNvPr id="302" name="Bent Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="5029200" y="1371600"/>
      <a:ext cx="3657600" cy="1371600"/>
    </a:xfrm>
    <a:prstGeom prst="bentConnector3">
      <a:avLst>
        <a:gd name="adj1" fmla="val 50000"/>
      </a:avLst>
    </a:prstGeom>
    <a:ln w="25400">
      <a:solidFill><a:srgbClr val="009900"/></a:solidFill>
      <a:tailEnd type="triangle"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree35.append(etree.fromstring(cxn3_xml))

# Curved connector
cxn4_xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvCxnSpPr>
    <p:cNvPr id="303" name="Curved Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="5029200" y="3200400"/>
      <a:ext cx="3657600" cy="1371600"/>
    </a:xfrm>
    <a:prstGeom prst="curvedConnector3">
      <a:avLst>
        <a:gd name="adj1" fmla="val 50000"/>
      </a:avLst>
    </a:prstGeom>
    <a:ln w="25400">
      <a:solidFill><a:srgbClr val="FF6600"/></a:solidFill>
      <a:headEnd type="oval"/>
      <a:tailEnd type="arrow"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree35.append(etree.fromstring(cxn4_xml))

lbl35 = slide35.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
lbl35.text_frame.paragraphs[0].text = "Left: straight + diagonal. Right: bent (green) + curved (orange)."
lbl35.text_frame.paragraphs[0].font.size = Pt(12)
lbl35.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ────────────────────────────────────────────────────────────────────────────
# Slide 36: Preset geometry shapes (triangle, diamond, arrow, star, heart, etc.)
# ────────────────────────────────────────────────────────────────────────────
slide36 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
spTree36 = slide36.shapes._spTree

# Helper: inject a preset geometry shape via raw XML
def make_prst_shape(name, prst, x_emu, y_emu, cx_emu, cy_emu, fill_hex, avlst_xml=''):
    av_xml = f'<a:avLst>{avlst_xml}</a:avLst>' if avlst_xml else '<a:avLst/>'
    return f'''<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="0" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{x_emu}" y="{y_emu}"/>
      <a:ext cx="{cx_emu}" cy="{cy_emu}"/>
    </a:xfrm>
    <a:prstGeom prst="{prst}">{av_xml}</a:prstGeom>
    <a:solidFill><a:srgbClr val="{fill_hex}"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr lang="en-US" sz="900"/><a:t>{prst}</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''

shapes_36 = [
    ("triangle", "FF6B6B"),
    ("diamond", "4ECDC4"),
    ("pentagon", "FFD93D"),
    ("hexagon", "6C5CE7"),
    ("rightArrow", "74B9FF"),
    ("star5", "FDCB6E"),
    ("heart", "E17055"),
    ("plus", "00CEC9"),
    ("flowChartDecision", "A29BFE"),
    ("chevron", "55EFC4"),
    ("parallelogram", "FF7675"),
    ("octagon", "FFEAA7"),
]

row_y = [Emu(457200), Emu(2286000)]  # 2 rows
col_x = [Emu(457200 + i * 1524000) for i in range(6)]  # 6 columns
sz = Emu(1143000)  # ~1.2 inch
for idx, (prst, fill) in enumerate(shapes_36):
    row = idx // 6
    col = idx % 6
    xml = make_prst_shape(f"prst_{prst}", prst, col_x[col], row_y[row], sz, sz, fill)
    spTree36.append(etree.fromstring(xml))

# Also add shapes with custom adj values
adj_shapes = [
    ("rightArrow_adj", "rightArrow", Emu(457200), Emu(4114800), Emu(2286000), Emu(914400), "74B9FF",
     '<a:gd name="adj1" fmla="val 70000"/><a:gd name="adj2" fmla="val 30000"/>'),
    ("star5_adj", "star5", Emu(3200400), Emu(4114800), Emu(1143000), Emu(1143000), "FDCB6E",
     '<a:gd name="adj" fmla="val 40000"/>'),
]
for name, prst, x, y, cx, cy, fill, avlst in adj_shapes:
    xml = make_prst_shape(name, prst, x, y, cx, cy, fill, avlst)
    spTree36.append(etree.fromstring(xml))

lbl36 = slide36.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl36.text_frame.paragraphs[0].text = "Preset geometry: triangle, diamond, pentagon, hexagon, rightArrow, star5, heart, plus, flowChartDecision, chevron, parallelogram, octagon"
lbl36.text_frame.paragraphs[0].font.size = Pt(10)
lbl36.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 37: Custom geometry (a:custGeom) ──
slide37 = prs.slides.add_slide(blank)
spTree37 = slide37.shapes._spTree

# Custom freeform shape 1: a star-like custom geometry with path coordinates
cust1_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="100" name="CustomStar"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="457200"/><a:ext cx="1828800" cy="1828800"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:pathLst>
        <a:path w="200" h="200">
          <a:moveTo><a:pt x="100" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="130" y="70"/></a:lnTo>
          <a:lnTo><a:pt x="200" y="80"/></a:lnTo>
          <a:lnTo><a:pt x="150" y="130"/></a:lnTo>
          <a:lnTo><a:pt x="160" y="200"/></a:lnTo>
          <a:lnTo><a:pt x="100" y="170"/></a:lnTo>
          <a:lnTo><a:pt x="40" y="200"/></a:lnTo>
          <a:lnTo><a:pt x="50" y="130"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="80"/></a:lnTo>
          <a:lnTo><a:pt x="70" y="70"/></a:lnTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="FFD700"/></a:solidFill>
    <a:ln w="19050"><a:solidFill><a:srgbClr val="B8860B"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree37.append(etree.fromstring(cust1_xml))

# Custom freeform shape 2: curved shape with cubicBezTo
cust2_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="101" name="CustomCurve"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="2743200" y="457200"/><a:ext cx="1828800" cy="1828800"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:pathLst>
        <a:path w="100" h="100">
          <a:moveTo><a:pt x="0" y="100"/></a:moveTo>
          <a:cubicBezTo>
            <a:pt x="0" y="0"/>
            <a:pt x="100" y="0"/>
            <a:pt x="100" y="100"/>
          </a:cubicBezTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="87CEEB"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="4169E1"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree37.append(etree.fromstring(cust2_xml))

# Custom shape 3: with guide formulas
cust3_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="102" name="CustomGuide"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="457200"/><a:ext cx="1828800" cy="1828800"/></a:xfrm>
    <a:custGeom>
      <a:avLst>
        <a:gd name="adj" fmla="val 50000"/>
      </a:avLst>
      <a:gdLst>
        <a:gd name="x1" fmla="*/ w adj 100000"/>
        <a:gd name="y1" fmla="*/ h adj 100000"/>
      </a:gdLst>
      <a:pathLst>
        <a:path>
          <a:moveTo><a:pt x="0" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="x1" y="0"/></a:lnTo>
          <a:lnTo><a:pt x="x1" y="y1"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="y1"/></a:lnTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="98FB98"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="228B22"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree37.append(etree.fromstring(cust3_xml))

lbl37 = slide37.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl37.text_frame.paragraphs[0].text = "Custom geometry (a:custGeom): freeform star, bezier curve, guide-based rect"
lbl37.text_frame.paragraphs[0].font.size = Pt(10)
lbl37.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 38: Gear shapes ────────────────────────────────────────────────────
slide38 = prs.slides.add_slide(blank)
spTree38 = slide38.shapes._spTree

# Gear6
gear6_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="201" name="Gear6"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="gear6"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="2F528F"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree38.append(etree.fromstring(gear6_xml))

# Gear9
gear9_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="202" name="Gear9"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="3429000" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="gear9"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="ED7D31"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="AE5A21"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree38.append(etree.fromstring(gear9_xml))

# Gear6 with custom adj (deeper teeth)
gear6_adj_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="203" name="Gear6Deep"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="6400800" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="gear6"><a:avLst><a:gd name="adj" fmla="val 50000"/></a:avLst></a:prstGeom>
    <a:solidFill><a:srgbClr val="70AD47"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="507E32"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree38.append(etree.fromstring(gear6_adj_xml))

lbl38 = slide38.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl38.text_frame.paragraphs[0].text = "Gear shapes: gear6 (default), gear9 (default), gear6 (adj=50000 deeper teeth)"
lbl38.text_frame.paragraphs[0].font.size = Pt(10)
lbl38.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 39: Text rectangles ────────────────────────────────────────────────
slide39 = prs.slides.add_slide(blank)
spTree39 = slide39.shapes._spTree

# Triangle with text
tri_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="301" name="Triangle"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="triangle"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFD700"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="B8860B"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Text in Triangle</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree39.append(etree.fromstring(tri_xml))

# Diamond with text
diamond_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="302" name="Diamond"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="3429000" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="diamond"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="87CEEB"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="4682B4"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Diamond</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree39.append(etree.fromstring(diamond_xml))

# Right arrow with text
arrow_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="303" name="RightArrow"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="6400800" y="914400"/><a:ext cx="2743200" cy="1371600"/></a:xfrm>
    <a:prstGeom prst="rightArrow"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FF6347"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="8B0000"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Arrow</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree39.append(etree.fromstring(arrow_xml))

lbl39 = slide39.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl39.text_frame.paragraphs[0].text = "Text rectangles: triangle (text in lower half), diamond (text inscribed), right arrow (text in shaft)"
lbl39.text_frame.paragraphs[0].font.size = Pt(10)
lbl39.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 40: Connection points ──────────────────────────────────────────────
slide40 = prs.slides.add_slide(blank)
spTree40 = slide40.shapes._spTree

# Shape 1: Rectangle (target for connectors)
rect1_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="401" name="Rect1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="914400" y="914400"/><a:ext cx="1828800" cy="1371600"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
  </p:spPr>
  <p:txBody><a:bodyPr anchor="ctr"/><a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>Start</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree40.append(etree.fromstring(rect1_xml))

# Shape 2: Rectangle (target for connectors)
rect2_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="402" name="Rect2"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5486400" y="914400"/><a:ext cx="1828800" cy="1371600"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="ED7D31"/></a:solidFill>
  </p:spPr>
  <p:txBody><a:bodyPr anchor="ctr"/><a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>End</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree40.append(etree.fromstring(rect2_xml))

# Connector with stCxn/endCxn references
cxn_xml = '''<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvCxnSpPr>
    <p:cNvPr id="403" name="Connector"/>
    <p:cNvCxnSpPr>
      <a:stCxn id="401" idx="3"/>
      <a:endCxn id="402" idx="1"/>
    </p:cNvCxnSpPr>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm><a:off x="2743200" y="1600200"/><a:ext cx="2743200" cy="0"/></a:xfrm>
    <a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>
    <a:ln w="25400">
      <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
      <a:tailEnd type="triangle"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree40.append(etree.fromstring(cxn_xml))

# Custom geometry with a:cxnLst
cust_cxn_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="404" name="CustomWithCxn"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="3200400" y="3429000"/><a:ext cx="2286000" cy="1828800"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:cxnLst>
        <a:cxn ang="0"><a:pos x="r" y="vc"/></a:cxn>
        <a:cxn ang="5400000"><a:pos x="hc" y="b"/></a:cxn>
        <a:cxn ang="10800000"><a:pos x="l" y="vc"/></a:cxn>
        <a:cxn ang="16200000"><a:pos x="hc" y="t"/></a:cxn>
      </a:cxnLst>
      <a:rect l="l" t="t" r="r" b="b"/>
      <a:pathLst>
        <a:path w="100" h="100">
          <a:moveTo><a:pt x="50" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="100" y="50"/></a:lnTo>
          <a:lnTo><a:pt x="50" y="100"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="50"/></a:lnTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="9370DB"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="4B0082"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree40.append(etree.fromstring(cust_cxn_xml))

lbl40 = slide40.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl40.text_frame.paragraphs[0].text = "Connection points: connector with stCxn/endCxn refs, custom geom with a:cxnLst"
lbl40.text_frame.paragraphs[0].font.size = Pt(10)
lbl40.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

