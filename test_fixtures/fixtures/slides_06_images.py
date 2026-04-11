"""Slides 43-51: image crop/alpha, external refs, lum, duotone, bg pattern, line gradient, hyperlinks, shape effects, 3D.

Mechanically carved from the legacy gen_test_features.py. Runs its
slide-building code at import time against the shared '_ctx.prs'.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
import base64, io, os, re, struct, tempfile, zipfile, zlib

from ._ctx import prs, blank, set_fill_xml

# ── Slide 43: Image crop (srcRect) + alpha (alphaModFix) ──────────────────────

slide43 = prs.slides.add_slide(blank)

# Create a test image: 2x2 pixel PNG with 4 colored quadrants (red/green/blue/yellow)
def make_test_png_4color():
    """Create a 4x4 PNG with colored quadrants: TL=red, TR=green, BL=blue, BR=yellow."""
    width, height = 4, 4
    def px(r, g, b):
        return bytes([r, g, b])
    rows = []
    for y in range(height):
        row = b'\x00'  # filter byte
        for x in range(width):
            if y < 2 and x < 2:
                row += px(255, 0, 0)      # TL = red
            elif y < 2 and x >= 2:
                row += px(0, 255, 0)      # TR = green
            elif y >= 2 and x < 2:
                row += px(0, 0, 255)      # BL = blue
            else:
                row += px(255, 255, 0)    # BR = yellow
        rows.append(row)
    raw = b''.join(rows)
    def png_chunk(chunk_type, data):
        c = chunk_type + data
        crc = struct.pack('>I', zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack('>I', len(data)) + c + crc
    ihdr = struct.pack('>IIBBBBB', width, height, 8, 2, 0, 0, 0)
    return (b'\x89PNG\r\n\x1a\n' +
            png_chunk(b'IHDR', ihdr) +
            png_chunk(b'IDAT', zlib.compress(raw)) +
            png_chunk(b'IEND', b''))

test_png = make_test_png_4color()
tmp_img43 = os.path.join(tempfile.gettempdir(), '_test_crop.png')
with open(tmp_img43, 'wb') as f:
    f.write(test_png)

img_part43, img_rid43 = slide43.part.get_or_add_image_part(tmp_img43)
os.unlink(tmp_img43)

# Shape 1: Picture with srcRect crop (crop 25% from each side — shows center)
ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

pic43_crop = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="100" name="CropPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid43}"/>
    <a:srcRect l="25000" t="25000" r="25000" b="25000"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="914400"/>
      <a:ext cx="2743200" cy="2743200"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 2: Picture with alphaModFix (50% opacity)
pic43_alpha = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="101" name="AlphaPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid43}">
      <a:alphaModFix amt="50000"/>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="3657600" y="914400"/>
      <a:ext cx="2743200" cy="2743200"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 3: Picture with crop + alpha combined (crop left 50%, alpha 75%)
pic43_both = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="102" name="CropAlphaPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid43}">
      <a:alphaModFix amt="75000"/>
    </a:blip>
    <a:srcRect l="50000"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="6858000" y="914400"/>
      <a:ext cx="2286000" cy="2743200"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 4: AutoShape with blipFill + srcRect crop
s43_autoshape = slide43.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(3), Inches(2.5))
s43_autoshape.text = ""
set_fill_xml(s43_autoshape, f'''<a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <a:blip r:embed="{img_rid43}"/>
  <a:srcRect l="0" t="50000" r="0" b="0"/>
  <a:stretch><a:fillRect/></a:stretch>
</a:blipFill>''')

spTree43 = slide43._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree43.append(etree.fromstring(pic43_crop))
spTree43.append(etree.fromstring(pic43_alpha))
spTree43.append(etree.fromstring(pic43_both))

lbl43 = slide43.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl43.text_frame.paragraphs[0].text = "Slide 43: Image crop (srcRect) + alpha (alphaModFix)"
lbl43.text_frame.paragraphs[0].font.size = Pt(18)
lbl43.text_frame.paragraphs[0].font.bold = True

# ── Slide 44: External image reference (TargetMode="External") ───────────────

slide44 = prs.slides.add_slide(blank)

# Use stable, CORS-friendly public image URLs
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Image 1: Wikimedia Commons PNG (small, stable, CORS-friendly)
ext_url1 = "https://upload.wikimedia.org/wikipedia/commons/thumb/4/47/PNG_transparency_demonstration_1.png/280px-PNG_transparency_demonstration_1.png"
ext_rid1 = slide44.part.relate_to(ext_url1, RT.IMAGE, is_external=True)

# Image 2: picsum.photos direct image (stable, specific ID = no redirect)
ext_url2 = "https://picsum.photos/id/237/200/300"
ext_rid2 = slide44.part.relate_to(ext_url2, RT.IMAGE, is_external=True)

pic44_ext1 = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="200" name="ExtPic1"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{ext_rid1}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="914400"/>
      <a:ext cx="3657600" cy="3657600"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

pic44_ext2 = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="201" name="ExtPic2"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{ext_rid2}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="4572000" y="914400"/>
      <a:ext cx="2743200" cy="3657600"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

spTree44 = slide44._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree44.append(etree.fromstring(pic44_ext1))
spTree44.append(etree.fromstring(pic44_ext2))

lbl44 = slide44.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl44.text_frame.paragraphs[0].text = "Slide 44: External image reference (TargetMode=External)"
lbl44.text_frame.paragraphs[0].font.size = Pt(18)
lbl44.text_frame.paragraphs[0].font.bold = True

lbl44b = slide44.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.5))
lbl44b.text_frame.paragraphs[0].text = f"Left: Wikimedia Commons PNG\nRight: picsum.photos (id=237)"
lbl44b.text_frame.paragraphs[0].font.size = Pt(12)

# ── Slide 45: Image effects — brightness/contrast (a:lum) + duotone (a:duotone) ──
slide45 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Reuse the test PNG for effects
tmp_img45 = os.path.join(tempfile.gettempdir(), '_test_effects.png')
with open(tmp_img45, 'wb') as f:
    f.write(test_png)
img_part45, img_rid45 = slide45.part.get_or_add_image_part(tmp_img45)
os.unlink(tmp_img45)
# Add second image for duotone
tmp_img45b = os.path.join(tempfile.gettempdir(), '_test_effects2.png')
with open(tmp_img45b, 'wb') as f:
    f.write(test_png)
img_part45b, img_rid45b = slide45.part.get_or_add_image_part(tmp_img45b)
os.unlink(tmp_img45b)
# Third image for combined brightness + contrast
tmp_img45c = os.path.join(tempfile.gettempdir(), '_test_effects3.png')
with open(tmp_img45c, 'wb') as f:
    f.write(test_png)
img_part45c, img_rid45c = slide45.part.get_or_add_image_part(tmp_img45c)
os.unlink(tmp_img45c)

# Shape 1: Picture with brightness +50% (a:lum bright="50000")
pic45_bright = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="300" name="BrightPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid45}">
      <a:lum bright="50000"/>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="914400"/>
      <a:ext cx="2286000" cy="2286000"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 2: Picture with contrast -30% (a:lum contrast="-30000")
pic45_contrast = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="301" name="ContrastPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid45b}">
      <a:lum contrast="-30000"/>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="3429000" y="914400"/>
      <a:ext cx="2286000" cy="2286000"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 3: Picture with bright+contrast combined (a:lum bright="20000" contrast="40000")
pic45_both = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="302" name="BrightContrastPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid45c}">
      <a:lum bright="20000" contrast="40000"/>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="6400800" y="914400"/>
      <a:ext cx="2286000" cy="2286000"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

spTree45 = slide45._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree45.append(etree.fromstring(pic45_bright))
spTree45.append(etree.fromstring(pic45_contrast))
spTree45.append(etree.fromstring(pic45_both))

lbl45 = slide45.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl45.text_frame.paragraphs[0].text = "Slide 45: Image effects — brightness/contrast (a:lum)"
lbl45.text_frame.paragraphs[0].font.size = Pt(18)
lbl45.text_frame.paragraphs[0].font.bold = True

lbl45b = slide45.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(1))
lbl45b.text_frame.paragraphs[0].text = "Left: bright +50% | Center: contrast -30% | Right: bright +20% contrast +40%"
lbl45b.text_frame.paragraphs[0].font.size = Pt(12)

# ── Slide 46: Duotone (a:duotone) + color change (a:clrChange) ──
slide46 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Image for duotone
tmp_img46 = os.path.join(tempfile.gettempdir(), '_test_duotone.png')
with open(tmp_img46, 'wb') as f:
    f.write(test_png)
img_part46, img_rid46 = slide46.part.get_or_add_image_part(tmp_img46)
os.unlink(tmp_img46)
# Image for clrChange
tmp_img46b = os.path.join(tempfile.gettempdir(), '_test_clrchange.png')
with open(tmp_img46b, 'wb') as f:
    f.write(test_png)
img_part46b, img_rid46b = slide46.part.get_or_add_image_part(tmp_img46b)
os.unlink(tmp_img46b)

# Shape 1: Picture with duotone (dark blue → light yellow)
pic46_duotone = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="400" name="DuotonePic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid46}">
      <a:duotone>
        <a:srgbClr val="000080"/>
        <a:srgbClr val="FFFF00"/>
      </a:duotone>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="914400"/>
      <a:ext cx="3657600" cy="3657600"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 2: Picture with clrChange (white → transparent)
pic46_clrchange = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="401" name="ClrChangePic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid46b}">
      <a:clrChange>
        <a:clrFrom><a:srgbClr val="FF0000"/></a:clrFrom>
        <a:clrTo><a:srgbClr val="00FF00"/></a:clrTo>
      </a:clrChange>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="4572000" y="914400"/>
      <a:ext cx="3657600" cy="3657600"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

spTree46 = slide46._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree46.append(etree.fromstring(pic46_duotone))
spTree46.append(etree.fromstring(pic46_clrchange))

lbl46 = slide46.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl46.text_frame.paragraphs[0].text = "Slide 46: Duotone (a:duotone) + Color change (a:clrChange)"
lbl46.text_frame.paragraphs[0].font.size = Pt(18)
lbl46.text_frame.paragraphs[0].font.bold = True

lbl46b = slide46.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1))
lbl46b.text_frame.paragraphs[0].text = "Left: Duotone (navy→yellow) | Right: clrChange (red→green, data preserved)"
lbl46b.text_frame.paragraphs[0].font.size = Pt(12)

########################################
# Slide 47: Background pattern fill
########################################
slide47 = prs.slides.add_slide(prs.slide_layouts[6])

# Inject pattern fill background via raw XML
ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
bg47_xml = f"""<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:bgPr>
    <a:pattFill prst="ltDnDiag">
      <a:fgClr><a:srgbClr val="3366CC"/></a:fgClr>
      <a:bgClr><a:srgbClr val="FFFFFF"/></a:bgClr>
    </a:pattFill>
    <a:effectLst/>
  </p:bgPr>
</p:bg>"""
csld47 = slide47._element.find(f'.//{ns_p}cSld')
# Remove existing bg if any
old_bg47 = csld47.find(f'{ns_p}bg')
if old_bg47 is not None:
    csld47.remove(old_bg47)
bg47_elem = etree.fromstring(bg47_xml)
csld47.insert(0, bg47_elem)

lbl47 = slide47.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl47.text_frame.paragraphs[0].text = "Slide 47: Background pattern fill (ltDnDiag)"
lbl47.text_frame.paragraphs[0].font.size = Pt(18)
lbl47.text_frame.paragraphs[0].font.bold = True

########################################
# Slide 48: Line gradient / pattern fill
########################################
slide48 = prs.slides.add_slide(prs.slide_layouts[6])

# Shape 1: Rectangle with gradient stroke
sp48_grad_xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="100" name="GradStrokeRect"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="3657600" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="F0F0F0"/></a:solidFill>
    <a:ln w="76200">
      <a:gradFill>
        <a:gsLst>
          <a:gs pos="0"><a:srgbClr val="FF0000"/></a:gs>
          <a:gs pos="50000"><a:srgbClr val="00FF00"/></a:gs>
          <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
        </a:gsLst>
        <a:lin ang="0" scaled="1"/>
      </a:gradFill>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Gradient stroke</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

# Shape 2: Ellipse with pattern stroke
sp48_patt_xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="101" name="PattStrokeEllipse"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="1371600"/><a:ext cx="3657600" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFCC"/></a:solidFill>
    <a:ln w="57150">
      <a:pattFill prst="smCheck">
        <a:fgClr><a:srgbClr val="990000"/></a:fgClr>
        <a:bgClr><a:srgbClr val="FFCC00"/></a:bgClr>
      </a:pattFill>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Pattern stroke</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

spTree48 = slide48._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree48.append(etree.fromstring(sp48_grad_xml))
spTree48.append(etree.fromstring(sp48_patt_xml))

lbl48 = slide48.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl48.text_frame.paragraphs[0].text = "Slide 48: Line gradient fill + Line pattern fill"
lbl48.text_frame.paragraphs[0].font.size = Pt(18)
lbl48.text_frame.paragraphs[0].font.bold = True

########################################
# Slide 49: Shape hyperlinks + color modifiers
########################################
slide49 = prs.slides.add_slide(prs.slide_layouts[6])

# Shape with hlinkClick on cNvPr
sp49_link_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="200" name="LinkedShape">
      <a:hlinkClick r:id="rId99"/>
    </p:cNvPr>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="3657600" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    <a:ln w="25400"><a:solidFill><a:srgbClr val="2F5597"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1600" b="1"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>Click me (shape link)</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

# Shape with color modifier: complement
sp49_comp_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="201" name="CompColor"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="1371600"/><a:ext cx="3657600" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FF6600"><a:comp/></a:srgbClr></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>comp(#FF6600) = complementary hue</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

# Shape with color modifier: inv
sp49_inv_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="202" name="InvColor"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="3657600"/><a:ext cx="3657600" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="3366CC"><a:inv/></a:srgbClr></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>inv(#3366CC) → #CC9933</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

# Shape with hueMod
sp49_huemod_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="203" name="HueModColor"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="3657600"/><a:ext cx="3657600" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FF0000"><a:hueMod val="50000"/></a:srgbClr></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>hueMod 50% of red</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

spTree49 = slide49._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree49.append(etree.fromstring(sp49_link_xml))
spTree49.append(etree.fromstring(sp49_comp_xml))
spTree49.append(etree.fromstring(sp49_inv_xml))
spTree49.append(etree.fromstring(sp49_huemod_xml))

lbl49 = slide49.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl49.text_frame.paragraphs[0].text = "Slide 49: Shape hyperlinks + Color modifiers (comp/inv/hueMod)"
lbl49.text_frame.paragraphs[0].font.size = Pt(18)
lbl49.text_frame.paragraphs[0].font.bold = True

# Add a fake relationship for the shape link (rId99)
# We need to add this to the slide49 rels
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
slide49_part = slide49.part
slide49_part.rels.get_or_add_ext_rel(RT.HYPERLINK, 'https://example.com/shape-link')

# ── Slide 50: Shape effects (outerShdw, innerShdw, glow, softEdge) ──
slide50 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

sp50_outer_shadow_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="301" name="OuterShadow"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    <a:effectLst>
      <a:outerShdw blurRad="152400" dist="114300" dir="5400000" algn="ctr" rotWithShape="0">
        <a:srgbClr val="000000"><a:alpha val="40000"/></a:srgbClr>
      </a:outerShdw>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Outer Shadow</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp50_inner_shadow_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="302" name="InnerShadow"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="3657600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="70AD47"/></a:solidFill>
    <a:effectLst>
      <a:innerShdw blurRad="114300" dist="76200" dir="2700000">
        <a:srgbClr val="000000"><a:alpha val="50000"/></a:srgbClr>
      </a:innerShdw>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Inner Shadow</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp50_glow_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="303" name="Glow"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="1371600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="ED7D31"/></a:solidFill>
    <a:effectLst>
      <a:glow rad="228600">
        <a:srgbClr val="FFC000"><a:alpha val="60000"/></a:srgbClr>
      </a:glow>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Glow Effect</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp50_softedge_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="304" name="SoftEdge"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="3657600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="A855F7"/></a:solidFill>
    <a:effectLst>
      <a:softEdge rad="317500"/>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Soft Edge</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

spTree50 = slide50._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree50.append(etree.fromstring(sp50_outer_shadow_xml))
spTree50.append(etree.fromstring(sp50_inner_shadow_xml))
spTree50.append(etree.fromstring(sp50_glow_xml))
spTree50.append(etree.fromstring(sp50_softedge_xml))

lbl50 = slide50.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl50.text_frame.paragraphs[0].text = "Slide 50: Shape effects (outerShdw / innerShdw / glow / softEdge)"
lbl50.text_frame.paragraphs[0].font.size = Pt(18)
lbl50.text_frame.paragraphs[0].font.bold = True

# ── Slide 51: 3D effects (bevel, scene3d, sp3d) + text shadow ──
slide51 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

sp51_bevel_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="401" name="Bevel3D"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    <a:scene3d>
      <a:camera prst="orthographicFront"/>
      <a:lightRig rig="threePt" dir="t"/>
    </a:scene3d>
    <a:sp3d prstMaterial="plastic">
      <a:bevelT w="152400" h="50800" prst="relaxedInset"/>
    </a:sp3d>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>3D Bevel</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp51_text_shadow_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="402" name="TextShadow"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="1371600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="2400" b="1">
      <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
      <a:effectLst>
        <a:outerShdw blurRad="38100" dist="38100" dir="5400000">
          <a:srgbClr val="000000"><a:alpha val="40000"/></a:srgbClr>
        </a:outerShdw>
      </a:effectLst>
    </a:rPr><a:t>Text with Shadow</a:t></a:r></a:p>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="2000">
      <a:solidFill><a:srgbClr val="FF4444"/></a:solidFill>
      <a:effectLst>
        <a:glow rad="101600">
          <a:srgbClr val="FF0000"><a:alpha val="50000"/></a:srgbClr>
        </a:glow>
      </a:effectLst>
    </a:rPr><a:t>Text with Glow</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp51_extrusion_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="403" name="Extrusion3D"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="3657600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="70AD47"/></a:solidFill>
    <a:scene3d>
      <a:camera prst="isometricLeftDown"/>
      <a:lightRig rig="balanced" dir="t"/>
    </a:scene3d>
    <a:sp3d extrusionH="76200" contourW="12700" prstMaterial="warmMatte">
      <a:bevelT w="101600" h="38100" prst="circle"/>
      <a:bevelB w="50800" h="25400"/>
      <a:extrusionClr><a:srgbClr val="5B8C3C"/></a:extrusionClr>
      <a:contourClr><a:srgbClr val="2E5B14"/></a:contourClr>
    </a:sp3d>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>3D Extrusion</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

spTree51 = slide51._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree51.append(etree.fromstring(sp51_bevel_xml))
spTree51.append(etree.fromstring(sp51_text_shadow_xml))
spTree51.append(etree.fromstring(sp51_extrusion_xml))

lbl51 = slide51.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl51.text_frame.paragraphs[0].text = "Slide 51: 3D effects (bevel/extrusion/scene3d) + text shadow/glow"
lbl51.text_frame.paragraphs[0].font.size = Pt(18)
lbl51.text_frame.paragraphs[0].font.bold = True

