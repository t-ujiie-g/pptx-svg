"""Slides 41-42: table cell merge, borders, diagonal borders, tblPr flags.

Mechanically carved from the legacy gen_test_features.py. Runs its
slide-building code at import time against the shared '_ctx.prs'.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
import base64, io, os, re, struct, tempfile, zipfile, zlib

from ._ctx import prs, blank

# ── Slide 41: Table cell merge + borders + margins + anchor ──────────────────
slide41 = prs.slides.add_slide(blank)
ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

# Build table XML directly for full control over cell merge/border/margin/anchor
table_xml = f'''<p:graphicFrame xmlns:a="{ns_a[1:-1]}" xmlns:p="{ns_p[1:-1]}" xmlns:r="{ns_r[1:-1]}">
  <p:nvGraphicFramePr>
    <p:cNvPr id="100" name="Table41"/>
    <p:cNvGraphicFramePr><a:graphicFrameLocks noGrp="1"/></p:cNvGraphicFramePr>
    <p:nvPr/>
  </p:nvGraphicFramePr>
  <p:xfrm>
    <a:off x="457200" y="457200"/>
    <a:ext cx="8229600" cy="4572000"/>
  </p:xfrm>
  <a:graphic>
    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
      <a:tbl>
        <a:tblPr firstRow="1" bandRow="1"/>
        <a:tblGrid>
          <a:gridCol w="2743200"/>
          <a:gridCol w="2743200"/>
          <a:gridCol w="2743200"/>
        </a:tblGrid>
        <!-- Row 1: horizontal merge (gridSpan=2) + custom borders -->
        <a:tr h="914400">
          <a:tc gridSpan="2">
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Merged 2 cols</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr marL="182880" marR="91440" marT="91440" marB="45720" anchor="ctr">
              <a:lnL w="25400"><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill></a:lnL>
              <a:lnR w="25400"><a:solidFill><a:srgbClr val="0000FF"/></a:solidFill></a:lnR>
              <a:lnT w="38100"><a:solidFill><a:srgbClr val="00AA00"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="FF8800"/></a:solidFill></a:lnB>
              <a:solidFill><a:srgbClr val="DDEEFF"/></a:solidFill>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Normal cell</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr anchor="b">
              <a:lnL w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:lnB>
              <a:solidFill><a:srgbClr val="FFFFDD"/></a:solidFill>
            </a:tcPr>
          </a:tc>
        </a:tr>
        <!-- Row 2: vertical merge start (rowSpan=2) + noFill border -->
        <a:tr h="914400">
          <a:tc rowSpan="2">
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Merged 2 rows</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr anchor="ctr">
              <a:lnL w="19050"><a:solidFill><a:srgbClr val="9900CC"/></a:solidFill></a:lnL>
              <a:lnR w="19050"><a:solidFill><a:srgbClr val="9900CC"/></a:solidFill></a:lnR>
              <a:lnT w="19050"><a:solidFill><a:srgbClr val="9900CC"/></a:solidFill></a:lnT>
              <a:lnB w="19050"><a:solidFill><a:srgbClr val="9900CC"/></a:solidFill></a:lnB>
              <a:solidFill><a:srgbClr val="F0E0FF"/></a:solidFill>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Cell B2</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnL w="12700"><a:noFill/></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnB>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Cell C2</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnL w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnB>
            </a:tcPr>
          </a:tc>
        </a:tr>
        <!-- Row 3: vertical merge continuation -->
        <a:tr h="914400">
          <a:tc vMerge="1">
            <a:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Cell B3 (large margin)</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr marL="274320" marT="137160">
              <a:lnL w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnB>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Cell C3 (top anchor)</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr anchor="t">
              <a:lnL w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnB>
              <a:solidFill><a:srgbClr val="EEFFEE"/></a:solidFill>
            </a:tcPr>
          </a:tc>
        </a:tr>
      </a:tbl>
    </a:graphicData>
  </a:graphic>
</p:graphicFrame>'''

spTree41 = slide41._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree41.append(etree.fromstring(table_xml))

lbl41 = slide41.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl41.text_frame.paragraphs[0].text = "Table: cell merge (gridSpan/rowSpan/vMerge) + borders + margins + anchor"
lbl41.text_frame.paragraphs[0].font.size = Pt(10)
lbl41.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 42: Table diagonal borders + tblPr flags + tblStyleId ──────────────
slide42 = prs.slides.add_slide(blank)

table42_xml = f'''<p:graphicFrame xmlns:a="{ns_a[1:-1]}" xmlns:p="{ns_p[1:-1]}" xmlns:r="{ns_r[1:-1]}">
  <p:nvGraphicFramePr>
    <p:cNvPr id="200" name="Table42"/>
    <p:cNvGraphicFramePr><a:graphicFrameLocks noGrp="1"/></p:cNvGraphicFramePr>
    <p:nvPr/>
  </p:nvGraphicFramePr>
  <p:xfrm>
    <a:off x="457200" y="457200"/>
    <a:ext cx="8229600" cy="3657600"/>
  </p:xfrm>
  <a:graphic>
    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
      <a:tbl>
        <a:tblPr firstRow="1" lastRow="1" firstCol="0" lastCol="0" bandRow="1" bandCol="0"
                 tblStyleId="{{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}}"/>
        <a:tblGrid>
          <a:gridCol w="2743200"/>
          <a:gridCol w="2743200"/>
          <a:gridCol w="2743200"/>
        </a:tblGrid>
        <!-- Row 1: diagonal borders -->
        <a:tr h="914400">
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200" b="1"/><a:t>TL-BR diagonal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnTlToBr w="19050"><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill></a:lnTlToBr>
              <a:solidFill><a:srgbClr val="FFF0F0"/></a:solidFill>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200" b="1"/><a:t>BL-TR diagonal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnBlToTr w="19050"><a:solidFill><a:srgbClr val="0000FF"/></a:solidFill></a:lnBlToTr>
              <a:solidFill><a:srgbClr val="F0F0FF"/></a:solidFill>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200" b="1"/><a:t>Both diagonals</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnTlToBr w="12700"><a:solidFill><a:srgbClr val="008800"/></a:solidFill></a:lnTlToBr>
              <a:lnBlToTr w="12700"><a:solidFill><a:srgbClr val="880000"/></a:solidFill></a:lnBlToTr>
              <a:solidFill><a:srgbClr val="FFFFF0"/></a:solidFill>
            </a:tcPr>
          </a:tc>
        </a:tr>
        <!-- Row 2: band row test -->
        <a:tr h="914400">
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Band row 1</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Normal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Normal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
        </a:tr>
        <!-- Row 3: band row test -->
        <a:tr h="914400">
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Band row 2</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Normal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Last row cell</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
        </a:tr>
      </a:tbl>
    </a:graphicData>
  </a:graphic>
</p:graphicFrame>'''

spTree42 = slide42._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree42.append(etree.fromstring(table42_xml))

lbl42 = slide42.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl42.text_frame.paragraphs[0].text = "Table: diagonal borders (lnTlToBr/lnBlToTr) + tblPr flags + tblStyleId"
lbl42.text_frame.paragraphs[0].font.size = Pt(10)
lbl42.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

