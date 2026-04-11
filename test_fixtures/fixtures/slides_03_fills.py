"""Slides 23-33: gradient/alpha/blip/pattern fills and stroke/arrow/join styles.

Mechanically carved from the legacy gen_test_features.py. Runs its
slide-building code at import time against the shared '_ctx.prs'.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
import base64, io, os, re, struct, tempfile, zipfile, zlib

from ._ctx import prs, blank

# ── Slide 23: Linear gradient fills ──────────────────────────────────────────

slide23 = prs.slides.add_slide(blank)

# Helper to inject gradFill XML into a shape's spPr (after prstGeom, replacing solidFill)
def set_gradient_fill(shape, grad_xml_str):
    """Replace any existing fill with a:gradFill in the shape's spPr."""
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    # spPr can be under a: or p: namespace depending on context
    spPr = shape._element.find(f'{ns_p}spPr')
    if spPr is None:
        spPr = shape._element.find(f'.//{ns_a}spPr')
    if spPr is None:
        print(f"WARNING: spPr not found on shape")
        return
    # Remove existing fills (always a: namespace children)
    for tag in ('solidFill', 'noFill', 'gradFill'):
        for el in spPr.findall(f'{ns_a}{tag}'):
            spPr.remove(el)
    # Insert after prstGeom (a: namespace child)
    prstGeom = spPr.find(f'{ns_a}prstGeom')
    grad_el = etree.fromstring(grad_xml_str)
    if prstGeom is not None:
        prstGeom.addnext(grad_el)
    else:
        spPr.append(grad_el)

# Shape 1: 3-stop linear gradient (left→right, ang=0)
s23a = slide23.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(3), Inches(2))
s23a.text = "Linear 0deg"
s23a.text_frame.paragraphs[0].font.size = Pt(18)
s23a.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s23a, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FF0000"/></a:gs>
    <a:gs pos="50000"><a:srgbClr val="FFFF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="1"/>
</a:gradFill>''')

# Shape 2: 2-stop linear gradient (top→bottom, ang=5400000 = 90°)
s23b = slide23.shapes.add_shape(1, Inches(4), Inches(0.5), Inches(3), Inches(2))
s23b.text = "Linear 90deg"
s23b.text_frame.paragraphs[0].font.size = Pt(18)
s23b.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s23b, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="003366"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="66CCFF"/></a:gs>
  </a:gsLst>
  <a:lin ang="5400000" scaled="1"/>
</a:gradFill>''')

# Shape 3: Diagonal linear gradient (45° = 2700000), rotWithShape="0"
s23c = slide23.shapes.add_shape(1, Inches(7.5), Inches(0.5), Inches(2), Inches(2))
s23c.text = "Linear 45deg"
s23c.text_frame.paragraphs[0].font.size = Pt(14)
set_gradient_fill(s23c, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="0">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="00FF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="006600"/></a:gs>
  </a:gsLst>
  <a:lin ang="2700000" scaled="1"/>
</a:gradFill>''')


# ── Slide 24: Radial/path gradient fills + gradient on ellipse ───────────────

slide24 = prs.slides.add_slide(blank)

# Shape 1: Radial circle gradient
s24a = slide24.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(3), Inches(2.5))
s24a.text = "Radial circle"
s24a.text_frame.paragraphs[0].font.size = Pt(18)
s24a.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s24a, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="FF6600"/></a:gs>
  </a:gsLst>
  <a:path path="circle">
    <a:fillToRect l="50000" t="50000" r="50000" b="50000"/>
  </a:path>
</a:gradFill>''')

# Shape 2: Radial rect gradient with off-center fillToRect
s24b = slide24.shapes.add_shape(1, Inches(4), Inches(0.5), Inches(3), Inches(2.5))
s24b.text = "Radial rect"
s24b.text_frame.paragraphs[0].font.size = Pt(18)
set_gradient_fill(s24b, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FFFF00"/></a:gs>
    <a:gs pos="50000"><a:srgbClr val="FF0000"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="660000"/></a:gs>
  </a:gsLst>
  <a:path path="rect">
    <a:fillToRect l="25000" t="25000" r="75000" b="75000"/>
  </a:path>
</a:gradFill>''')

# Shape 3: Gradient on ellipse (oval)
from pptx.enum.shapes import MSO_SHAPE
s24c = slide24.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(0.5), Inches(2), Inches(2.5))
s24c.text = "Ellipse grad"
s24c.text_frame.paragraphs[0].font.size = Pt(14)
s24c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s24c, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="9933FF"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="330066"/></a:gs>
  </a:gsLst>
  <a:lin ang="16200000" scaled="1"/>
</a:gradFill>''')


# ── Slide 25: Gradient background ────────────────────────────────────────────

slide25 = prs.slides.add_slide(blank)

# Inject gradient background via XML
ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
cSld25 = slide25._element.find(f'{ns_p}cSld')
spTree25 = cSld25.find(f'{ns_p}spTree')
bg_xml = f'''<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:bgPr>
    <a:gradFill>
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="1B2838"/></a:gs>
        <a:gs pos="50000"><a:srgbClr val="2A475E"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="66C0F4"/></a:gs>
      </a:gsLst>
      <a:lin ang="5400000" scaled="0"/>
    </a:gradFill>
    <a:effectLst/>
  </p:bgPr>
</p:bg>'''
bg_el = etree.fromstring(bg_xml)
cSld25.insert(list(cSld25).index(spTree25), bg_el)

# Text shape on gradient background (noFill)
s25a = slide25.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
s25a.text = "Gradient Background"
s25a.text_frame.paragraphs[0].font.size = Pt(36)
s25a.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
s25a.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ── Slide 26: Alpha/transparency ─────────────────────────────────────────────

slide26 = prs.slides.add_slide(blank)

def set_fill_xml(shape, fill_xml_str):
    """Replace any existing fill with custom fill XML in the shape's spPr."""
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    spPr = shape._element.find(f'{ns_p}spPr')
    if spPr is None:
        spPr = shape._element.find(f'.//{ns_a}spPr')
    if spPr is None:
        return
    for tag in ('solidFill', 'noFill', 'gradFill', 'blipFill', 'pattFill'):
        for el in spPr.findall(f'{ns_a}{tag}'):
            spPr.remove(el)
    prstGeom = spPr.find(f'{ns_a}prstGeom')
    fill_el = etree.fromstring(fill_xml_str)
    if prstGeom is not None:
        prstGeom.addnext(fill_el)
    else:
        spPr.append(fill_el)

# Shape 1: Semi-transparent red solid fill (alpha 50%)
s26a = slide26.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(3), Inches(2))
s26a.text = "Alpha 50%"
s26a.text_frame.paragraphs[0].font.size = Pt(18)
set_fill_xml(s26a, '''<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:srgbClr val="FF0000"><a:alpha val="50000"/></a:srgbClr>
</a:solidFill>''')

# Shape 2: Semi-transparent gradient stop
s26b = slide26.shapes.add_shape(1, Inches(4), Inches(0.5), Inches(3), Inches(2))
s26b.text = "Alpha Gradient"
s26b.text_frame.paragraphs[0].font.size = Pt(18)
s26b.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s26b, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="0000FF"><a:alpha val="80000"/></a:srgbClr></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"><a:alpha val="20000"/></a:srgbClr></a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="1"/>
</a:gradFill>''')

# ── Slide 27: Image fill on AutoShape ────────────────────────────────────────

slide27 = prs.slides.add_slide(blank)

# We need to add an image relationship and then inject a:blipFill into p:spPr
# First, add a simple shape
s27a = slide27.shapes.add_shape(1, Inches(1), Inches(1), Inches(4), Inches(3))
s27a.text = ""

# Add image relationship (use slide's existing rels) — use a small 1px PNG
# Minimal 1x1 red PNG
mini_png = base64.b64decode(
    'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8/5+hHgAHggJ/PchI7wAAAABJRU5ErkJggg=='
)
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Write temp image
tmp_img = os.path.join(tempfile.gettempdir(), '_test_blip.png')
with open(tmp_img, 'wb') as f:
    f.write(mini_png)

# Add image as a relationship on the slide's part
from pptx.opc.package import Part
img_part, img_rid = slide27.part.get_or_add_image_part(tmp_img)

# Inject a:blipFill into p:spPr
set_fill_xml(s27a, f'''<a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <a:blip r:embed="{img_rid}"/>
  <a:stretch><a:fillRect/></a:stretch>
</a:blipFill>''')

os.unlink(tmp_img)

# ── Slide 28: Pattern fill ───────────────────────────────────────────────────

slide28 = prs.slides.add_slide(blank)

# Shape 1: ltDnDiag pattern
s28a = slide28.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(2.5), Inches(2))
s28a.text = "ltDnDiag"
s28a.text_frame.paragraphs[0].font.size = Pt(14)
set_fill_xml(s28a, '''<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="ltDnDiag">
  <a:fgClr><a:srgbClr val="000000"/></a:fgClr>
  <a:bgClr><a:srgbClr val="FFFFFF"/></a:bgClr>
</a:pattFill>''')

# Shape 2: smCheck pattern
s28b = slide28.shapes.add_shape(1, Inches(3.5), Inches(0.5), Inches(2.5), Inches(2))
s28b.text = "smCheck"
s28b.text_frame.paragraphs[0].font.size = Pt(14)
set_fill_xml(s28b, '''<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="smCheck">
  <a:fgClr><a:srgbClr val="FF0000"/></a:fgClr>
  <a:bgClr><a:srgbClr val="FFFF00"/></a:bgClr>
</a:pattFill>''')

# Shape 3: dkHorz pattern
s28c = slide28.shapes.add_shape(1, Inches(6.5), Inches(0.5), Inches(2.5), Inches(2))
s28c.text = "dkHorz"
s28c.text_frame.paragraphs[0].font.size = Pt(14)
set_fill_xml(s28c, '''<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="dkHorz">
  <a:fgClr><a:srgbClr val="003366"/></a:fgClr>
  <a:bgClr><a:srgbClr val="CCCCCC"/></a:bgClr>
</a:pattFill>''')

# ── Slide 29: Gradient tileFlip ──────────────────────────────────────────────

slide29 = prs.slides.add_slide(blank)

# Shape 1: tileFlip="x"
s29a = slide29.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(2.5), Inches(2))
s29a.text = "tileFlip=x"
s29a.text_frame.paragraphs[0].font.size = Pt(14)
set_gradient_fill(s29a, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1" tileFlip="x">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FF0000"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="1"/>
</a:gradFill>''')

# Shape 2: tileFlip="y"
s29b = slide29.shapes.add_shape(1, Inches(3.5), Inches(0.5), Inches(2.5), Inches(2))
s29b.text = "tileFlip=y"
s29b.text_frame.paragraphs[0].font.size = Pt(14)
set_gradient_fill(s29b, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1" tileFlip="y">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="00FF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="FF00FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="5400000" scaled="1"/>
</a:gradFill>''')

# Shape 3: tileFlip="xy"
s29c = slide29.shapes.add_shape(1, Inches(6.5), Inches(0.5), Inches(2.5), Inches(2))
s29c.text = "tileFlip=xy"
s29c.text_frame.paragraphs[0].font.size = Pt(14)
set_gradient_fill(s29c, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1" tileFlip="xy">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FFFF00"/></a:gs>
    <a:gs pos="50000"><a:srgbClr val="FF6600"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0066FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="2700000" scaled="1"/>
</a:gradFill>''')

# ── Slide 30: Additional pattern fills ──────────────────────────────────────

slide30 = prs.slides.add_slide(blank)

patterns_30 = [
    ("pct50", "333333", "CCCCCC"),
    ("dnDiag", "000080", "FFFFFF"),
    ("cross", "FF0000", "FFFFCC"),
    ("lgCheck", "008000", "FFFFFF"),
    ("solidDmnd", "800080", "FFE0FF"),
    ("trellis", "004040", "E0FFFF"),
]

for idx, (prst, fg_c, bg_c) in enumerate(patterns_30):
    col = idx % 3
    row = idx // 3
    left = Inches(0.5 + col * 3.0)
    top = Inches(0.5 + row * 2.5)
    sh = slide30.shapes.add_shape(1, left, top, Inches(2.5), Inches(2))
    sh.text = prst
    sh.text_frame.paragraphs[0].font.size = Pt(14)
    set_fill_xml(sh, f'''<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="{prst}">
  <a:fgClr><a:srgbClr val="{fg_c}"/></a:fgClr>
  <a:bgClr><a:srgbClr val="{bg_c}"/></a:bgClr>
</a:pattFill>''')

# ── Slide 31: Image fill tile ───────────────────────────────────────────────

slide31 = prs.slides.add_slide(blank)

# Re-create temp image for tile test
tmp_img2 = os.path.join(tempfile.gettempdir(), '_test_blip_tile.png')
with open(tmp_img2, 'wb') as f:
    f.write(mini_png)

img_part2, img_rid2 = slide31.part.get_or_add_image_part(tmp_img2)

# Shape with tile fill
s31a = slide31.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(4), Inches(3))
s31a.text = "Tile fill"
s31a.text_frame.paragraphs[0].font.size = Pt(14)
set_fill_xml(s31a, f'''<a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <a:blip r:embed="{img_rid2}"/>
  <a:tile tx="0" ty="0" sx="50000" sy="50000" flip="xy" algn="tl"/>
</a:blipFill>''')

os.unlink(tmp_img2)

# ── Helper: set line XML on a shape ────────────────────────────────────────

def set_line_xml(shape, ln_xml_str):
    """Replace any existing a:ln with custom line XML in the shape's spPr."""
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    spPr = shape._element.find(f'{ns_p}spPr')
    if spPr is None:
        spPr = shape._element.find(f'.//{ns_a}spPr')
    if spPr is None:
        return
    for el in spPr.findall(f'{ns_a}ln'):
        spPr.remove(el)
    ln_el = etree.fromstring(ln_xml_str)
    spPr.append(ln_el)

# ── Helper: create a line shape (p:sp with a:prstGeom prst="line") ─────────

def add_line_shape(slide, left, top, width, height, ln_xml_str):
    """Add a p:sp with line geometry + custom a:ln XML."""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    # Create via add_shape(1=rect) then change geometry to line
    s = slide.shapes.add_shape(1, left, top, width, height)
    spPr = s._element.find(f'{{{ns_p}}}spPr')
    if spPr is None:
        spPr = s._element.find(f'.//{{{ns_a}}}spPr')
    # Replace prstGeom with line
    for pg in spPr.findall(f'{{{ns_a}}}prstGeom'):
        spPr.remove(pg)
    pg = etree.SubElement(spPr, f'{{{ns_a}}}prstGeom')
    pg.set('prst', 'line')
    etree.SubElement(pg, f'{{{ns_a}}}avLst')
    # Remove fill (lines don't have fill)
    for tag in ('solidFill', 'noFill', 'gradFill'):
        for el in spPr.findall(f'{{{ns_a}}}{tag}'):
            spPr.remove(el)
    noFill = etree.SubElement(spPr, f'{{{ns_a}}}noFill')
    # Set line XML
    for el in spPr.findall(f'{{{ns_a}}}ln'):
        spPr.remove(el)
    ln_el = etree.fromstring(ln_xml_str)
    spPr.append(ln_el)
    return s

# ── Slide 32: Stroke styles on lines ─────────────────────────────────────

slide32 = prs.slides.add_slide(blank)

title32 = slide32.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
title32.text_frame.paragraphs[0].text = "Slide 32: Line Stroke Styles"
title32.text_frame.paragraphs[0].font.size = Pt(20)
title32.text_frame.paragraphs[0].font.bold = True

# Dash styles on lines
dash_info = [
    ('dash', '0000FF'),
    ('dot', '0066CC'),
    ('dashDot', '009900'),
    ('lgDash', 'CC6600'),
    ('sysDot', '990099'),
    ('sysDash', 'CC0000'),
]
for i, (dash, color) in enumerate(dash_info):
    y = 1.0 + i * 0.9
    add_line_shape(slide32, Inches(1), Inches(y), Inches(8), Emu(0),
        f'''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
          <a:prstDash val="{dash}"/>
        </a:ln>''')
    # Label
    lbl = slide32.shapes.add_textbox(Inches(0.2), Inches(y - 0.25), Inches(0.8), Inches(0.3))
    lbl.text_frame.paragraphs[0].text = dash
    lbl.text_frame.paragraphs[0].font.size = Pt(9)
    lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Custom dash (a:custDash with a:ds elements)
add_line_shape(slide32, Inches(1), Inches(1.0 + 6 * 0.9), Inches(8), Emu(0),
    '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="FF6600"/></a:solidFill>
      <a:custDash>
        <a:ds d="400000" sp="100000"/>
        <a:ds d="100000" sp="100000"/>
      </a:custDash>
    </a:ln>''')
lbl_cd = slide32.shapes.add_textbox(Inches(0.2), Inches(1.0 + 6 * 0.9 - 0.25), Inches(0.8), Inches(0.3))
lbl_cd.text_frame.paragraphs[0].text = "custDash"
lbl_cd.text_frame.paragraphs[0].font.size = Pt(9)
lbl_cd.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide 33: Arrows + Line cap/join ──────────────────────────────────────

slide33 = prs.slides.add_slide(blank)

title33 = slide33.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
title33.text_frame.paragraphs[0].text = "Slide 33: Arrow Lines"
title33.text_frame.paragraphs[0].font.size = Pt(20)
title33.text_frame.paragraphs[0].font.bold = True

# Line with triangle head + stealth tail
add_line_shape(slide33, Inches(0.5), Inches(1.2), Inches(4), Emu(0),
    '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>
      <a:headEnd type="triangle" w="med" len="med"/>
      <a:tailEnd type="stealth" w="lg" len="lg"/>
    </a:ln>''')
lbl33a = slide33.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4), Inches(0.3))
lbl33a.text_frame.paragraphs[0].text = "triangle head + stealth tail"
lbl33a.text_frame.paragraphs[0].font.size = Pt(9)
lbl33a.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Line with diamond head + oval tail
add_line_shape(slide33, Inches(0.5), Inches(2.0), Inches(4), Emu(0),
    '''<a:ln w="19050" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="008000"/></a:solidFill>
      <a:headEnd type="diamond" w="med" len="med"/>
      <a:tailEnd type="oval" w="med" len="med"/>
    </a:ln>''')
lbl33b = slide33.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4), Inches(0.3))
lbl33b.text_frame.paragraphs[0].text = "diamond head + oval tail"
lbl33b.text_frame.paragraphs[0].font.size = Pt(9)
lbl33b.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Line with arrow (open) head + triangle tail
add_line_shape(slide33, Inches(0.5), Inches(2.8), Inches(4), Emu(0),
    '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="0000FF"/></a:solidFill>
      <a:headEnd type="arrow" w="med" len="med"/>
      <a:tailEnd type="triangle" w="sm" len="sm"/>
    </a:ln>''')
lbl33c = slide33.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(4), Inches(0.3))
lbl33c.text_frame.paragraphs[0].text = "arrow head + triangle tail (sm)"
lbl33c.text_frame.paragraphs[0].font.size = Pt(9)
lbl33c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Dashed line with round cap + round join
add_line_shape(slide33, Inches(5), Inches(1.2), Inches(4.5), Emu(0),
    '''<a:ln w="38100" cap="rnd" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="800080"/></a:solidFill>
      <a:prstDash val="dash"/>
      <a:round/>
    </a:ln>''')
lbl33d = slide33.shapes.add_textbox(Inches(5), Inches(1.3), Inches(4.5), Inches(0.3))
lbl33d.text_frame.paragraphs[0].text = "dash + round cap + round join"
lbl33d.text_frame.paragraphs[0].font.size = Pt(9)
lbl33d.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Line with square cap + bevel join + lgDash
add_line_shape(slide33, Inches(5), Inches(2.0), Inches(4.5), Emu(0),
    '''<a:ln w="25400" cap="sq" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="006699"/></a:solidFill>
      <a:prstDash val="lgDash"/>
      <a:bevel/>
    </a:ln>''')
lbl33e = slide33.shapes.add_textbox(Inches(5), Inches(2.1), Inches(4.5), Inches(0.3))
lbl33e.text_frame.paragraphs[0].text = "lgDash + sq cap + bevel join"
lbl33e.text_frame.paragraphs[0].font.size = Pt(9)
lbl33e.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Miter join with limit
add_line_shape(slide33, Inches(5), Inches(2.8), Inches(4.5), Emu(0),
    '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="FF6600"/></a:solidFill>
      <a:miter lim="800000"/>
    </a:ln>''')
lbl33f = slide33.shapes.add_textbox(Inches(5), Inches(2.9), Inches(4.5), Inches(0.3))
lbl33f.text_frame.paragraphs[0].text = "miter join (lim=800000)"
lbl33f.text_frame.paragraphs[0].font.size = Pt(9)
lbl33f.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Dashed rect (stroke on rect shape) — keeps one rect for dash+cap test
s33rect = slide33.shapes.add_shape(1, Inches(0.5), Inches(4.0), Inches(4), Inches(2))
s33rect.text = "Rect: dashDot + rnd cap"
s33rect.text_frame.paragraphs[0].font.size = Pt(12)
set_line_xml(s33rect, '''<a:ln w="25400" cap="rnd" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:solidFill><a:srgbClr val="FF00FF"/></a:solidFill>
  <a:prstDash val="dashDot"/>
  <a:round/>
</a:ln>''')

# Compound line: double
s33cmpd = slide33.shapes.add_shape(1, Inches(5), Inches(4.0), Inches(4.5), Inches(0.8))
s33cmpd.text = "Compound: dbl"
s33cmpd.text_frame.paragraphs[0].font.size = Pt(12)
set_line_xml(s33cmpd, '''<a:ln w="38100" cmpd="dbl" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
</a:ln>''')

# noFill line (stroke_no_fill)
s33nf = slide33.shapes.add_shape(1, Inches(5), Inches(5.2), Inches(4.5), Inches(0.8))
s33nf.text = "Line noFill"
s33nf.text_frame.paragraphs[0].font.size = Pt(12)
set_line_xml(s33nf, '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:noFill/>
</a:ln>''')

