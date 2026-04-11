"""Slides 52-73: standard charts (bar/line/pie/donut/scatter/area/radar/bubble/stock/surface/ofPie/3D/trendline/errBars/composite/stacked).

Mechanically carved from the legacy gen_test_features.py. Runs its
slide-building code at import time against the shared '_ctx.prs'.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
import base64, io, os, re, struct, tempfile, zipfile, zlib

from ._ctx import prs, blank

# ── Slide 52: Bar/Column chart ────────────────────────────────────────────────

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

slide52 = prs.slides.add_slide(blank)

chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Sales 2024', (120, 150, 180, 200))
chart_data.add_series('Sales 2025', (100, 130, 160, 190))

chart_frame = slide52.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(5), Inches(4), chart_data
)
chart52 = chart_frame.chart
chart52.has_legend = True

lbl52 = slide52.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl52.text_frame.paragraphs[0].text = "Slide 52: Column chart (clustered)"
lbl52.text_frame.paragraphs[0].font.size = Pt(18)
lbl52.text_frame.paragraphs[0].font.bold = True

# ── Slide 53: Line chart ─────────────────────────────────────────────────────

slide53 = prs.slides.add_slide(blank)

chart_data2 = CategoryChartData()
chart_data2.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May']
chart_data2.add_series('Revenue', (80, 95, 110, 105, 130))
chart_data2.add_series('Cost', (60, 65, 70, 75, 80))

chart_frame2 = slide53.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(1), Inches(1.5),
    Inches(5), Inches(4), chart_data2
)
chart53 = chart_frame2.chart
chart53.has_legend = True

lbl53 = slide53.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl53.text_frame.paragraphs[0].text = "Slide 53: Line chart"
lbl53.text_frame.paragraphs[0].font.size = Pt(18)
lbl53.text_frame.paragraphs[0].font.bold = True

# ── Slide 54: Pie chart ──────────────────────────────────────────────────────

slide54 = prs.slides.add_slide(blank)

chart_data3 = CategoryChartData()
chart_data3.categories = ['Desktop', 'Mobile', 'Tablet', 'Other']
chart_data3.add_series('Market Share', (45, 35, 15, 5))

chart_frame3 = slide54.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(2), Inches(1.5),
    Inches(5), Inches(4), chart_data3
)
chart54 = chart_frame3.chart
chart54.has_legend = True

lbl54 = slide54.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl54.text_frame.paragraphs[0].text = "Slide 54: Pie chart"
lbl54.text_frame.paragraphs[0].font.size = Pt(18)
lbl54.text_frame.paragraphs[0].font.bold = True

# ── Slide 55: Bar chart (horizontal) + Donut chart ──────────────────────────

slide55 = prs.slides.add_slide(blank)

chart_data4 = CategoryChartData()
chart_data4.categories = ['Product A', 'Product B', 'Product C']
chart_data4.add_series('Units Sold', (250, 180, 320))

chart_frame4 = slide55.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.5),
    Inches(4), Inches(4), chart_data4
)

chart_data5 = CategoryChartData()
chart_data5.categories = ['Yes', 'No', 'Maybe']
chart_data5.add_series('Responses', (60, 25, 15))

chart_frame5 = slide55.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, Inches(5), Inches(1.5),
    Inches(4), Inches(4), chart_data5
)

lbl55 = slide55.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl55.text_frame.paragraphs[0].text = "Slide 55: Bar chart (horizontal) + Donut chart"
lbl55.text_frame.paragraphs[0].font.size = Pt(18)
lbl55.text_frame.paragraphs[0].font.bold = True

# ── Slide 56: Column chart with data labels ──────────────────────────────────

slide56 = prs.slides.add_slide(blank)

chart_data6 = CategoryChartData()
chart_data6.categories = ['North', 'South', 'East', 'West']
chart_data6.add_series('Revenue', (320, 280, 190, 410))

chart_frame6 = slide56.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(5), Inches(4), chart_data6
)
chart56 = chart_frame6.chart
chart56.has_legend = True

# Enable data labels with values
from pptx.util import Emu as _Emu
plot = chart56.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.show_value = True
data_labels.show_category_name = False

lbl56 = slide56.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl56.text_frame.paragraphs[0].text = "Slide 56: Column chart with data labels"
lbl56.text_frame.paragraphs[0].font.size = Pt(18)
lbl56.text_frame.paragraphs[0].font.bold = True

# ── Slide 57: Pie chart with data point colors + percentage labels ────────────

slide57 = prs.slides.add_slide(blank)

chart_data7 = CategoryChartData()
chart_data7.categories = ['Chrome', 'Firefox', 'Safari', 'Edge', 'Other']
chart_data7.add_series('Browser Share', (65, 10, 18, 5, 2))

chart_frame7 = slide57.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(2), Inches(1.5),
    Inches(5), Inches(4), chart_data7
)
chart57 = chart_frame7.chart
chart57.has_legend = True

# Set per-slice colors via c:dPt
# Access the chart part XML directly
chart_part57 = chart57.part
chart_xml57 = chart_part57._element
nsmap = {
    'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}
pie_chart = chart_xml57.findall('.//c:pieChart', nsmap)[0]
ser_elem = pie_chart.findall('c:ser', nsmap)[0]

# Add dPt elements with custom colors
dpt_colors = [
    (0, '4285F4'),  # Chrome blue
    (1, 'FF7139'),  # Firefox orange
    (2, '000000'),  # Safari black
    (3, '0078D4'),  # Edge blue
    (4, '888888'),  # Other gray
]
c_ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
for pt_idx, color_hex in dpt_colors:
    dpt = etree.SubElement(ser_elem, f'{{{c_ns}}}dPt')
    idx_elem = etree.SubElement(dpt, f'{{{c_ns}}}idx')
    idx_elem.set('val', str(pt_idx))
    spPr = etree.SubElement(dpt, f'{{{c_ns}}}spPr')
    solidFill = etree.SubElement(spPr, f'{{{a_ns}}}solidFill')
    srgb = etree.SubElement(solidFill, f'{{{a_ns}}}srgbClr')
    srgb.set('val', color_hex)

# Enable percentage labels
plot57 = chart57.plots[0]
plot57.has_data_labels = True
dl57 = plot57.data_labels
dl57.show_percentage = True
dl57.show_value = False
dl57.show_category_name = True

lbl57 = slide57.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl57.text_frame.paragraphs[0].text = "Slide 57: Pie chart with custom colors + % labels"
lbl57.text_frame.paragraphs[0].font.size = Pt(18)
lbl57.text_frame.paragraphs[0].font.bold = True

# ── Slide 58: Line chart with series spPr colors + data labels ────────────────

slide58 = prs.slides.add_slide(blank)

chart_data8 = CategoryChartData()
chart_data8.categories = ['Week 1', 'Week 2', 'Week 3', 'Week 4']
chart_data8.add_series('Actual', (10, 25, 40, 55))
chart_data8.add_series('Target', (15, 30, 45, 60))

chart_frame8 = slide58.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(1), Inches(1.5),
    Inches(5), Inches(4), chart_data8
)
chart58 = chart_frame8.chart
chart58.has_legend = True

# Set explicit series colors
from pptx.chart.series import LineSeries
for i, ser in enumerate(chart58.series):
    ser_format = ser.format
    if i == 0:
        ser_format.line.color.rgb = RGBColor(0x00, 0x88, 0x00)  # Green
    else:
        ser_format.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)  # Red

# Enable data labels showing values
plot58 = chart58.plots[0]
plot58.has_data_labels = True
dl58 = plot58.data_labels
dl58.show_value = True

lbl58 = slide58.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl58.text_frame.paragraphs[0].text = "Slide 58: Line chart with series colors + data labels"
lbl58.text_frame.paragraphs[0].font.size = Pt(18)
lbl58.text_frame.paragraphs[0].font.bold = True

# ── Slide 59: Scatter chart + Area chart ─────────────────────────────────────

from pptx.chart.data import XyChartData

slide59 = prs.slides.add_slide(blank)

xy_data = XyChartData()
s1 = xy_data.add_series('Group A')
s1.add_data_point(10, 20)
s1.add_data_point(30, 50)
s1.add_data_point(50, 40)
s1.add_data_point(70, 80)
s2 = xy_data.add_series('Group B')
s2.add_data_point(15, 60)
s2.add_data_point(35, 30)
s2.add_data_point(55, 70)
s2.add_data_point(75, 45)

chart_frame59a = slide59.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER, Inches(0.3), Inches(1.5),
    Inches(4.5), Inches(4), xy_data
)
chart59a = chart_frame59a.chart
chart59a.has_legend = True

area_data = CategoryChartData()
area_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May']
area_data.add_series('Downloads', (120, 180, 150, 200, 250))
area_data.add_series('Uploads', (80, 100, 130, 110, 170))

chart_frame59b = slide59.shapes.add_chart(
    XL_CHART_TYPE.AREA, Inches(5.2), Inches(1.5),
    Inches(4.5), Inches(4), area_data
)
chart59b = chart_frame59b.chart
chart59b.has_legend = True

lbl59 = slide59.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl59.text_frame.paragraphs[0].text = "Slide 59: Scatter chart + Area chart"
lbl59.text_frame.paragraphs[0].font.size = Pt(18)
lbl59.text_frame.paragraphs[0].font.bold = True

# ── Slide 60: Radar chart ────────────────────────────────────────────────────

slide60 = prs.slides.add_slide(blank)

radar_data = CategoryChartData()
radar_data.categories = ['Str', 'Dex', 'Con', 'Int', 'Wis', 'Cha']
radar_data.add_series('Character A', (15, 12, 14, 10, 8, 16))
radar_data.add_series('Character B', (10, 16, 8, 15, 14, 12))

chart_frame60 = slide60.shapes.add_chart(
    XL_CHART_TYPE.RADAR, Inches(2), Inches(1.5),
    Inches(5), Inches(4.5), radar_data
)
chart60 = chart_frame60.chart
chart60.has_legend = True

lbl60 = slide60.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl60.text_frame.paragraphs[0].text = "Slide 60: Radar chart"
lbl60.text_frame.paragraphs[0].font.size = Pt(18)
lbl60.text_frame.paragraphs[0].font.bold = True

# ── Slide 61: Bubble chart ──────────────────────────────────────────────────

from pptx.chart.data import BubbleChartData

slide61 = prs.slides.add_slide(blank)

bubble_data = BubbleChartData()
bs1 = bubble_data.add_series('Product X')
bs1.add_data_point(10, 40, 15)
bs1.add_data_point(25, 60, 30)
bs1.add_data_point(40, 30, 10)
bs1.add_data_point(55, 70, 25)
bs2 = bubble_data.add_series('Product Y')
bs2.add_data_point(15, 55, 20)
bs2.add_data_point(30, 35, 12)
bs2.add_data_point(50, 50, 35)

chart_frame61 = slide61.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE, Inches(1.5), Inches(1.5),
    Inches(6), Inches(4.5), bubble_data
)
chart61 = chart_frame61.chart
chart61.has_legend = True

lbl61 = slide61.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl61.text_frame.paragraphs[0].text = "Slide 61: Bubble chart"
lbl61.text_frame.paragraphs[0].font.size = Pt(18)
lbl61.text_frame.paragraphs[0].font.bold = True

# ── Slide 62: Stock chart (OHLC) ───────────────────────────────────────────

slide62 = prs.slides.add_slide(blank)

# python-pptx doesn't have direct stock chart support, so inject via lxml
# Create a bar chart as placeholder, then replace with stock chart XML
stock_placeholder = CategoryChartData()
stock_placeholder.categories = ['Day 1', 'Day 2', 'Day 3', 'Day 4', 'Day 5']
stock_placeholder.add_series('Open', (100, 105, 98, 110, 108))

chart_frame62 = slide62.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), stock_placeholder
)
chart62 = chart_frame62.chart

# Replace chart XML with a stock chart
chart_part62 = chart62.part
chart_elem62 = chart_part62._element
ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Find plotArea and remove existing barChart
plot_area = chart_elem62.find(f'.//{{{ns_c}}}plotArea')
for bar in plot_area.findall(f'{{{ns_c}}}barChart'):
    plot_area.remove(bar)

# Build stock chart element
stock_xml = f"""<c:stockChart xmlns:c="{ns_c}" xmlns:a="{ns_a}">
  <c:ser>
    <c:idx val="0"/><c:order val="0"/>
    <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>Open</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$6</c:f><c:strCache><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>Day 1</c:v></c:pt><c:pt idx="1"><c:v>Day 2</c:v></c:pt>
      <c:pt idx="2"><c:v>Day 3</c:v></c:pt><c:pt idx="3"><c:v>Day 4</c:v></c:pt>
      <c:pt idx="4"><c:v>Day 5</c:v></c:pt></c:strCache></c:strRef></c:cat>
    <c:val><c:numRef><c:f>Sheet1!$B$2:$B$6</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="5"/>
      <c:pt idx="0"><c:v>100</c:v></c:pt><c:pt idx="1"><c:v>105</c:v></c:pt>
      <c:pt idx="2"><c:v>98</c:v></c:pt><c:pt idx="3"><c:v>110</c:v></c:pt>
      <c:pt idx="4"><c:v>108</c:v></c:pt></c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="1"/><c:order val="1"/>
    <c:tx><c:strRef><c:f>Sheet1!$C$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>High</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$C$2:$C$6</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="5"/>
      <c:pt idx="0"><c:v>115</c:v></c:pt><c:pt idx="1"><c:v>118</c:v></c:pt>
      <c:pt idx="2"><c:v>112</c:v></c:pt><c:pt idx="3"><c:v>120</c:v></c:pt>
      <c:pt idx="4"><c:v>122</c:v></c:pt></c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="2"/><c:order val="2"/>
    <c:tx><c:strRef><c:f>Sheet1!$D$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>Low</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$D$2:$D$6</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="5"/>
      <c:pt idx="0"><c:v>95</c:v></c:pt><c:pt idx="1"><c:v>100</c:v></c:pt>
      <c:pt idx="2"><c:v>90</c:v></c:pt><c:pt idx="3"><c:v>105</c:v></c:pt>
      <c:pt idx="4"><c:v>102</c:v></c:pt></c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="3"/><c:order val="3"/>
    <c:tx><c:strRef><c:f>Sheet1!$E$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>Close</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$E$2:$E$6</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="5"/>
      <c:pt idx="0"><c:v>108</c:v></c:pt><c:pt idx="1"><c:v>102</c:v></c:pt>
      <c:pt idx="2"><c:v>107</c:v></c:pt><c:pt idx="3"><c:v>115</c:v></c:pt>
      <c:pt idx="4"><c:v>118</c:v></c:pt></c:numCache></c:numRef></c:val>
  </c:ser>
  <c:axId val="10000001"/>
  <c:axId val="10000002"/>
</c:stockChart>"""

stock_elem = etree.fromstring(stock_xml)

# Find existing axes and remove to add stock-compatible ones
for ax in plot_area.findall(f'{{{ns_c}}}catAx') + plot_area.findall(f'{{{ns_c}}}valAx'):
    plot_area.remove(ax)

# Insert stock chart after any layout element
layout_elem = plot_area.find(f'{{{ns_c}}}layout')
if layout_elem is not None:
    idx_pos = list(plot_area).index(layout_elem) + 1
    plot_area.insert(idx_pos, stock_elem)
else:
    plot_area.insert(0, stock_elem)

# Add catAx and valAx for stock chart
cat_ax_xml = f"""<c:catAx xmlns:c="{ns_c}">
  <c:axId val="10000001"/>
  <c:scaling><c:orientation val="minMax"/></c:scaling>
  <c:delete val="0"/>
  <c:axPos val="b"/>
  <c:crossAx val="10000002"/>
</c:catAx>"""
val_ax_xml = f"""<c:valAx xmlns:c="{ns_c}">
  <c:axId val="10000002"/>
  <c:scaling><c:orientation val="minMax"/></c:scaling>
  <c:delete val="0"/>
  <c:axPos val="l"/>
  <c:crossAx val="10000001"/>
  <c:majorGridlines/>
</c:valAx>"""
plot_area.append(etree.fromstring(cat_ax_xml))
plot_area.append(etree.fromstring(val_ax_xml))

# Set legend
chart_node62 = chart_elem62.find(f'.//{{{ns_c}}}chart')
legend_xml = f"""<c:legend xmlns:c="{ns_c}">
  <c:legendPos val="r"/>
  <c:overlay val="0"/>
</c:legend>"""
chart_node62.append(etree.fromstring(legend_xml))

lbl62 = slide62.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl62.text_frame.paragraphs[0].text = "Slide 62: Stock chart (OHLC)"
lbl62.text_frame.paragraphs[0].font.size = Pt(18)
lbl62.text_frame.paragraphs[0].font.bold = True

# ── Slide 63: All chart types overview ──────────────────────────────────────

slide63 = prs.slides.add_slide(blank)

# Mini column chart (top-left)
mini_col = CategoryChartData()
mini_col.categories = ['A', 'B', 'C']
mini_col.add_series('S1', (30, 50, 40))

slide63.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.3), Inches(1.2),
    Inches(3), Inches(2.5), mini_col
)

# Mini line chart (top-center)
mini_line = CategoryChartData()
mini_line.categories = ['A', 'B', 'C', 'D']
mini_line.add_series('S1', (10, 30, 20, 40))

slide63.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(3.5), Inches(1.2),
    Inches(3), Inches(2.5), mini_line
)

# Mini pie chart (top-right)
mini_pie = CategoryChartData()
mini_pie.categories = ['X', 'Y', 'Z']
mini_pie.add_series('S1', (40, 35, 25))

slide63.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(6.7), Inches(1.2),
    Inches(3), Inches(2.5), mini_pie
)

# Mini bar chart (bottom-left)
mini_bar = CategoryChartData()
mini_bar.categories = ['P', 'Q', 'R']
mini_bar.add_series('S1', (25, 40, 35))

slide63.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.3), Inches(4.0),
    Inches(3), Inches(2.5), mini_bar
)

# Mini scatter chart (bottom-center)
mini_scatter = XyChartData()
ms1 = mini_scatter.add_series('S1')
ms1.add_data_point(5, 10)
ms1.add_data_point(15, 30)
ms1.add_data_point(25, 20)
ms1.add_data_point(35, 40)

slide63.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER, Inches(3.5), Inches(4.0),
    Inches(3), Inches(2.5), mini_scatter
)

# Mini donut chart (bottom-right)
mini_donut = CategoryChartData()
mini_donut.categories = ['A', 'B', 'C']
mini_donut.add_series('S1', (50, 30, 20))

slide63.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, Inches(6.7), Inches(4.0),
    Inches(3), Inches(2.5), mini_donut
)

lbl63 = slide63.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl63.text_frame.paragraphs[0].text = "Slide 63: All chart types overview"
lbl63.text_frame.paragraphs[0].font.size = Pt(18)
lbl63.text_frame.paragraphs[0].font.bold = True

# ── Slide 64: Line chart with linear trendline ──────────────────────────────

slide64 = prs.slides.add_slide(blank)

trend_data = CategoryChartData()
trend_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
trend_data.add_series('Sales', (20, 35, 28, 45, 52, 60))

chart_frame64 = slide64.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), trend_data
)
chart64 = chart_frame64.chart
chart64.has_legend = True

# Add trendline via lxml (python-pptx has limited trendline support)
chart_part64 = chart64.part
chart_elem64 = chart_part64._element
ns_c64 = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
ns_a64 = 'http://schemas.openxmlformats.org/drawingml/2006/main'

line_chart64 = chart_elem64.findall(f'.//{{{ns_c64}}}lineChart')[0]
ser64 = line_chart64.findall(f'{{{ns_c64}}}ser')[0]

trendline_xml = f"""<c:trendline xmlns:c="{ns_c64}" xmlns:a="{ns_a64}">
  <c:trendlineType val="linear"/>
  <c:spPr>
    <a:ln w="12700">
      <a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>
      <a:prstDash val="dash"/>
    </a:ln>
  </c:spPr>
</c:trendline>"""
ser64.append(etree.fromstring(trendline_xml))

lbl64 = slide64.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl64.text_frame.paragraphs[0].text = "Slide 64: Line chart with linear trendline"
lbl64.text_frame.paragraphs[0].font.size = Pt(18)
lbl64.text_frame.paragraphs[0].font.bold = True

# ── Slide 65: Column chart with error bars ──────────────────────────────────

slide65 = prs.slides.add_slide(blank)

err_data = CategoryChartData()
err_data.categories = ['Group A', 'Group B', 'Group C', 'Group D']
err_data.add_series('Measurement', (85, 92, 78, 95))

chart_frame65 = slide65.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), err_data
)
chart65 = chart_frame65.chart
chart65.has_legend = True

# Add error bars via lxml
chart_part65 = chart65.part
chart_elem65 = chart_part65._element
ns_c65 = ns_c64

bar_chart65 = chart_elem65.findall(f'.//{{{ns_c65}}}barChart')[0]
ser65 = bar_chart65.findall(f'{{{ns_c65}}}ser')[0]

errbar_xml = f"""<c:errBars xmlns:c="{ns_c65}">
  <c:errDir val="y"/>
  <c:errBarType val="both"/>
  <c:errValType val="fixedVal"/>
  <c:val val="8"/>
</c:errBars>"""
ser65.append(etree.fromstring(errbar_xml))

lbl65 = slide65.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl65.text_frame.paragraphs[0].text = "Slide 65: Column chart with error bars (±8)"
lbl65.text_frame.paragraphs[0].font.size = Pt(18)
lbl65.text_frame.paragraphs[0].font.bold = True

# ── Slide 66: Composite chart (column + line on same plot) ───────────────────

slide66 = prs.slides.add_slide(blank)

# Create a column chart first
combo_data = CategoryChartData()
combo_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
combo_data.add_series('Revenue', (300, 350, 400, 450))

chart_frame66 = slide66.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), combo_data
)
chart66 = chart_frame66.chart
chart66.has_legend = True

# Add a lineChart group to the plotArea via lxml
chart_part66 = chart66.part
chart_elem66 = chart_part66._element
ns_c66 = ns_c64

plot_area66 = chart_elem66.find(f'.//{{{ns_c66}}}plotArea')

# Find existing barChart's axId values
bar_chart66 = plot_area66.find(f'{{{ns_c66}}}barChart')
ax_ids66 = bar_chart66.findall(f'{{{ns_c66}}}axId')
ax_id_vals = [a.get('val') for a in ax_ids66]

line_chart_xml = f"""<c:lineChart xmlns:c="{ns_c66}" xmlns:a="{ns_a64}">
  <c:grouping val="standard"/>
  <c:ser>
    <c:idx val="1"/><c:order val="1"/>
    <c:tx><c:strRef><c:f>Sheet1!$C$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>Profit</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:spPr><a:ln w="25400"><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill></a:ln></c:spPr>
    <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$5</c:f><c:strCache><c:ptCount val="4"/>
      <c:pt idx="0"><c:v>Q1</c:v></c:pt><c:pt idx="1"><c:v>Q2</c:v></c:pt>
      <c:pt idx="2"><c:v>Q3</c:v></c:pt><c:pt idx="3"><c:v>Q4</c:v></c:pt>
    </c:strCache></c:strRef></c:cat>
    <c:val><c:numRef><c:f>Sheet1!$C$2:$C$5</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="4"/>
      <c:pt idx="0"><c:v>80</c:v></c:pt><c:pt idx="1"><c:v>120</c:v></c:pt>
      <c:pt idx="2"><c:v>150</c:v></c:pt><c:pt idx="3"><c:v>200</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:axId val="{ax_id_vals[0] if ax_id_vals else '10000001'}"/>
  <c:axId val="{ax_id_vals[1] if len(ax_id_vals) > 1 else '10000002'}"/>
</c:lineChart>"""

# Insert lineChart after barChart
bar_idx = list(plot_area66).index(bar_chart66)
plot_area66.insert(bar_idx + 1, etree.fromstring(line_chart_xml))

lbl66 = slide66.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl66.text_frame.paragraphs[0].text = "Slide 66: Composite chart (column + line)"
lbl66.text_frame.paragraphs[0].font.size = Pt(18)
lbl66.text_frame.paragraphs[0].font.bold = True

# ── Slide 67: Surface chart (2D heatmap approximation) ─────────────────────
slide67 = prs.slides.add_slide(blank)

# Create a bar chart as placeholder, then replace with surface chart XML
surf_placeholder = CategoryChartData()
surf_placeholder.categories = ['X1', 'X2', 'X3', 'X4', 'X5']
surf_placeholder.add_series('Y1', (10, 20, 30, 40, 50))

chart_frame67 = slide67.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), surf_placeholder
)
chart67 = chart_frame67.chart
chart_part67 = chart67.part
chart_elem67 = chart_part67._element
ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Find plotArea and remove existing barChart
plot_area67 = chart_elem67.find(f'.//{{{ns_c}}}plotArea')
for bar in plot_area67.findall(f'{{{ns_c}}}barChart'):
    plot_area67.remove(bar)

# Build surface chart with 4 series × 5 data points (grid)
surface_xml = f"""<c:surfaceChart xmlns:c="{ns_c}" xmlns:a="{ns_a}">
  <c:wireframe val="0"/>
  <c:ser>
    <c:idx val="0"/><c:order val="0"/>
    <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Row 1</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$6</c:f><c:strCache><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>X1</c:v></c:pt><c:pt idx="1"><c:v>X2</c:v></c:pt><c:pt idx="2"><c:v>X3</c:v></c:pt><c:pt idx="3"><c:v>X4</c:v></c:pt><c:pt idx="4"><c:v>X5</c:v></c:pt>
    </c:strCache></c:strRef></c:cat>
    <c:val><c:numRef><c:f>Sheet1!$B$2:$B$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>10</c:v></c:pt><c:pt idx="1"><c:v>25</c:v></c:pt><c:pt idx="2"><c:v>40</c:v></c:pt><c:pt idx="3"><c:v>30</c:v></c:pt><c:pt idx="4"><c:v>15</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="1"/><c:order val="1"/>
    <c:tx><c:strRef><c:f>Sheet1!$C$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Row 2</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$C$2:$C$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>20</c:v></c:pt><c:pt idx="1"><c:v>35</c:v></c:pt><c:pt idx="2"><c:v>50</c:v></c:pt><c:pt idx="3"><c:v>45</c:v></c:pt><c:pt idx="4"><c:v>30</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="2"/><c:order val="2"/>
    <c:tx><c:strRef><c:f>Sheet1!$D$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Row 3</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$D$2:$D$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>35</c:v></c:pt><c:pt idx="1"><c:v>50</c:v></c:pt><c:pt idx="2"><c:v>60</c:v></c:pt><c:pt idx="3"><c:v>55</c:v></c:pt><c:pt idx="4"><c:v>40</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="3"/><c:order val="3"/>
    <c:tx><c:strRef><c:f>Sheet1!$E$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Row 4</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$E$2:$E$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>25</c:v></c:pt><c:pt idx="1"><c:v>40</c:v></c:pt><c:pt idx="2"><c:v>45</c:v></c:pt><c:pt idx="3"><c:v>50</c:v></c:pt><c:pt idx="4"><c:v>35</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:bandFmts/>
  <c:axId val="111111"/>
  <c:axId val="222222"/>
  <c:axId val="333333"/>
</c:surfaceChart>"""

# Remove existing axes
for ax_tag in ['catAx', 'valAx', 'serAx']:
    for ax in plot_area67.findall(f'{{{ns_c}}}{ax_tag}'):
        plot_area67.remove(ax)

# Insert surface chart
plot_area67.append(etree.fromstring(surface_xml))

# Add axes for surface chart
surf_axes = f"""<c:catAx xmlns:c="{ns_c}"><c:axId val="111111"/><c:scaling><c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="b"/><c:crossAx val="222222"/></c:catAx>"""
plot_area67.append(etree.fromstring(surf_axes))
surf_val_ax = f"""<c:valAx xmlns:c="{ns_c}"><c:axId val="222222"/><c:scaling><c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="l"/><c:crossAx val="111111"/><c:majorGridlines/></c:valAx>"""
plot_area67.append(etree.fromstring(surf_val_ax))
# serAx for surface chart
surf_ser_ax = f"""<c:serAx xmlns:c="{ns_c}"><c:axId val="333333"/><c:scaling><c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="b"/><c:crossAx val="222222"/></c:serAx>"""
plot_area67.append(etree.fromstring(surf_ser_ax))

# Add view3D to chart element
chart_node67 = chart_elem67.find(f'{{{ns_c}}}chart')
view3d_xml = f"""<c:view3D xmlns:c="{ns_c}"><c:rotX val="15"/><c:rotY val="20"/><c:depthPercent val="100"/><c:rAngAx val="1"/><c:perspective val="30"/></c:view3D>"""
chart_node67.insert(0, etree.fromstring(view3d_xml))

lbl67 = slide67.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl67.text_frame.paragraphs[0].text = "Slide 67: Surface chart (2D heatmap)"
lbl67.text_frame.paragraphs[0].font.size = Pt(18)
lbl67.text_frame.paragraphs[0].font.bold = True

# ── Slide 68: Pie of pie chart (ofPieChart) ─────────────────────────────────
slide68 = prs.slides.add_slide(blank)

# Create a pie chart as placeholder, then replace with ofPieChart XML
ofpie_placeholder = CategoryChartData()
ofpie_placeholder.categories = ['Product A', 'Product B', 'Product C', 'Product D', 'Product E']
ofpie_placeholder.add_series('Sales', (40, 25, 15, 12, 8))

chart_frame68 = slide68.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), ofpie_placeholder
)
chart68 = chart_frame68.chart
chart_part68 = chart68.part
chart_elem68 = chart_part68._element

# Find plotArea and remove existing pieChart
plot_area68 = chart_elem68.find(f'.//{{{ns_c}}}plotArea')
for pie in plot_area68.findall(f'{{{ns_c}}}pieChart'):
    plot_area68.remove(pie)

# Build ofPieChart element
ofpie_xml = f"""<c:ofPieChart xmlns:c="{ns_c}" xmlns:a="{ns_a}">
  <c:ofPieType val="pie"/>
  <c:varyColors val="1"/>
  <c:ser>
    <c:idx val="0"/><c:order val="0"/>
    <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Sales</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$6</c:f><c:strCache><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>Product A</c:v></c:pt><c:pt idx="1"><c:v>Product B</c:v></c:pt>
      <c:pt idx="2"><c:v>Product C</c:v></c:pt><c:pt idx="3"><c:v>Product D</c:v></c:pt>
      <c:pt idx="4"><c:v>Product E</c:v></c:pt>
    </c:strCache></c:strRef></c:cat>
    <c:val><c:numRef><c:f>Sheet1!$B$2:$B$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>40</c:v></c:pt><c:pt idx="1"><c:v>25</c:v></c:pt>
      <c:pt idx="2"><c:v>15</c:v></c:pt><c:pt idx="3"><c:v>12</c:v></c:pt>
      <c:pt idx="4"><c:v>8</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:gapWidth val="150"/>
  <c:splitPos val="2"/>
</c:ofPieChart>"""
plot_area68.append(etree.fromstring(ofpie_xml))

lbl68 = slide68.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl68.text_frame.paragraphs[0].text = "Slide 68: Pie of pie chart (ofPieChart)"
lbl68.text_frame.paragraphs[0].font.size = Pt(18)
lbl68.text_frame.paragraphs[0].font.bold = True

# ── Slide 69: 3D bar chart with view3D ────────────────────────────────────
slide69 = prs.slides.add_slide(blank)

chart_data69 = CategoryChartData()
chart_data69.categories = ['East', 'West', 'North', 'South']
chart_data69.add_series('Revenue', (320, 280, 190, 250))
chart_data69.add_series('Cost', (200, 180, 140, 160))

chart_frame69 = slide69.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), chart_data69
)

# Replace barChart with bar3DChart and add view3D via lxml
chart69 = chart_frame69.chart
chart_part69 = chart69.part
chart_elem69 = chart_part69._element
plot_area69 = chart_elem69.find(f'.//{{{ns_c}}}plotArea')
for bar in plot_area69.findall(f'{{{ns_c}}}barChart'):
    # Convert barChart to bar3DChart by changing tag
    bar.tag = f'{{{ns_c}}}bar3DChart'

# Add view3D element
chart_node69 = chart_elem69.find(f'{{{ns_c}}}chart')
view3d_69 = f"""<c:view3D xmlns:c="{ns_c}"><c:rotX val="20"/><c:rotY val="30"/><c:depthPercent val="150"/><c:rAngAx val="1"/><c:perspective val="40"/></c:view3D>"""
chart_node69.insert(0, etree.fromstring(view3d_69))

lbl69 = slide69.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl69.text_frame.paragraphs[0].text = "Slide 69: 3D bar chart (view3D preserved)"
lbl69.text_frame.paragraphs[0].font.size = Pt(18)
lbl69.text_frame.paragraphs[0].font.bold = True

# ── Slide 70: Text outline (a:rPr/a:ln) ──────────────────────────────────────
slide70 = prs.slides.add_slide(blank)
lbl70 = slide70.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl70.text_frame.paragraphs[0].text = "Slide 70: Text outline (a:rPr/a:ln)"
lbl70.text_frame.paragraphs[0].font.size = Pt(18)
lbl70.text_frame.paragraphs[0].font.bold = True

ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Text box with outlined text
tb70a = slide70.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
tf70a = tb70a.text_frame
p70a = tf70a.paragraphs[0]
r70a = p70a.add_run()
r70a.text = "Red Outlined Text"
r70a.font.size = Pt(48)
r70a.font.bold = True
# Inject a:ln into a:rPr via lxml
rpr70a = r70a._r.find(f'{{{ns_a}}}rPr')
ln70a = etree.SubElement(rpr70a, f'{{{ns_a}}}ln', attrib={'w': '25400'})
sf70a = etree.SubElement(ln70a, f'{{{ns_a}}}solidFill')
etree.SubElement(sf70a, f'{{{ns_a}}}srgbClr', attrib={'val': 'FF0000'})

# Thinner outline with different color
tb70b = slide70.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(8), Inches(1.5))
tf70b = tb70b.text_frame
p70b = tf70b.paragraphs[0]
r70b = p70b.add_run()
r70b.text = "Blue Outlined (thin)"
r70b.font.size = Pt(36)
r70b.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
rpr70b = r70b._r.find(f'{{{ns_a}}}rPr')
ln70b = etree.SubElement(rpr70b, f'{{{ns_a}}}ln', attrib={'w': '12700'})
sf70b = etree.SubElement(ln70b, f'{{{ns_a}}}solidFill')
etree.SubElement(sf70b, f'{{{ns_a}}}srgbClr', attrib={'val': '0000FF'})

# ── Slide 71: Text gradient fill (a:rPr/a:gradFill) ─────────────────────────
slide71 = prs.slides.add_slide(blank)
lbl71 = slide71.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl71.text_frame.paragraphs[0].text = "Slide 71: Text gradient fill (a:rPr/a:gradFill)"
lbl71.text_frame.paragraphs[0].font.size = Pt(18)
lbl71.text_frame.paragraphs[0].font.bold = True

# Text with gradient fill
tb71a = slide71.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2))
tf71a = tb71a.text_frame
p71a = tf71a.paragraphs[0]
r71a = p71a.add_run()
r71a.text = "Gradient Text"
r71a.font.size = Pt(60)
r71a.font.bold = True
# Inject a:gradFill into a:rPr
rpr71a = r71a._r.find(f'{{{ns_a}}}rPr')
grad_xml71 = f"""<a:gradFill xmlns:a="{ns_a}">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FF0000"/></a:gs>
    <a:gs pos="50000"><a:srgbClr val="FFFF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="1"/>
</a:gradFill>"""
rpr71a.append(etree.fromstring(grad_xml71))

# Second text with solid fill + gradient (gradient takes priority in display)
tb71b = slide71.shapes.add_textbox(Inches(0.5), Inches(4), Inches(8), Inches(2))
tf71b = tb71b.text_frame
p71b = tf71b.paragraphs[0]
r71b = p71b.add_run()
r71b.text = "Green-Blue Grad"
r71b.font.size = Pt(48)
rpr71b = r71b._r.find(f'{{{ns_a}}}rPr')
grad_xml71b = f"""<a:gradFill xmlns:a="{ns_a}">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="00FF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="5400000" scaled="1"/>
</a:gradFill>"""
rpr71b.append(etree.fromstring(grad_xml71b))

# ── Slide 72: Text warp (a:prstTxWarp) ──────────────────────────────────────
slide72 = prs.slides.add_slide(blank)
lbl72 = slide72.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl72.text_frame.paragraphs[0].text = "Slide 72: Text warp (a:prstTxWarp)"
lbl72.text_frame.paragraphs[0].font.size = Pt(18)
lbl72.text_frame.paragraphs[0].font.bold = True

ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# Text box with wave warp
tb72a = slide72.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(2))
tf72a = tb72a.text_frame
p72a = tf72a.paragraphs[0]
r72a = p72a.add_run()
r72a.text = "Wave Text"
r72a.font.size = Pt(36)
r72a.font.bold = True
# Inject prstTxWarp into bodyPr
sp72a = tb72a._element
body_pr72a = sp72a.find(f'.//{{{ns_a}}}bodyPr')
warp72a = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textWave1">
  <a:avLst>
    <a:gd name="adj" fmla="val 19773"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72a.insert(0, etree.fromstring(warp72a))

# Text box with arch up warp
tb72b = slide72.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(2))
tf72b = tb72b.text_frame
p72b = tf72b.paragraphs[0]
r72b = p72b.add_run()
r72b.text = "Arch Up"
r72b.font.size = Pt(36)
r72b.font.bold = True
sp72b = tb72b._element
body_pr72b = sp72b.find(f'.//{{{ns_a}}}bodyPr')
warp72b = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textArchUp">
  <a:avLst>
    <a:gd name="adj" fmla="val 10800000"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72b.insert(0, etree.fromstring(warp72b))

# Text box with no adjust values (textDeflate)
tb72c = slide72.shapes.add_textbox(Inches(0.5), Inches(4), Inches(8), Inches(2))
tf72c = tb72c.text_frame
p72c = tf72c.paragraphs[0]
r72c = p72c.add_run()
r72c.text = "Deflate Text"
r72c.font.size = Pt(36)
sp72c = tb72c._element
body_pr72c = sp72c.find(f'.//{{{ns_a}}}bodyPr')
warp72c = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textDeflate">
  <a:avLst/>
</a:prstTxWarp>"""
body_pr72c.insert(0, etree.fromstring(warp72c))

# ── Slide 72b: Additional text warp presets ──────────────────────────────────
slide72b = prs.slides.add_slide(blank)
lbl72b = slide72b.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl72b.text_frame.paragraphs[0].text = "Slide 72b: Additional text warp presets"
lbl72b.text_frame.paragraphs[0].font.size = Pt(18)
lbl72b.text_frame.paragraphs[0].font.bold = True

# textSlantUp
tb72d = slide72b.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(1.5))
tf72d = tb72d.text_frame
r72d = tf72d.paragraphs[0].add_run()
r72d.text = "Slant Up"
r72d.font.size = Pt(36)
r72d.font.bold = True
sp72d = tb72d._element
body_pr72d = sp72d.find(f'.//{{{ns_a}}}bodyPr')
warp72d = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textSlantUp">
  <a:avLst>
    <a:gd name="adj" fmla="val 55556"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72d.insert(0, etree.fromstring(warp72d))

# textCurveUp
tb72e = slide72b.shapes.add_textbox(Inches(5), Inches(1), Inches(4), Inches(1.5))
tf72e = tb72e.text_frame
r72e = tf72e.paragraphs[0].add_run()
r72e.text = "Curve Up"
r72e.font.size = Pt(36)
r72e.font.bold = True
r72e.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
sp72e = tb72e._element
body_pr72e = sp72e.find(f'.//{{{ns_a}}}bodyPr')
warp72e = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textCurveUp">
  <a:avLst>
    <a:gd name="adj" fmla="val 45977"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72e.insert(0, etree.fromstring(warp72e))

# textInflate
tb72f = slide72b.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4), Inches(1.5))
tf72f = tb72f.text_frame
r72f = tf72f.paragraphs[0].add_run()
r72f.text = "Inflate"
r72f.font.size = Pt(36)
r72f.font.bold = True
r72f.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
sp72f = tb72f._element
body_pr72f = sp72f.find(f'.//{{{ns_a}}}bodyPr')
warp72f = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textInflate">
  <a:avLst/>
</a:prstTxWarp>"""
body_pr72f.insert(0, etree.fromstring(warp72f))

# textChevron
tb72g = slide72b.shapes.add_textbox(Inches(5), Inches(3), Inches(4), Inches(1.5))
tf72g = tb72g.text_frame
r72g = tf72g.paragraphs[0].add_run()
r72g.text = "Chevron"
r72g.font.size = Pt(36)
r72g.font.bold = True
r72g.font.color.rgb = RGBColor(0x00, 0xB0, 0x50)
sp72g = tb72g._element
body_pr72g = sp72g.find(f'.//{{{ns_a}}}bodyPr')
warp72g = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textChevron">
  <a:avLst>
    <a:gd name="adj" fmla="val 25000"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72g.insert(0, etree.fromstring(warp72g))

# textCircle
tb72h = slide72b.shapes.add_textbox(Inches(2.5), Inches(5), Inches(4), Inches(2.5))
tf72h = tb72h.text_frame
r72h = tf72h.paragraphs[0].add_run()
r72h.text = "Circle Text Example"
r72h.font.size = Pt(28)
r72h.font.bold = True
sp72h = tb72h._element
body_pr72h = sp72h.find(f'.//{{{ns_a}}}bodyPr')
warp72h = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textCircle">
  <a:avLst>
    <a:gd name="adj" fmla="val 10800000"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72h.insert(0, etree.fromstring(warp72h))

# ── Slide 73: Stacked / Percent-stacked bar charts ──────────────────────────

slide73 = prs.slides.add_slide(blank)

# Percent-stacked horizontal bar (like gender breakdown: 58.1% / 41.9%)
chart_data73a = CategoryChartData()
chart_data73a.categories = ['Gender']
chart_data73a.add_series('Male', (58.1,))
chart_data73a.add_series('Female', (41.9,))

chart_frame73a = slide73.shapes.add_chart(
    XL_CHART_TYPE.BAR_STACKED_100, Inches(0.5), Inches(1.5),
    Inches(4), Inches(3), chart_data73a
)

# Stacked vertical column chart (multiple categories)
chart_data73b = CategoryChartData()
chart_data73b.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data73b.add_series('Product A', (100, 120, 90, 150))
chart_data73b.add_series('Product B', (80, 70, 110, 60))
chart_data73b.add_series('Product C', (50, 40, 30, 80))

chart_frame73b = slide73.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED, Inches(5), Inches(1.5),
    Inches(4), Inches(3), chart_data73b
)

# Percent-stacked vertical column chart
chart_data73c = CategoryChartData()
chart_data73c.categories = ['East', 'West', 'North']
chart_data73c.add_series('2024', (300, 200, 150))
chart_data73c.add_series('2025', (250, 350, 200))

chart_frame73c = slide73.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED_100, Inches(2.5), Inches(5),
    Inches(4), Inches(2), chart_data73c
)

lbl73 = slide73.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl73.text_frame.paragraphs[0].text = "Slide 73: Stacked / Percent-stacked bar charts"
lbl73.text_frame.paragraphs[0].font.size = Pt(18)
lbl73.text_frame.paragraphs[0].font.bold = True

