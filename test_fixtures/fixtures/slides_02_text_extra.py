"""Slides 15-22: capitalization, color map override, CS/sym fonts, rotation, vertical, hyperlinks, image bullet.

Mechanically carved from the legacy gen_test_features.py. Runs its
slide-building code at import time against the shared '_ctx.prs'.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
import base64, io, os, re, struct, tempfile, zipfile, zlib

from ._ctx import prs, blank, nsmap

# ── Slide 15: Capitalization (cap) ──────────────────────────────────────
slide15 = prs.slides.add_slide(blank)

title15 = slide15.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title15.text_frame.paragraphs[0].text = "Slide 15: Capitalization (a:rPr cap)"
title15.text_frame.paragraphs[0].font.size = Pt(24)
title15.text_frame.paragraphs[0].font.bold = True

# Box with cap="all"
box15a = slide15.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2.5))
tf15a = box15a.text_frame
tf15a.word_wrap = True
fill15a = box15a.fill
fill15a.solid()
fill15a.fore_color.rgb = RGBColor(230, 245, 255)

p15a = tf15a.paragraphs[0]
run15a = p15a.add_run()
run15a.text = "cap=all: This Should Be All Caps"
run15a.font.size = Pt(18)
rPr15a = run15a._r.find('a:rPr', nsmap)
if rPr15a is None:
    rPr15a = etree.SubElement(run15a._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run15a._r.insert(0, rPr15a)
rPr15a.set('cap', 'all')

p15a2 = tf15a.add_paragraph()
run15a2 = p15a2.add_run()
run15a2.text = "Mixed: "
run15a2.font.size = Pt(16)
run15a2b = p15a2.add_run()
run15a2b.text = "all caps part"
run15a2b.font.size = Pt(16)
rPr15a2b = run15a2b._r.find('a:rPr', nsmap)
if rPr15a2b is None:
    rPr15a2b = etree.SubElement(run15a2b._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run15a2b._r.insert(0, rPr15a2b)
rPr15a2b.set('cap', 'all')
run15a2c = p15a2.add_run()
run15a2c.text = " and normal part"
run15a2c.font.size = Pt(16)

p15a3 = tf15a.add_paragraph()
run15a3 = p15a3.add_run()
run15a3.text = "No cap: Regular text for comparison"
run15a3.font.size = Pt(16)

# Box with cap="small"
box15b = slide15.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(2.5))
tf15b = box15b.text_frame
tf15b.word_wrap = True
fill15b = box15b.fill
fill15b.solid()
fill15b.fore_color.rgb = RGBColor(255, 245, 230)

p15b = tf15b.paragraphs[0]
run15b = p15b.add_run()
run15b.text = "cap=small: Small Caps Text"
run15b.font.size = Pt(18)
rPr15b = run15b._r.find('a:rPr', nsmap)
if rPr15b is None:
    rPr15b = etree.SubElement(run15b._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run15b._r.insert(0, rPr15b)
rPr15b.set('cap', 'small')

p15b2 = tf15b.add_paragraph()
run15b2 = p15b2.add_run()
run15b2.text = "No cap: Regular text"
run15b2.font.size = Pt(16)

p15b3 = tf15b.add_paragraph()
run15b3 = p15b3.add_run()
run15b3.text = "Bold Small Caps"
run15b3.font.size = Pt(18)
run15b3.font.bold = True
rPr15b3 = run15b3._r.find('a:rPr', nsmap)
if rPr15b3 is None:
    rPr15b3 = etree.SubElement(run15b3._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run15b3._r.insert(0, rPr15b3)
rPr15b3.set('cap', 'small')

# Info text
info15 = slide15.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
tf15info = info15.text_frame
tf15info.word_wrap = True
p15info = tf15info.paragraphs[0]
p15info.text = "cap=all: lowercase a-z should display as UPPERCASE A-Z. cap=small: should display with SVG font-variant=small-caps."
p15info.font.size = Pt(12)
p15info.font.color.rgb = RGBColor(100, 100, 100)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 16: Color map override (p:clrMapOvr)
# ═══════════════════════════════════════════════════════════════════════════════
slide16 = prs.slides.add_slide(prs.slide_layouts[0])

nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

# Remove existing p:clrMapOvr (python-pptx adds a:masterClrMapping by default)
existing_ovr = slide16._element.findall(
    '{http://schemas.openxmlformats.org/presentationml/2006/main}clrMapOvr'
)
for ovr in existing_ovr:
    slide16._element.remove(ovr)

# Add p:clrMapOvr with a:overrideClrMapping that swaps bg1↔tx1 (dk1↔lt1)
# This makes the slide use dark background + light text
clrMapOvr = etree.SubElement(
    slide16._element,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}clrMapOvr'
)
overrideMapping = etree.SubElement(
    clrMapOvr,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}overrideClrMapping',
    attrib={
        'bg1': 'dk1',
        'tx1': 'lt1',
        'bg2': 'lt2',
        'tx2': 'dk2',
        'accent1': 'accent1',
        'accent2': 'accent2',
        'accent3': 'accent3',
        'accent4': 'accent4',
        'accent5': 'accent5',
        'accent6': 'accent6',
        'hlink': 'hlink',
        'folHlink': 'folHlink',
    }
)

# Add a textbox with scheme-colored text (should resolve to light color on dark bg)
tb16 = slide16.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
tf16 = tb16.text_frame
tf16.word_wrap = True
p16 = tf16.paragraphs[0]
p16.text = "Color Map Override: bg1=dk1, tx1=lt1"
p16.font.size = Pt(28)
p16.font.bold = True
# Use scheme color tx1 (which should resolve to lt1 = white on this slide)
rPr16 = p16.runs[0]._r.find('a:rPr', nsmap)
if rPr16 is None:
    rPr16 = etree.SubElement(p16.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p16.runs[0]._r.insert(0, rPr16)
solidFill16 = etree.SubElement(rPr16, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
schemeClr16 = etree.SubElement(solidFill16, '{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr', attrib={'val': 'tx1'})

# Add a second textbox with explicit description
tb16b = slide16.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
tf16b = tb16b.text_frame
tf16b.word_wrap = True
p16b = tf16b.paragraphs[0]
p16b.text = "This slide has clrMapOvr swapping bg1/tx1. Text using scheme color tx1 should appear light (white) because tx1 maps to lt1."
p16b.font.size = Pt(16)
# Also use scheme color for this text
rPr16b = p16b.runs[0]._r.find('a:rPr', nsmap)
if rPr16b is None:
    rPr16b = etree.SubElement(p16b.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p16b.runs[0]._r.insert(0, rPr16b)
solidFill16b = etree.SubElement(rPr16b, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
schemeClr16b = etree.SubElement(solidFill16b, '{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr', attrib={'val': 'tx1'})

# Add slide background using scheme color bg1 (which maps to dk1 = dark)
bg16 = slide16._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld', nsmap)
if bg16 is None:
    bg16 = slide16._element.find('p:cSld', nsmap)
# Insert p:bg before p:cSld
pBg = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
bgPr = etree.SubElement(pBg, '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr')
solidFillBg = etree.SubElement(bgPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
# Use a dark solid color for background to make the slide visually dark
srgbClrBg = etree.SubElement(solidFillBg, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr', attrib={'val': '1A1A2E'})
# Insert bg element at the right position in slide XML
cSld = slide16._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
if cSld is not None:
    slide16._element.insert(list(slide16._element).index(cSld), pBg)

# ── Slide 17: CS/Sym fonts + kerning ─────────────────────────────────────────

slide17 = prs.slides.add_slide(blank)

# Textbox with Complex Script font (a:cs)
tb17a = slide17.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
tf17a = tb17a.text_frame
tf17a.word_wrap = True
p17a = tf17a.paragraphs[0]
p17a.text = "Complex Script Font (Arabic style)"
p17a.font.size = Pt(24)
p17a.font.bold = True
# Add a:cs element to rPr
rPr17a = p17a.runs[0]._r.find('a:rPr', nsmap)
if rPr17a is None:
    rPr17a = etree.SubElement(p17a.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p17a.runs[0]._r.insert(0, rPr17a)
etree.SubElement(rPr17a, '{http://schemas.openxmlformats.org/drawingml/2006/main}cs', attrib={'typeface': 'Arial'})

# Textbox with Symbol font (a:sym)
tb17b = slide17.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
tf17b = tb17b.text_frame
p17b = tf17b.paragraphs[0]
p17b.text = "Symbol Font Text"
p17b.font.size = Pt(24)
rPr17b = p17b.runs[0]._r.find('a:rPr', nsmap)
if rPr17b is None:
    rPr17b = etree.SubElement(p17b.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p17b.runs[0]._r.insert(0, rPr17b)
etree.SubElement(rPr17b, '{http://schemas.openxmlformats.org/drawingml/2006/main}sym', attrib={'typeface': 'Wingdings'})

# Textbox with kerning
tb17c = slide17.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
tf17c = tb17c.text_frame
p17c = tf17c.paragraphs[0]
p17c.text = "Kerning enabled at 1200 hundredths-pt"
p17c.font.size = Pt(24)
rPr17c = p17c.runs[0]._r.find('a:rPr', nsmap)
if rPr17c is None:
    rPr17c = etree.SubElement(p17c.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p17c.runs[0]._r.insert(0, rPr17c)
rPr17c.set('kern', '1200')

# ── Slide 18: Text rotation + tab stops ──────────────────────────────────────

slide18 = prs.slides.add_slide(blank)

# Textbox with rotated text (bodyPr rot)
tb18a = slide18.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(2))
tf18a = tb18a.text_frame
p18a = tf18a.paragraphs[0]
p18a.text = "Rotated text (45 degrees)"
p18a.font.size = Pt(20)
# Set rot on bodyPr (45 degrees = 2700000 in 60000ths)
bodyPr18a = tb18a._element.find('.//a:bodyPr', nsmap)
bodyPr18a.set('rot', '2700000')

# Textbox with tab stops
tb18b = slide18.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(3))
tf18b = tb18b.text_frame
tf18b.word_wrap = True
p18b = tf18b.paragraphs[0]
p18b.text = "Col1\tCol2\tCol3"
p18b.font.size = Pt(18)
# Add tab stops via XML
pPr18b = p18b._pPr
if pPr18b is None:
    pPr18b = etree.SubElement(p18b._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p18b._p.insert(0, pPr18b)
tabLst = etree.SubElement(pPr18b, '{http://schemas.openxmlformats.org/drawingml/2006/main}tabLst')
etree.SubElement(tabLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}tab', attrib={'pos': '2743200', 'algn': 'l'})  # 3 inches
etree.SubElement(tabLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}tab', attrib={'pos': '5486400', 'algn': 'r'})  # 6 inches

# ── Slide 19: Vertical text + text columns ──────────────────────────────────

slide19 = prs.slides.add_slide(blank)

# Textbox with vertical text (vert="vert")
tb19a = slide19.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(5))
tf19a = tb19a.text_frame
p19a = tf19a.paragraphs[0]
p19a.text = "Vertical text"
p19a.font.size = Pt(20)
bodyPr19a = tb19a._element.find('.//a:bodyPr', nsmap)
bodyPr19a.set('vert', 'vert')

# Textbox with eaVert
tb19b = slide19.shapes.add_textbox(Inches(3), Inches(0.5), Inches(2), Inches(5))
tf19b = tb19b.text_frame
p19b = tf19b.paragraphs[0]
p19b.text = "EA Vertical"
p19b.font.size = Pt(20)
bodyPr19b = tb19b._element.find('.//a:bodyPr', nsmap)
bodyPr19b.set('vert', 'eaVert')

# Textbox with text columns (numCol + spcCol)
tb19c = slide19.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(4), Inches(5))
tf19c = tb19c.text_frame
tf19c.word_wrap = True
p19c = tf19c.paragraphs[0]
p19c.text = "This is column text that spans multiple columns. The text should flow from the first column to the second column when it reaches the bottom."
p19c.font.size = Pt(14)
bodyPr19c = tb19c._element.find('.//a:bodyPr', nsmap)
bodyPr19c.set('numCol', '2')
bodyPr19c.set('spcCol', '457200')  # 0.5 inch spacing

# ── Slide 20: Hyperlink + RTL ────────────────────────────────────────────────

slide20 = prs.slides.add_slide(blank)

# Textbox with hyperlink
tb20a = slide20.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
tf20a = tb20a.text_frame
tf20a.word_wrap = True
p20a = tf20a.paragraphs[0]
r20a = p20a.runs[0] if len(p20a.runs) > 0 else p20a.add_run()
r20a.text = "Click here to visit example.com"
r20a.font.size = Pt(20)
r20a.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
r20a.font.underline = True
# Add hlinkClick via XML — need to add a relationship first
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
rel = slide20.part.relate_to('https://example.com', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
rPr20a = r20a._r.find('a:rPr', nsmap)
if rPr20a is None:
    rPr20a = etree.SubElement(r20a._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    r20a._r.insert(0, rPr20a)
etree.SubElement(rPr20a, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rel})

# RTL paragraph
tb20b = slide20.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(2))
tf20b = tb20b.text_frame
tf20b.word_wrap = True
p20b = tf20b.paragraphs[0]
p20b.text = "RTL paragraph text (right-to-left)"
p20b.font.size = Pt(20)
pPr20b = p20b._pPr
if pPr20b is None:
    pPr20b = etree.SubElement(p20b._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p20b._p.insert(0, pPr20b)
pPr20b.set('rtl', '1')

# ── Slide 21: Image bullet (a:buBlip) ───────────────────────────────────────

slide21 = prs.slides.add_slide(blank)

# For image bullets, we need a small embedded image
# We'll create a tiny 1x1 red PNG in memory
def make_tiny_png(r, g, b):
    """Create a minimal 1x1 PNG."""
    # IHDR
    ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
    ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
    # IDAT
    raw = bytes([0, r, g, b])
    compressed = zlib.compress(raw)
    idat_crc = zlib.crc32(b'IDAT' + compressed) & 0xffffffff
    idat = struct.pack('>I', len(compressed)) + b'IDAT' + compressed + struct.pack('>I', idat_crc)
    # IEND
    iend_crc = zlib.crc32(b'IEND') & 0xffffffff
    iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
    return b'\x89PNG\r\n\x1a\n' + ihdr + idat + iend

png_data = make_tiny_png(255, 0, 0)  # red dot
png_stream = io.BytesIO(png_data)

# Add image as a relationship to the slide
from pptx.opc.package import Part
from pptx.opc.constants import CONTENT_TYPE as CT
image_part, rId_img = slide21.part.get_or_add_image_part(png_stream)

# Create textbox with bullet paragraphs
tb21 = slide21.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
tf21 = tb21.text_frame
tf21.word_wrap = True
p21 = tf21.paragraphs[0]
p21.text = "Image bullet paragraph 1"
p21.font.size = Pt(20)

# Add image bullet via XML: <a:buBlip><a:blip r:embed="rIdN"/></a:buBlip>
pPr21 = p21._pPr
if pPr21 is None:
    pPr21 = etree.SubElement(p21._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p21._p.insert(0, pPr21)
buBlip = etree.SubElement(pPr21, '{http://schemas.openxmlformats.org/drawingml/2006/main}buBlip')
etree.SubElement(buBlip, '{http://schemas.openxmlformats.org/drawingml/2006/main}blip',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed': rId_img})

# Add second paragraph with same image bullet
p21b_elem = etree.SubElement(tf21._txBody, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
pPr21b = etree.SubElement(p21b_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
buBlip21b = etree.SubElement(pPr21b, '{http://schemas.openxmlformats.org/drawingml/2006/main}buBlip')
etree.SubElement(buBlip21b, '{http://schemas.openxmlformats.org/drawingml/2006/main}blip',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed': rId_img})
r21b = etree.SubElement(p21b_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
rPr21b = etree.SubElement(r21b, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr', attrib={'lang': 'en-US', 'sz': '2000'})
t21b = etree.SubElement(r21b, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
t21b.text = "Image bullet paragraph 2"

# ── Slide 22: Hover link + link color ──────────────────────────────────────────
slide22 = prs.slides.add_slide(blank)
tb22 = slide22.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
tf22 = tb22.text_frame
tf22.word_wrap = True

# Paragraph 1: hlinkClick (normal hyperlink)
p22a = tf22.paragraphs[0]
p22a.text = ""
r22a = p22a.add_run()
r22a.text = "Click link"
r22a.font.size = Pt(24)
# Add hlinkClick via XML
rPr22a = r22a._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
if rPr22a is None:
    rPr22a = etree.SubElement(r22a._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    r22a._r.insert(0, rPr22a)
# Add relationship for click link
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
rId_click = slide22.part.relate_to('https://example.com/click', RT.HYPERLINK, is_external=True)
etree.SubElement(rPr22a, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rId_click})

# Paragraph 2: hlinkHover (mouse-over hyperlink)
p22b = tf22.add_paragraph()
r22b = p22b.add_run()
r22b.text = "Hover link"
r22b.font.size = Pt(24)
rPr22b = r22b._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
if rPr22b is None:
    rPr22b = etree.SubElement(r22b._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    r22b._r.insert(0, rPr22b)
rId_hover = slide22.part.relate_to('https://example.com/hover', RT.HYPERLINK, is_external=True)
etree.SubElement(rPr22b, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkHover',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rId_hover})

# Paragraph 3: both hlinkClick + hlinkHover
p22c = tf22.add_paragraph()
r22c = p22c.add_run()
r22c.text = "Both links"
r22c.font.size = Pt(24)
rPr22c = r22c._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
if rPr22c is None:
    rPr22c = etree.SubElement(r22c._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    r22c._r.insert(0, rPr22c)
rId_click2 = slide22.part.relate_to('https://example.com/both-click', RT.HYPERLINK, is_external=True)
rId_hover2 = slide22.part.relate_to('https://example.com/both-hover', RT.HYPERLINK, is_external=True)
etree.SubElement(rPr22c, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rId_click2})
etree.SubElement(rPr22c, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkHover',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rId_hover2})

