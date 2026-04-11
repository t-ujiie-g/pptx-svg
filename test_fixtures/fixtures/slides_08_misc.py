"""Slides 74-88: notes, comments, SmartArt fallback, OLE, media, OMML math, transition/timing, hidden slide, WMF, blur/prstShdw/fillOverlay, justified text.

Mechanically carved from the legacy gen_test_features.py. Runs its
slide-building code at import time against the shared '_ctx.prs'.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
import base64, io, os, re, struct, tempfile, zipfile, zlib

from ._ctx import prs, blank

# ── Slide 74: Speaker notes + comments ──────────────────────────────────────

slide74 = prs.slides.add_slide(blank)

lbl74 = slide74.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl74.text_frame.paragraphs[0].text = "Slide 74: Speaker notes + comments"
lbl74.text_frame.paragraphs[0].font.size = Pt(18)
lbl74.text_frame.paragraphs[0].font.bold = True

body74 = slide74.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(3))
body74.text_frame.paragraphs[0].text = (
    "This slide has speaker notes and comments attached.\n"
    "Notes and comments are metadata not rendered on slide, "
    "but preserved in round-trip export."
)
body74.text_frame.paragraphs[0].font.size = Pt(14)
body74.text_frame.word_wrap = True

# Add speaker notes
notes_slide74 = slide74.notes_slide
notes_tf = notes_slide74.notes_text_frame
notes_tf.text = ""
p1 = notes_tf.paragraphs[0]
p1.text = "These are the speaker notes for slide 74."
p2 = notes_tf.add_paragraph()
p2.text = "Second paragraph of notes with key points."
p3 = notes_tf.add_paragraph()
p3.text = "Remember to mention the round-trip preservation."

# Add comments via raw XML (python-pptx has no public comment API)
from pptx.opc.package import Part as OpcPart
from pptx.opc.packuri import PackURI

authors_xml = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:cmAuthorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
    '<p:cmAuthor id="1" name="Test User" initials="TU" lastIdx="2" clrIdx="0"/>'
    '</p:cmAuthorLst>'
)

# Add commentAuthors part
authors_part = OpcPart(
    PackURI('/ppt/commentAuthors.xml'),
    'application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml',
    prs.part.package,
    authors_xml.encode('utf-8'),
)
prs.part.relate_to(authors_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors')

# Add comments part for slide 74
comment_xml = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:cmLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
    '<p:cm authorId="1" dt="2025-01-15T10:30:00.000" idx="1">'
    '<p:pos x="100" y="200"/>'
    '<p:text>This is a test comment on slide 74.</p:text>'
    '</p:cm>'
    '<p:cm authorId="1" dt="2025-01-15T11:00:00.000" idx="2">'
    '<p:pos x="300" y="400"/>'
    '<p:text>Second comment with review feedback.</p:text>'
    '</p:cm>'
    '</p:cmLst>'
)

slide74_idx = len(prs.slides)
comment_part = OpcPart(
    PackURI(f'/ppt/comments/comment{slide74_idx}.xml'),
    'application/vnd.openxmlformats-officedocument.presentationml.comments+xml',
    prs.part.package,
    comment_xml.encode('utf-8'),
)
slide74.part.relate_to(comment_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments')

# ── Slide 75: SmartArt fallback (mc:AlternateContent) ────────────────────────

slide75 = prs.slides.add_slide(blank)

# We inject raw XML with mc:AlternateContent containing:
#   mc:Choice  → dummy p:graphicFrame referencing diagram (preserved for round-trip)
#   mc:Fallback → p:grpSp with 3 rectangles representing a simple process flow
slide75_sp_tree = slide75.shapes._spTree
# Register mc: namespace prefix on the root <p:sld> element so lxml uses mc: prefix
sld_root = slide75_sp_tree.getparent()  # p:cSld -> p:sld
if sld_root is not None and sld_root.getparent() is not None:
    sld_root = sld_root.getparent()
MC_NS = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
DGM_NS = 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
etree.register_namespace('mc', MC_NS)
etree.register_namespace('dgm', DGM_NS)

mc_ac = etree.SubElement(slide75_sp_tree, f'{{{MC_NS}}}AlternateContent')

# mc:Choice — dummy SmartArt reference (for round-trip preservation)
mc_choice = etree.SubElement(mc_ac, f'{{{MC_NS}}}Choice', attrib={'Requires': 'dgm'})
choice_gf = etree.SubElement(mc_choice,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame')
choice_nvgf = etree.SubElement(choice_gf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvGraphicFramePr')
choice_cnvpr = etree.SubElement(choice_nvgf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr',
    attrib={'id': '100', 'name': 'SmartArt Diagram'})
choice_cnvgf = etree.SubElement(choice_nvgf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvGraphicFramePr')
choice_nvpr = etree.SubElement(choice_nvgf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
choice_xfrm = etree.SubElement(choice_gf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}xfrm')
etree.SubElement(choice_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}off',
    attrib={'x': '457200', 'y': '1371600'})
etree.SubElement(choice_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}ext',
    attrib={'cx': '8229600', 'y': '4572000'})
# Dummy graphic with dgm:relIds
choice_graphic = etree.SubElement(choice_gf,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}graphic')
choice_gd = etree.SubElement(choice_graphic,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData',
    attrib={'uri': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'})
etree.SubElement(choice_gd,
    f'{{{DGM_NS}}}relIds',
    attrib={
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}dm': 'rId10',
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}lo': 'rId11',
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}qs': 'rId12',
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}cs': 'rId13',
    })

# mc:Fallback — pre-rendered shapes (3 rectangles as a simple process flow)
mc_fallback = etree.SubElement(mc_ac,
    f'{{{MC_NS}}}Fallback')

# Create a group shape containing 3 process boxes
grp_sp = etree.SubElement(mc_fallback,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}grpSp')
grp_nv = etree.SubElement(grp_sp,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvGrpSpPr')
etree.SubElement(grp_nv,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr',
    attrib={'id': '101', 'name': 'SmartArt Fallback Group'})
etree.SubElement(grp_nv,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvGrpSpPr')
etree.SubElement(grp_nv,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')

grp_sp_pr = etree.SubElement(grp_sp,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr')
grp_xfrm = etree.SubElement(grp_sp_pr,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
etree.SubElement(grp_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}off',
    attrib={'x': '457200', 'y': '1371600'})
etree.SubElement(grp_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}ext',
    attrib={'cx': '8229600', 'cy': '2743200'})
etree.SubElement(grp_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}chOff',
    attrib={'x': '0', 'y': '0'})
etree.SubElement(grp_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}chExt',
    attrib={'cx': '8229600', 'cy': '2743200'})

# Helper to add a rectangle shape inside the group
def add_smartart_box(parent, sp_id, name, x, y, cx, cy, text, fill_hex):
    sp = etree.SubElement(parent,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sp')
    nv = etree.SubElement(sp,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
    etree.SubElement(nv,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr',
        attrib={'id': str(sp_id), 'name': name})
    etree.SubElement(nv,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvSpPr')
    etree.SubElement(nv,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
    spPr = etree.SubElement(sp,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
    xfrm = etree.SubElement(spPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    etree.SubElement(xfrm,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}off',
        attrib={'x': str(x), 'y': str(y)})
    etree.SubElement(xfrm,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}ext',
        attrib={'cx': str(cx), 'cy': str(cy)})
    etree.SubElement(spPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom',
        attrib={'prst': 'roundRect'}).append(
        etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst'))
    solid = etree.SubElement(spPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    etree.SubElement(solid,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
        attrib={'val': fill_hex})
    # Text body
    txBody = etree.SubElement(sp,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}txBody')
    bodyPr = etree.SubElement(txBody,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr',
        attrib={'anchor': 'ctr', 'wrap': 'square'})
    etree.SubElement(txBody,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
    p = etree.SubElement(txBody,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
    pPr = etree.SubElement(p,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr',
        attrib={'algn': 'ctr'})
    r = etree.SubElement(p,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
    rPr = etree.SubElement(r,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr',
        attrib={'lang': 'en-US', 'sz': '1800', 'b': '1'})
    solid_font = etree.SubElement(rPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    etree.SubElement(solid_font,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
        attrib={'val': 'FFFFFF'})
    t = etree.SubElement(r,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
    t.text = text

# Three process boxes: Plan → Build → Ship
add_smartart_box(grp_sp, 102, 'Step 1', 0, 457200, 2286000, 1828800, 'Plan', '4472C4')
add_smartart_box(grp_sp, 103, 'Step 2', 2971800, 457200, 2286000, 1828800, 'Build', 'ED7D31')
add_smartart_box(grp_sp, 104, 'Step 3', 5943600, 457200, 2286000, 1828800, 'Ship', '70AD47')

# Also add a title textbox
title75 = slide75.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title75.text_frame.paragraphs[0].text = "Slide 75: SmartArt Fallback (mc:AlternateContent)"
title75.text_frame.paragraphs[0].font.size = Pt(24)
title75.text_frame.paragraphs[0].font.bold = True

# ── Slide 76: OLE embedded object (p:oleObj with fallback image) ──────────────

slide76 = prs.slides.add_slide(blank)

# Create a small 2x2 red PNG as the OLE fallback image
def make_tiny_png(r, g, b, w=2, h=2):
    """Create a minimal RGBA PNG."""
    raw = b''
    for _ in range(h):
        raw += b'\x00'  # filter: none
        for _ in range(w):
            raw += struct.pack('BBBB', r, g, b, 255)
    compressed = zlib.compress(raw)
    def chunk(ctype, data):
        c = ctype + data
        return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
    ihdr = struct.pack('>IIBBBBB', w, h, 8, 6, 0, 0, 0)  # 8-bit RGBA
    return b'\x89PNG\r\n\x1a\n' + chunk(b'IHDR', ihdr) + chunk(b'IDAT', compressed) + chunk(b'IEND', b'')

ole_fallback_png = make_tiny_png(200, 50, 50)

# Add the PNG as an image part via the slide's relationships
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part as OpcPart
from pptx.opc.packuri import PackURI

slide76_idx = len(prs.slides)
ole_img_partname = PackURI(f'/ppt/media/oleImage{slide76_idx}.png')
ole_img_part = OpcPart(
    ole_img_partname,
    'image/png',
    prs.part.package,
    ole_fallback_png,
)
ole_img_rel = slide76.part.relate_to(ole_img_part, RT.IMAGE)

# Also create a dummy OLE binary part (empty, just for structure)
ole_bin_partname = PackURI(f'/ppt/embeddings/oleObject{slide76_idx}.bin')
ole_bin_part = OpcPart(
    ole_bin_partname,
    'application/vnd.openxmlformats-officedocument.oleObject',
    prs.part.package,
    b'\x00' * 16,  # Minimal dummy data
)
ole_bin_rel = slide76.part.relate_to(
    ole_bin_part,
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject',
)

# Inject p:graphicFrame with p:oleObj into spTree
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

slide76_sp_tree = slide76.shapes._spTree
gf = etree.SubElement(slide76_sp_tree, f'{{{P_NS}}}graphicFrame')

# nvGraphicFramePr
nv = etree.SubElement(gf, f'{{{P_NS}}}nvGraphicFramePr')
etree.SubElement(nv, f'{{{P_NS}}}cNvPr', attrib={'id': '200', 'name': 'OLE Object'})
etree.SubElement(nv, f'{{{P_NS}}}cNvGraphicFramePr')
etree.SubElement(nv, f'{{{P_NS}}}nvPr')

# xfrm
xfrm = etree.SubElement(gf, f'{{{P_NS}}}xfrm')
etree.SubElement(xfrm, f'{{{A_NS}}}off', attrib={'x': '914400', 'y': '1828800'})
etree.SubElement(xfrm, f'{{{A_NS}}}ext', attrib={'cx': '4572000', 'cy': '3429000'})

# a:graphic > a:graphicData (OLE URI)
graphic = etree.SubElement(gf, f'{{{A_NS}}}graphic')
gdata = etree.SubElement(graphic, f'{{{A_NS}}}graphicData',
    attrib={'uri': 'http://schemas.openxmlformats.org/presentationml/2006/ole'})

# p:oleObj with r:id and fallback p:pic
ole_obj = etree.SubElement(gdata, f'{{{P_NS}}}oleObj', attrib={
    f'{{{R_NS}}}id': ole_bin_rel,
    'imgW': '4572000',
    'imgH': '3429000',
    'progId': 'Excel.Sheet.12',
    'name': 'Embedded Spreadsheet',
})
etree.SubElement(ole_obj, f'{{{P_NS}}}embed')

# p:pic inside oleObj (fallback image)
ole_pic = etree.SubElement(ole_obj, f'{{{P_NS}}}pic')
ole_pic_nv = etree.SubElement(ole_pic, f'{{{P_NS}}}nvPicPr')
etree.SubElement(ole_pic_nv, f'{{{P_NS}}}cNvPr', attrib={'id': '201', 'name': 'OLE Fallback'})
etree.SubElement(ole_pic_nv, f'{{{P_NS}}}cNvPicPr')
etree.SubElement(ole_pic_nv, f'{{{P_NS}}}nvPr')
ole_pic_bf = etree.SubElement(ole_pic, f'{{{P_NS}}}blipFill')
etree.SubElement(ole_pic_bf, f'{{{A_NS}}}blip', attrib={f'{{{R_NS}}}embed': ole_img_rel})
etree.SubElement(ole_pic_bf, f'{{{A_NS}}}stretch').append(
    etree.Element(f'{{{A_NS}}}fillRect'))
ole_pic_sp = etree.SubElement(ole_pic, f'{{{P_NS}}}spPr')
ole_pic_xfrm = etree.SubElement(ole_pic_sp, f'{{{A_NS}}}xfrm')
etree.SubElement(ole_pic_xfrm, f'{{{A_NS}}}off', attrib={'x': '0', 'y': '0'})
etree.SubElement(ole_pic_xfrm, f'{{{A_NS}}}ext', attrib={'cx': '4572000', 'cy': '3429000'})
etree.SubElement(ole_pic_sp, f'{{{A_NS}}}prstGeom', attrib={'prst': 'rect'}).append(
    etree.Element(f'{{{A_NS}}}avLst'))

# Title
title76 = slide76.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title76.text_frame.paragraphs[0].text = "Slide 76: OLE Embedded Object (p:oleObj)"
title76.text_frame.paragraphs[0].font.size = Pt(24)
title76.text_frame.paragraphs[0].font.bold = True

# ── Slide 77: Media (video with poster frame) ────────────────────────────────

slide77 = prs.slides.add_slide(blank)

# Create a poster frame image (blue rectangle) for the video
poster_png = make_tiny_png(30, 100, 200, w=4, h=3)

slide77_idx = len(prs.slides)
poster_partname = PackURI(f'/ppt/media/posterFrame{slide77_idx}.png')
poster_part = OpcPart(
    poster_partname,
    'image/png',
    prs.part.package,
    poster_png,
)
poster_rel = slide77.part.relate_to(poster_part, RT.IMAGE)

# Create a dummy video part
video_partname = PackURI(f'/ppt/media/video{slide77_idx}.mp4')
video_part = OpcPart(
    video_partname,
    'video/mp4',
    prs.part.package,
    b'\x00' * 32,  # Minimal dummy data
)
video_rel = slide77.part.relate_to(
    video_part,
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/video',
)

# Inject p:pic with a:videoFile into spTree
slide77_sp_tree = slide77.shapes._spTree

vid_pic = etree.SubElement(slide77_sp_tree, f'{{{P_NS}}}pic')

# nvPicPr
vid_nv = etree.SubElement(vid_pic, f'{{{P_NS}}}nvPicPr')
etree.SubElement(vid_nv, f'{{{P_NS}}}cNvPr', attrib={'id': '300', 'name': 'Video Placeholder'})
etree.SubElement(vid_nv, f'{{{P_NS}}}cNvPicPr')
vid_nvpr = etree.SubElement(vid_nv, f'{{{P_NS}}}nvPr')
etree.SubElement(vid_nvpr, f'{{{A_NS}}}videoFile', attrib={f'{{{R_NS}}}link': video_rel})

# blipFill (poster frame)
vid_bf = etree.SubElement(vid_pic, f'{{{P_NS}}}blipFill')
etree.SubElement(vid_bf, f'{{{A_NS}}}blip', attrib={f'{{{R_NS}}}embed': poster_rel})
etree.SubElement(vid_bf, f'{{{A_NS}}}stretch').append(
    etree.Element(f'{{{A_NS}}}fillRect'))

# spPr
vid_sp = etree.SubElement(vid_pic, f'{{{P_NS}}}spPr')
vid_xfrm = etree.SubElement(vid_sp, f'{{{A_NS}}}xfrm')
etree.SubElement(vid_xfrm, f'{{{A_NS}}}off', attrib={'x': '1371600', 'y': '1828800'})
etree.SubElement(vid_xfrm, f'{{{A_NS}}}ext', attrib={'cx': '6858000', 'cy': '3886200'})
etree.SubElement(vid_sp, f'{{{A_NS}}}prstGeom', attrib={'prst': 'rect'}).append(
    etree.Element(f'{{{A_NS}}}avLst'))

# Title
title77 = slide77.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title77.text_frame.paragraphs[0].text = "Slide 77: Media (Video with Poster Frame)"
title77.text_frame.paragraphs[0].font.size = Pt(24)
title77.text_frame.paragraphs[0].font.bold = True

# ── Slide 78: Math equation (OMML) ────────────────────────────────────────────

slide78 = prs.slides.add_slide(blank)

# Title
title78 = slide78.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
title78.text_frame.paragraphs[0].text = "Slide 78: Math Equation (OMML)"
title78.text_frame.paragraphs[0].font.size = Pt(28)
title78.text_frame.paragraphs[0].font.bold = True

# Add a textbox, then inject OMML XML directly into the slide XML
math_tb = slide78.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
math_tb.text_frame.paragraphs[0].text = "MATH_PLACEHOLDER"
math_tb.text_frame.paragraphs[0].font.size = Pt(24)

# ── Slide 79: Transition + Timing (round-trip preservation) ───────────────────

slide79 = prs.slides.add_slide(blank)

title79 = slide79.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title79.text_frame.paragraphs[0].text = "Slide 79: Transition + Timing (round-trip)"
title79.text_frame.paragraphs[0].font.size = Pt(24)
title79.text_frame.paragraphs[0].font.bold = True

desc79 = slide79.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
desc79.text_frame.paragraphs[0].text = "This slide has a fade transition and timing data for round-trip."
desc79.text_frame.paragraphs[0].font.size = Pt(16)

# ── Slide 80: Hidden slide ───────────────────────────────────────────────────

slide80 = prs.slides.add_slide(blank)

title80 = slide80.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title80.text_frame.paragraphs[0].text = "Slide 80: Hidden Slide"
title80.text_frame.paragraphs[0].font.size = Pt(24)
title80.text_frame.paragraphs[0].font.bold = True

desc80 = slide80.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
desc80.text_frame.paragraphs[0].text = "This slide is marked as hidden (show='0')."
desc80.text_frame.paragraphs[0].font.size = Pt(16)

# ── Slide 81: WMF → SVG conversion ───────────────────────────────────────────

slide81 = prs.slides.add_slide(blank)

title81 = slide81.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title81.text_frame.paragraphs[0].text = "Slide 81: WMF Image (WMF → SVG conversion)"
title81.text_frame.paragraphs[0].font.size = Pt(24)
title81.text_frame.paragraphs[0].font.bold = True

desc81 = slide81.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
desc81.text_frame.paragraphs[0].text = "WMF images below are converted to SVG at runtime."
desc81.text_frame.paragraphs[0].font.size = Pt(14)

# Placeholder — will be replaced with p:pic referencing a WMF via ZIP patching
wmf_placeholder = slide81.shapes.add_textbox(Inches(1), Inches(2.5), Inches(3), Inches(3))
wmf_placeholder.text_frame.paragraphs[0].text = "WMF_PLACEHOLDER"
wmf_placeholder.text_frame.paragraphs[0].font.size = Pt(12)

# ── Slide 82: OMML — Large operator (integral) + delimiters ─────────────────

slide82 = prs.slides.add_slide(blank)

title82 = slide82.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title82.text_frame.paragraphs[0].text = "Slide 82: OMML — Large Operators + Delimiters"
title82.text_frame.paragraphs[0].font.size = Pt(24)
title82.text_frame.paragraphs[0].font.bold = True

math82 = slide82.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
math82.text_frame.paragraphs[0].text = "NARY_PLACEHOLDER"
math82.text_frame.paragraphs[0].font.size = Pt(24)

# ── Slide 83: OMML — Matrix + Delimiters ────────────────────────────────────

slide83 = prs.slides.add_slide(blank)

title83 = slide83.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title83.text_frame.paragraphs[0].text = "Slide 83: OMML — Matrix + Delimiters"
title83.text_frame.paragraphs[0].font.size = Pt(24)
title83.text_frame.paragraphs[0].font.bold = True

math83 = slide83.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
math83.text_frame.paragraphs[0].text = "MATRIX_PLACEHOLDER"
math83.text_frame.paragraphs[0].font.size = Pt(24)

# ── Slide 84: OMML — Accent + Bar + SubSup ──────────────────────────────────

slide84 = prs.slides.add_slide(blank)

title84 = slide84.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title84.text_frame.paragraphs[0].text = "Slide 84: OMML — Accent + Bar + SubSup"
title84.text_frame.paragraphs[0].font.size = Pt(24)
title84.text_frame.paragraphs[0].font.bold = True

math84 = slide84.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
math84.text_frame.paragraphs[0].text = "ACC_BAR_PLACEHOLDER"
math84.text_frame.paragraphs[0].font.size = Pt(24)

# ── Slide 85: Blur effect (a:blur) ───────────────────────────────────────────
slide85 = prs.slides.add_slide(blank)

title85 = slide85.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title85.text_frame.paragraphs[0].text = "Slide 85: Blur Effect (a:blur)"
title85.text_frame.paragraphs[0].font.size = Pt(24)
title85.text_frame.paragraphs[0].font.bold = True

# ── Slide 86: Preset shadow (a:prstShdw) ─────────────────────────────────────
slide86 = prs.slides.add_slide(blank)

title86 = slide86.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title86.text_frame.paragraphs[0].text = "Slide 86: Preset Shadow (a:prstShdw)"
title86.text_frame.paragraphs[0].font.size = Pt(24)
title86.text_frame.paragraphs[0].font.bold = True

# ── Slide 87: Fill overlay (a:fillOverlay) ────────────────────────────────────
slide87 = prs.slides.add_slide(blank)

title87 = slide87.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title87.text_frame.paragraphs[0].text = "Slide 87: Fill Overlay (a:fillOverlay)"
title87.text_frame.paragraphs[0].font.size = Pt(24)
title87.text_frame.paragraphs[0].font.bold = True

# Add shapes with solid fills that will get fillOverlay via patching
fo_labels = ["over", "mult", "screen", "darken", "lighten"]
fo_colors = ["0000FF", "00AA00", "FF6600", "8800CC", "CC0000"]
for i, (label, base_clr) in enumerate(zip(fo_labels, fo_colors)):
    left = Inches(0.5 + i * 1.8)
    box = slide87.shapes.add_textbox(left, Inches(1.5), Inches(1.5), Inches(1.5))
    box.text_frame.paragraphs[0].text = label
    box.text_frame.paragraphs[0].font.size = Pt(14)
    box.text_frame.paragraphs[0].font.bold = True
    from pptx.util import Pt as PtUtil
    from pptx.dml.color import RGBColor
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(
        int(base_clr[0:2], 16), int(base_clr[2:4], 16), int(base_clr[4:6], 16)
    )

# ── Slide 88: Justified text (algn="just") ────────────────────────────────────
slide88 = prs.slides.add_slide(blank)

title88 = slide88.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title88.text_frame.paragraphs[0].text = "Slide 88: Justified Text (algn=\"just\")"
title88.text_frame.paragraphs[0].font.size = Pt(24)
title88.text_frame.paragraphs[0].font.bold = True

# Add text box with justified paragraphs
from pptx.enum.text import PP_ALIGN
just_box = slide88.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(3))
tf = just_box.text_frame
tf.word_wrap = True
p1 = tf.paragraphs[0]
p1.text = "This is a justified paragraph with enough text to wrap across multiple lines. The word spacing should be adjusted so that each line extends to fill the full width of the text box evenly."
p1.alignment = PP_ALIGN.JUSTIFY
p1.font.size = Pt(16)

p2 = tf.add_paragraph()
p2.text = "Short last line."
p2.alignment = PP_ALIGN.JUSTIFY
p2.font.size = Pt(16)

# Japanese justified text (CJK — uses letter-spacing fallback)
just_box2 = slide88.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(2))
tf2 = just_box2.text_frame
tf2.word_wrap = True
p3 = tf2.paragraphs[0]
p3.text = "均等割り付けのテスト文章です。日本語テキストはスペースがないため文字間隔で調整されます。"
p3.alignment = PP_ALIGN.JUSTIFY
p3.font.size = Pt(16)

