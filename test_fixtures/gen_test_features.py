#!/usr/bin/env python3
"""
Generate a test PPTX that exercises all recently implemented features.

Slides:
  1. Text body vertical alignment (top / center / bottom)
  2. Paragraph spacing (spcBef / spcAft) + indent (marL / indent)
  3. Bullet characters + auto-numbering
  4. Run decorations (underline / strikethrough / superscript / subscript)
  5. Body insets (lIns / tIns / rIns / bIns)
  6. Master/Layout inheritance: slide with NO explicit bg (inherits master bg)
  7. Master/Layout inheritance: placeholder shapes with text style defaults
  8. Master/Layout inheritance: slide WITH explicit bg (overrides master)
  9. East Asian fonts + font theme references (+mj-lt/+mn-lt/+mj-ea/+mn-ea)
 10. Line spacing (a:lnSpc) — EMU and percentage
 11. Character spacing (a:rPr spc) + lstStyle
 12. normAutofit (fontScale / lnSpcReduction)
 13. Text wrapping (long text, CJK, wrap="none")
 14. Bullet formatting (a:buFont, a:buSzPct, a:buSzPts, a:buClr)
 15. Capitalization (a:rPr cap="all" / "small")
 16. Color map override (p:clrMapOvr — dark bg + light text via bg1↔dk1 swap)
 17. CS/Sym fonts + kerning (a:cs, a:sym, kern attr)
 18. Text rotation (bodyPr rot) + tab stops (a:tabLst)
 19. Vertical text (bodyPr vert) + text columns (numCol/spcCol)
 20. Hyperlink (a:hlinkClick) + RTL (a:pPr rtl)
 21. Image bullet (a:buBlip)
 22. Hover link + link color (a:hlinkHover)
 23. Linear gradient fills (a:gradFill + a:lin at 0°/90°/45°)
 24. Radial/path gradient fills (a:path circle/rect) + gradient on ellipse
 25. Gradient background (p:bg → p:bgPr → a:gradFill)
 26. Alpha/transparency (semi-transparent solid fill, semi-transparent gradient stop)
 27. Image fill on AutoShape (a:blipFill inside p:spPr)
 28. Pattern fill (a:pattFill — ltDnDiag / smCheck / dkHorz)
 29. Gradient tileFlip (tileFlip="x" / "y" / "xy")
 30. Additional pattern fills (pct50 / dnDiag / cross / lgCheck / solidDmnd / trellis)
 31. Image fill tile (a:tile with sx/sy/flip/algn)
 32. Stroke dash styles (dash/dot/dashDot/lgDash/sysDot/sysDash + custDash)
 33. Arrows (headEnd/tailEnd), line join (round/bevel/miter), line cap (rnd/sq), compound line (dbl), noFill
 34. Group shapes (p:grpSp) — simple group + nested group
 35. Connectors (p:cxnSp) — straight, diagonal, bent, curved
 36. Preset geometry shapes (triangle, diamond, pentagon, hexagon, arrow, star, heart, etc.)
 37. Custom geometry (a:custGeom) — freeform shapes with guide formulas and paths
 38. Gear shapes (gear6, gear9) — preset geometry with accurate tooth paths
 39. Text rectangles (a:rect) — text positioned within non-rectangular shapes
 40. Connection points (a:cxnLst) — custom geometry with connection points + connector refs
 41. Table cell merge (gridSpan + rowSpan/vMerge) + borders + margins + anchor
 42. Table diagonal borders + tblPr flags (firstRow/lastRow/bandRow/bandCol) + tblStyleId
 43. Image crop (srcRect) + alpha (alphaModFix) — p:pic crop/alpha, AutoShape blipFill crop
 44. External image reference (TargetMode="External") — Wikimedia + picsum.photos
 45. Image effects — brightness/contrast (a:lum bright/contrast)
 46. Duotone (a:duotone) + color change (a:clrChange)
 47. Background pattern fill (a:pattFill in p:bgPr)
 48. Line gradient fill (a:gradFill in a:ln) + line pattern fill (a:pattFill in a:ln)
 49. Shape hyperlinks (a:hlinkClick/a:hlinkHover on p:cNvPr) + color modifiers (comp/inv/hueMod)
 50. Shape effects (a:effectLst — outerShdw / innerShdw / glow / softEdge)
 51. 3D effects (a:scene3d / a:sp3d / bevelT / extrusion) + text shadow/glow (a:rPr/a:effectLst)
 52. Column chart (clustered, 2 series)
 53. Line chart (2 series, with markers)
 54. Pie chart (single series, with legend)
 55. Bar chart (horizontal) + Donut chart
 56. Column chart with data labels (c:dLbls showVal)
 57. Pie chart with custom data point colors (c:dPt) + percentage labels
 58. Line chart with series spPr colors + data labels
 59. Scatter chart + Area chart
 60. Radar chart
 61. Bubble chart (c:bubbleChart with xVal/yVal/bubbleSize)
 62. Stock chart (c:stockChart — Open/High/Low/Close)
 63. All chart types overview (mini charts for visual regression)
 64. Line chart with linear trendline (c:trendline)
 65. Column chart with error bars (c:errBars)
 66. Composite chart (column + line on same plot area)
 67. Surface chart (c:surfaceChart — 2D heatmap approximation + view3D)
 68. Pie of pie chart (c:ofPieChart — rendered as standard pie)
 69. 3D bar chart (c:bar3DChart — 2D rendering + view3D preserved)
 70. Text outline (a:rPr/a:ln — stroke on text glyphs)
 71. Text gradient fill (a:rPr/a:gradFill — gradient on text)
 72. Text warp (a:prstTxWarp — preset text warp with adjust values)
 72b. Additional text warp presets (slant, curve, inflate, chevron, circle)
 73. Stacked / Percent-stacked bar charts (BAR_STACKED_100, COLUMN_STACKED, COLUMN_STACKED_100)
 74. Speaker notes (p:notes) + comments (p:cmAuthorLst / p:cmLst)
 75. SmartArt fallback (mc:AlternateContent with mc:Choice + mc:Fallback group shapes)
 76. OLE embedded object (p:oleObj with fallback image in p:graphicFrame)
 77. Media (video with poster frame — a:videoFile in p:nvPr)
 78. Math equation (OMML m:oMathPara / m:oMath in text body)
 79. Transition + Timing (p:transition + p:timing round-trip preservation)
 80. Hidden slide (p:sld show="0")
 81. WMF image (WMF → SVG conversion via wmfToSvg)
"""

import base64
import io
import os
import re
import struct
import tempfile
import zipfile
import zlib

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # blank layout

# ── Slide 1: Vertical alignment ─────────────────────────────────────────────

slide1 = prs.slides.add_slide(blank)

def add_box(slide, left, top, width, height, text, anchor, fill_rgb=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text = text
    tf = txBox.text_frame
    tf.word_wrap = True
    # Set anchor (vertical alignment)
    if anchor == "top":
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        # anchor is set via xml manipulation below
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 0)
    # Set fill
    if fill_rgb:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*fill_rgb)
    # Set body anchor
    bodyPr = txBox.text_frame._txBody.find(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    anchor_map = {"top": "t", "center": "ctr", "bottom": "b"}
    bodyPr.set('anchor', anchor_map.get(anchor, "t"))
    return txBox

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_box.text_frame.paragraphs[0].text = "Slide 1: Vertical Alignment (top / center / bottom)"
title_box.text_frame.paragraphs[0].font.size = Pt(24)
title_box.text_frame.paragraphs[0].font.bold = True

add_box(slide1, Inches(0.5), Inches(1.2), Inches(2.8), Inches(2.5),
        "TOP aligned text\n(anchor=t)", "top", (220, 230, 240))
add_box(slide1, Inches(3.6), Inches(1.2), Inches(2.8), Inches(2.5),
        "CENTER aligned text\n(anchor=ctr)", "center", (200, 220, 240))
add_box(slide1, Inches(6.7), Inches(1.2), Inches(2.8), Inches(2.5),
        "BOTTOM aligned text\n(anchor=b)", "bottom", (180, 210, 240))

# Also add boxes with explicit insets
inset_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4), Inches(2.5))
inset_box.text_frame.paragraphs[0].text = "Large insets (lIns=36pt, tIns=36pt)"
inset_box.text_frame.paragraphs[0].font.size = Pt(12)
fill = inset_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 230, 200)
bodyPr = inset_box.text_frame._txBody.find(
    '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
bodyPr.set('lIns', str(Emu(Pt(36))))
bodyPr.set('tIns', str(Emu(Pt(36))))
bodyPr.set('rIns', str(Emu(Pt(36))))
bodyPr.set('bIns', str(Emu(Pt(36))))

# ── Slide 2: Paragraph spacing + indent ──────────────────────────────────────

slide2 = prs.slides.add_slide(blank)

title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title2.text_frame.paragraphs[0].text = "Slide 2: Paragraph Spacing & Indent"
title2.text_frame.paragraphs[0].font.size = Pt(24)
title2.text_frame.paragraphs[0].font.bold = True

box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
tf2 = box2.text_frame
tf2.word_wrap = True
fill2 = box2.fill
fill2.solid()
fill2.fore_color.rgb = RGBColor(245, 245, 250)

# First paragraph - normal
p0 = tf2.paragraphs[0]
p0.text = "Paragraph 1: No special spacing or indent."
p0.font.size = Pt(14)

nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# Second paragraph - spacing before 24pt
p1 = tf2.add_paragraph()
p1.text = "Paragraph 2: spcBef=24pt (gap above this line)"
p1.font.size = Pt(14)
pPr1 = p1._p.find('a:pPr', nsmap)
if pPr1 is None:
    pPr1 = etree.SubElement(p1._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p1._p.insert(0, pPr1)
spcBef = etree.SubElement(pPr1, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcBef')
spcPts = etree.SubElement(spcBef, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
spcPts.set('val', '2400')  # 24pt in hundredths

# Third paragraph - indent marL
p2 = tf2.add_paragraph()
p2.text = "Paragraph 3: marL=1in (indented from left)"
p2.font.size = Pt(14)
pPr2 = p2._p.find('a:pPr', nsmap)
if pPr2 is None:
    pPr2 = etree.SubElement(p2._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p2._p.insert(0, pPr2)
pPr2.set('marL', str(914400))  # 1 inch in EMU

# Fourth paragraph - spacing after 18pt
p3 = tf2.add_paragraph()
p3.text = "Paragraph 4: spcAft=18pt (gap below this line)"
p3.font.size = Pt(14)
pPr3 = p3._p.find('a:pPr', nsmap)
if pPr3 is None:
    pPr3 = etree.SubElement(p3._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p3._p.insert(0, pPr3)
spcAft = etree.SubElement(pPr3, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcAft')
spcPtsA = etree.SubElement(spcAft, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
spcPtsA.set('val', '1800')  # 18pt

# Fifth paragraph - after the gap
p4 = tf2.add_paragraph()
p4.text = "Paragraph 5: After the 18pt gap. Should be visibly separated from P4."
p4.font.size = Pt(14)

# ── Slide 3: Bullets ────────────────────────────────────────────────────────

slide3 = prs.slides.add_slide(blank)

title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title3.text_frame.paragraphs[0].text = "Slide 3: Bullets & Auto-Numbering"
title3.text_frame.paragraphs[0].font.size = Pt(24)
title3.text_frame.paragraphs[0].font.bold = True

# Character bullets
box3a = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2.5))
tf3a = box3a.text_frame
tf3a.word_wrap = True
fill3a = box3a.fill
fill3a.solid()
fill3a.fore_color.rgb = RGBColor(230, 245, 230)

bullet_items = [
    ("First bullet item", "\u2022"),   # •
    ("Second bullet item", "\u2022"),
    ("Third with dash", "\u2013"),      # –
    ("Fourth with arrow", "\u25B6"),    # ▶
]

for i, (text, char) in enumerate(bullet_items):
    if i == 0:
        p = tf3a.paragraphs[0]
    else:
        p = tf3a.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    pPr.set('marL', str(457200))   # 0.5in
    pPr.set('indent', str(-228600))  # -0.25in (hanging indent)
    buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
    buChar.set('char', char)

# Auto-numbering
box3b = slide3.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(2.5))
tf3b = box3b.text_frame
tf3b.word_wrap = True
fill3b = box3b.fill
fill3b.solid()
fill3b.fore_color.rgb = RGBColor(230, 230, 245)

auto_items = ["Arabic period", "Second item", "Third item", "Fourth item"]
for i, text in enumerate(auto_items):
    if i == 0:
        p = tf3b.paragraphs[0]
    else:
        p = tf3b.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    pPr.set('marL', str(457200))
    pPr.set('indent', str(-228600))
    buAutoNum = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
    buAutoNum.set('type', 'arabicPeriod')

# Roman numerals
box3c = slide3.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4), Inches(2.5))
tf3c = box3c.text_frame
tf3c.word_wrap = True
fill3c = box3c.fill
fill3c.solid()
fill3c.fore_color.rgb = RGBColor(245, 240, 230)

roman_items = ["Roman upper item", "Another item", "Third one"]
for i, text in enumerate(roman_items):
    if i == 0:
        p = tf3c.paragraphs[0]
    else:
        p = tf3c.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    pPr.set('marL', str(457200))
    pPr.set('indent', str(-228600))
    buAutoNum = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
    buAutoNum.set('type', 'romanUcPeriod')

# Alpha lowercase
box3d = slide3.shapes.add_textbox(Inches(5), Inches(4.2), Inches(4.5), Inches(2.5))
tf3d = box3d.text_frame
tf3d.word_wrap = True
fill3d = box3d.fill
fill3d.solid()
fill3d.fore_color.rgb = RGBColor(240, 245, 245)

alpha_items = ["Alpha lower a", "Alpha lower b", "Alpha lower c"]
for i, text in enumerate(alpha_items):
    if i == 0:
        p = tf3d.paragraphs[0]
    else:
        p = tf3d.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    pPr.set('marL', str(457200))
    pPr.set('indent', str(-228600))
    buAutoNum = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
    buAutoNum.set('type', 'alphaLcParenR')

# ── Slide 4: Run decorations ───────────────────────────────────────────────

slide4 = prs.slides.add_slide(blank)

title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title4.text_frame.paragraphs[0].text = "Slide 4: Run Decorations"
title4.text_frame.paragraphs[0].font.size = Pt(24)
title4.text_frame.paragraphs[0].font.bold = True

box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
tf4 = box4.text_frame
tf4.word_wrap = True

# Underline (single)
p_u = tf4.paragraphs[0]
run_u = p_u.add_run()
run_u.text = "This text has single underline"
run_u.font.size = Pt(18)
run_u.font.underline = True

# Strikethrough
p_s = tf4.add_paragraph()
run_s = p_s.add_run()
run_s.text = "This text has strikethrough"
run_s.font.size = Pt(18)
# Set strike via XML
rPr = run_s._r.find('a:rPr', nsmap)
if rPr is None:
    rPr = etree.SubElement(run_s._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run_s._r.insert(0, rPr)
rPr.set('strike', 'sngStrike')

# Both underline and strikethrough
p_us = tf4.add_paragraph()
run_us = p_us.add_run()
run_us.text = "Both underline AND strikethrough"
run_us.font.size = Pt(18)
run_us.font.underline = True
rPr_us = run_us._r.find('a:rPr', nsmap)
rPr_us.set('strike', 'sngStrike')

# Superscript
p_sup = tf4.add_paragraph()
run_normal = p_sup.add_run()
run_normal.text = "E = mc"
run_normal.font.size = Pt(18)
run_super = p_sup.add_run()
run_super.text = "2"
run_super.font.size = Pt(18)
rPr_sup = run_super._r.find('a:rPr', nsmap)
if rPr_sup is None:
    rPr_sup = etree.SubElement(run_super._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run_super._r.insert(0, rPr_sup)
rPr_sup.set('baseline', '30000')

# Subscript
p_sub = tf4.add_paragraph()
run_h2 = p_sub.add_run()
run_h2.text = "H"
run_h2.font.size = Pt(18)
run_2 = p_sub.add_run()
run_2.text = "2"
run_2.font.size = Pt(18)
rPr_2 = run_2._r.find('a:rPr', nsmap)
if rPr_2 is None:
    rPr_2 = etree.SubElement(run_2._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run_2._r.insert(0, rPr_2)
rPr_2.set('baseline', '-25000')
run_o = p_sub.add_run()
run_o.text = "O (water)"
run_o.font.size = Pt(18)

# Mixed decorations in one paragraph
p_mix = tf4.add_paragraph()
run_b = p_mix.add_run()
run_b.text = "Bold "
run_b.font.bold = True
run_b.font.size = Pt(18)
run_i = p_mix.add_run()
run_i.text = "Italic "
run_i.font.italic = True
run_i.font.size = Pt(18)
run_bi = p_mix.add_run()
run_bi.text = "BoldItalicUnderline"
run_bi.font.bold = True
run_bi.font.italic = True
run_bi.font.underline = True
run_bi.font.size = Pt(18)

# ── Slide 5: Combined test ──────────────────────────────────────────────────

slide5 = prs.slides.add_slide(blank)
# Dark background
bg = slide5.background
fill_bg = bg.fill
fill_bg.solid()
fill_bg.fore_color.rgb = RGBColor(27, 58, 107)

title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf5t = title5.text_frame
tf5t.paragraphs[0].text = "Slide 5: Combined Features"
tf5t.paragraphs[0].font.size = Pt(28)
tf5t.paragraphs[0].font.bold = True
tf5t.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
bodyPr5 = tf5t._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
bodyPr5.set('anchor', 'ctr')

# Box with center-anchored, spaced, bulleted content
box5 = slide5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
tf5 = box5.text_frame
tf5.word_wrap = True
fill5 = box5.fill
fill5.solid()
fill5.fore_color.rgb = RGBColor(240, 240, 250)

bodyPr5b = tf5._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
bodyPr5b.set('anchor', 'ctr')
bodyPr5b.set('lIns', str(Emu(Pt(18))))
bodyPr5b.set('tIns', str(Emu(Pt(18))))

items5 = [
    "Center-anchored box with bullets",
    "Second item with spacing",
    "Third item with underline run",
]
for i, text in enumerate(items5):
    if i == 0:
        p = tf5.paragraphs[0]
    else:
        p = tf5.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    pPr.set('marL', str(457200))
    pPr.set('indent', str(-228600))
    buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
    buChar.set('char', '\u2022')
    # Add spacing before
    spcBef = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcBef')
    spcPts = etree.SubElement(spcBef, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
    spcPts.set('val', '1200')  # 12pt

# ── Set up slide master for inheritance testing ────────────────────────────
# Modify the slide master to have:
#   - A dark navy background (#1B3A6B)
#   - Title style: 44pt, centered, white text
#   - Body style: 24pt, left-aligned, light gray text

sld_master = prs.slide_masters[0]
nsmap_p = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
           'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
           'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

# Set master background
master_xml = sld_master._element
# Remove existing bg if any
for old_bg in master_xml.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}bg'):
    master_xml.remove(old_bg)
# Find cSld or create it
cSld = master_xml.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
if cSld is None:
    cSld = etree.SubElement(master_xml, '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
# Remove existing bg from cSld
for old_bg in cSld.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}bg'):
    cSld.remove(old_bg)
# Add background: navy blue
bg_elem = etree.SubElement(cSld, '{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
cSld.insert(0, bg_elem)  # bg should be first child of cSld
bgPr = etree.SubElement(bg_elem, '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr')
solidFill_bg = etree.SubElement(bgPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
srgbClr_bg = etree.SubElement(solidFill_bg, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
srgbClr_bg.set('val', '1B3A6B')
effectLst = etree.SubElement(bgPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')

# Set up txStyles on the master
# Remove existing txStyles if any
for old_ts in master_xml.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}txStyles'):
    master_xml.remove(old_ts)

txStyles = etree.SubElement(master_xml, '{http://schemas.openxmlformats.org/presentationml/2006/main}txStyles')

# Title style: 44pt, centered, white
titleStyle = etree.SubElement(txStyles, '{http://schemas.openxmlformats.org/presentationml/2006/main}titleStyle')
lvl1pPr_t = etree.SubElement(titleStyle, '{http://schemas.openxmlformats.org/drawingml/2006/main}lvl1pPr')
lvl1pPr_t.set('algn', 'ctr')
defRPr_t = etree.SubElement(lvl1pPr_t, '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
defRPr_t.set('sz', '4400')  # 44pt
solidFill_t = etree.SubElement(defRPr_t, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
srgbClr_t = etree.SubElement(solidFill_t, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
srgbClr_t.set('val', 'FFFFFF')

# Body style: lvl1=24pt left white, lvl2=20pt left lightgray
bodyStyle = etree.SubElement(txStyles, '{http://schemas.openxmlformats.org/presentationml/2006/main}bodyStyle')
lvl1pPr_b = etree.SubElement(bodyStyle, '{http://schemas.openxmlformats.org/drawingml/2006/main}lvl1pPr')
lvl1pPr_b.set('algn', 'l')
lvl1pPr_b.set('marL', '342900')
buChar_b1 = etree.SubElement(lvl1pPr_b, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
buChar_b1.set('char', '\u2022')
defRPr_b1 = etree.SubElement(lvl1pPr_b, '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
defRPr_b1.set('sz', '2400')  # 24pt
solidFill_b1 = etree.SubElement(defRPr_b1, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
srgbClr_b1 = etree.SubElement(solidFill_b1, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
srgbClr_b1.set('val', 'FFFFFF')

lvl2pPr_b = etree.SubElement(bodyStyle, '{http://schemas.openxmlformats.org/drawingml/2006/main}lvl2pPr')
lvl2pPr_b.set('algn', 'l')
lvl2pPr_b.set('marL', '685800')
buChar_b2 = etree.SubElement(lvl2pPr_b, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
buChar_b2.set('char', '\u2013')
defRPr_b2 = etree.SubElement(lvl2pPr_b, '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
defRPr_b2.set('sz', '2000')  # 20pt
solidFill_b2 = etree.SubElement(defRPr_b2, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
srgbClr_b2 = etree.SubElement(solidFill_b2, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
srgbClr_b2.set('val', 'CCCCCC')

# Other style: 18pt, gray
otherStyle = etree.SubElement(txStyles, '{http://schemas.openxmlformats.org/presentationml/2006/main}otherStyle')
lvl1pPr_o = etree.SubElement(otherStyle, '{http://schemas.openxmlformats.org/drawingml/2006/main}lvl1pPr')
defRPr_o = etree.SubElement(lvl1pPr_o, '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
defRPr_o.set('sz', '1800')
solidFill_o = etree.SubElement(defRPr_o, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
srgbClr_o = etree.SubElement(solidFill_o, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
srgbClr_o.set('val', '999999')

# ── Slide 6: No explicit background — should inherit navy from master ──────
slide6 = prs.slides.add_slide(blank)
# Do NOT set any background on this slide

title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf6t = title6.text_frame
tf6t.paragraphs[0].text = "Slide 6: Inherits master navy bg"
tf6t.paragraphs[0].font.size = Pt(24)
tf6t.paragraphs[0].font.bold = True
tf6t.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

info6 = slide6.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3))
tf6i = info6.text_frame
tf6i.word_wrap = True
p6 = tf6i.paragraphs[0]
p6.text = "Navy bg (#1B3A6B) = inheritance OK"
p6.font.size = Pt(18)
p6.font.color.rgb = RGBColor(200, 220, 255)
p6b = tf6i.add_paragraph()
p6b.text = "White bg = inheritance FAILED"
p6b.font.size = Pt(18)
p6b.font.color.rgb = RGBColor(255, 150, 150)

# ── Slide 7: Placeholder shapes with text style defaults ──────────────────
# Use title+content layout (index 1) to get real placeholders
title_content_layout = prs.slide_layouts[1]  # Title and Content layout
slide7 = prs.slides.add_slide(title_content_layout)

# The title placeholder should pick up master titleStyle (44pt ctr white)
title_ph = slide7.placeholders[0]  # Title placeholder
title_ph.text = "Slide 7: Placeholder Title (inherits 44pt ctr white)"

# The content placeholder should pick up master bodyStyle
body_ph = slide7.placeholders[1]  # Content placeholder
body_ph.text_frame.clear()
p7a = body_ph.text_frame.paragraphs[0]
p7a.text = "Level 0: Should inherit 24pt white from master bodyStyle"
# Don't set font size or color — let it inherit!

p7b = body_ph.text_frame.add_paragraph()
p7b.text = "Level 1: Should inherit 20pt lightgray from master bodyStyle lvl2"
p7b.level = 1

p7c = body_ph.text_frame.add_paragraph()
p7c.text = "Level 0 again: back to 24pt white"

# Also add a non-placeholder textbox for comparison
info7 = slide7.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
tf7i = info7.text_frame
tf7i.word_wrap = True
p7info = tf7i.paragraphs[0]
p7info.text = "^ Title should be 44pt centered white. Body bullets should be 24pt/20pt white/gray."
p7info.font.size = Pt(12)
p7info.font.color.rgb = RGBColor(180, 200, 255)

# ── Slide 8: Explicit background overrides master ─────────────────────────
slide8 = prs.slides.add_slide(blank)
# Set an explicit green background — should NOT inherit master navy
bg8 = slide8.background
fill_bg8 = bg8.fill
fill_bg8.solid()
fill_bg8.fore_color.rgb = RGBColor(34, 139, 34)  # Forest green

title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf8t = title8.text_frame
tf8t.paragraphs[0].text = "Slide 8: Explicit green bg (#228B22) — overrides master"
tf8t.paragraphs[0].font.size = Pt(20)
tf8t.paragraphs[0].font.bold = True
tf8t.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

info8 = slide8.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3))
tf8i = info8.text_frame
tf8i.word_wrap = True
p8 = tf8i.paragraphs[0]
p8.text = "This slide has an explicit green background. It should NOT show the master's navy."
p8.font.size = Pt(18)
p8.font.color.rgb = RGBColor(220, 255, 220)
p8b = tf8i.add_paragraph()
p8b.text = "If you see green, explicit bg override is working correctly."
p8b.font.size = Pt(18)
p8b.font.color.rgb = RGBColor(255, 255, 255)

# ── Set theme EA fonts ────────────────────────────────────────────────────
# Modify theme1.xml to add East Asian font entries
# Access theme via slide master's relationship
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
theme_part = None
for rel in sld_master.part.rels.values():
    if 'theme' in rel.reltype:
        theme_part = rel.target_part
        break

if theme_part is not None:
    theme_elem = etree.fromstring(theme_part.blob)
    font_scheme = theme_elem.find(f'.//{{{ns_a}}}fontScheme')
    if font_scheme is not None:
        # Add a:ea to majorFont
        major_font_elem = font_scheme.find(f'{{{ns_a}}}majorFont')
        if major_font_elem is not None:
            for old_ea in major_font_elem.findall(f'{{{ns_a}}}ea'):
                major_font_elem.remove(old_ea)
            ea_major = etree.SubElement(major_font_elem, f'{{{ns_a}}}ea')
            ea_major.set('typeface', 'Yu Gothic')
        # Add a:ea to minorFont
        minor_font_elem = font_scheme.find(f'{{{ns_a}}}minorFont')
        if minor_font_elem is not None:
            for old_ea in minor_font_elem.findall(f'{{{ns_a}}}ea'):
                minor_font_elem.remove(old_ea)
            ea_minor = etree.SubElement(minor_font_elem, f'{{{ns_a}}}ea')
            ea_minor.set('typeface', 'Yu Gothic')
    # Write back the modified theme XML
    theme_part._blob = etree.tostring(theme_elem, xml_declaration=True, encoding='UTF-8', standalone=True)
    print("Theme EA fonts set: Yu Gothic")

# ── Slide 9: East Asian fonts + font theme references ───────────────────
slide9 = prs.slides.add_slide(blank)

title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title9.text_frame.paragraphs[0].text = "Slide 9: East Asian Fonts & Theme References"
title9.text_frame.paragraphs[0].font.size = Pt(24)
title9.text_frame.paragraphs[0].font.bold = True

# Box with explicit EA font (a:ea typeface="MS PGothic")
box9a = slide9.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2))
tf9a = box9a.text_frame
tf9a.word_wrap = True
fill9a = box9a.fill
fill9a.solid()
fill9a.fore_color.rgb = RGBColor(240, 240, 255)
p9a = tf9a.paragraphs[0]
run9a = p9a.add_run()
run9a.text = "日本語テキスト（MS PGothic）"
run9a.font.size = Pt(18)
# Set a:ea typeface directly in XML
rPr9a = run9a._r.find('a:rPr', nsmap)
if rPr9a is None:
    rPr9a = etree.SubElement(run9a._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run9a._r.insert(0, rPr9a)
ea9a = etree.SubElement(rPr9a, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
ea9a.set('typeface', 'MS PGothic')
latin9a = etree.SubElement(rPr9a, '{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
latin9a.set('typeface', 'Arial')

# Second para: theme font reference +mj-ea and +mj-lt
p9b = tf9a.add_paragraph()
run9b = p9b.add_run()
run9b.text = "テーマ参照: +mj-ea / +mj-lt"
run9b.font.size = Pt(16)
rPr9b = run9b._r.find('a:rPr', nsmap)
if rPr9b is None:
    rPr9b = etree.SubElement(run9b._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run9b._r.insert(0, rPr9b)
ea9b = etree.SubElement(rPr9b, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
ea9b.set('typeface', '+mj-ea')
latin9b = etree.SubElement(rPr9b, '{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
latin9b.set('typeface', '+mj-lt')

# Third para: minor font theme reference +mn-ea and +mn-lt
p9c = tf9a.add_paragraph()
run9c = p9c.add_run()
run9c.text = "マイナーフォント: +mn-ea / +mn-lt"
run9c.font.size = Pt(14)
rPr9c = run9c._r.find('a:rPr', nsmap)
if rPr9c is None:
    rPr9c = etree.SubElement(run9c._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run9c._r.insert(0, rPr9c)
ea9c = etree.SubElement(rPr9c, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
ea9c.set('typeface', '+mn-ea')
latin9c = etree.SubElement(rPr9c, '{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
latin9c.set('typeface', '+mn-lt')

# Box with mixed Latin + EA text in same run
box9b = slide9.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(2))
tf9b = box9b.text_frame
tf9b.word_wrap = True
fill9b = box9b.fill
fill9b.solid()
fill9b.fore_color.rgb = RGBColor(255, 245, 235)
p9d = tf9b.paragraphs[0]
run9d = p9d.add_run()
run9d.text = "Mixed: ABC + あいう + 123"
run9d.font.size = Pt(18)
rPr9d = run9d._r.find('a:rPr', nsmap)
if rPr9d is None:
    rPr9d = etree.SubElement(run9d._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run9d._r.insert(0, rPr9d)
ea9d = etree.SubElement(rPr9d, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
ea9d.set('typeface', 'Meiryo')
latin9d = etree.SubElement(rPr9d, '{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
latin9d.set('typeface', 'Segoe UI')

# Info text
info9 = slide9.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
tf9info = info9.text_frame
tf9info.word_wrap = True
p9info = tf9info.paragraphs[0]
p9info.text = "Verify: EA fonts should appear in font-family. Theme refs (+mj-ea/+mn-ea) should resolve to Yu Gothic."
p9info.font.size = Pt(12)
p9info.font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 10: Line spacing ──────────────────────────────────────────────
slide10 = prs.slides.add_slide(blank)

title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title10.text_frame.paragraphs[0].text = "Slide 10: Line Spacing (a:lnSpc)"
title10.text_frame.paragraphs[0].font.size = Pt(24)
title10.text_frame.paragraphs[0].font.bold = True

# Box with percentage-based line spacing (150%)
box10a = slide10.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(3))
tf10a = box10a.text_frame
tf10a.word_wrap = True
fill10a = box10a.fill
fill10a.solid()
fill10a.fore_color.rgb = RGBColor(230, 245, 255)

for i, text in enumerate(["Line 1: 150% spacing", "Line 2: wider gap above", "Line 3: still 150%"]):
    if i == 0:
        p = tf10a.paragraphs[0]
    else:
        p = tf10a.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
    spcPct.set('val', '150000')  # 150%

# Box with point-based line spacing (36pt absolute)
box10b = slide10.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(3))
tf10b = box10b.text_frame
tf10b.word_wrap = True
fill10b = box10b.fill
fill10b.solid()
fill10b.fore_color.rgb = RGBColor(255, 240, 230)

for i, text in enumerate(["Line 1: 36pt absolute", "Line 2: fixed height", "Line 3: still 36pt"]):
    if i == 0:
        p = tf10b.paragraphs[0]
    else:
        p = tf10b.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    spcPts = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
    spcPts.set('val', '3600')  # 36pt

# Box with tight line spacing (80%)
box10c = slide10.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(4), Inches(2.5))
tf10c = box10c.text_frame
tf10c.word_wrap = True
fill10c = box10c.fill
fill10c.solid()
fill10c.fore_color.rgb = RGBColor(245, 255, 240)

for i, text in enumerate(["Tight: 80% spacing", "Lines should be close together", "Third tight line"]):
    if i == 0:
        p = tf10c.paragraphs[0]
    else:
        p = tf10c.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    lnSpc = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    spcPct = etree.SubElement(lnSpc, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
    spcPct.set('val', '80000')  # 80%

# ── Slide 11: Character spacing + lstStyle ──────────────────────────────
slide11 = prs.slides.add_slide(blank)

title11 = slide11.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title11.text_frame.paragraphs[0].text = "Slide 11: Character Spacing & lstStyle"
title11.text_frame.paragraphs[0].font.size = Pt(24)
title11.text_frame.paragraphs[0].font.bold = True

# Box with character spacing variations
box11a = slide11.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(2.5))
tf11a = box11a.text_frame
tf11a.word_wrap = True
fill11a = box11a.fill
fill11a.solid()
fill11a.fore_color.rgb = RGBColor(245, 240, 255)

# Normal spacing
p11_0 = tf11a.paragraphs[0]
run11_0 = p11_0.add_run()
run11_0.text = "Normal spacing (spc=0)"
run11_0.font.size = Pt(16)

# Wide spacing (spc=300 = 3pt)
p11_1 = tf11a.add_paragraph()
run11_1 = p11_1.add_run()
run11_1.text = "Wide spacing (spc=300, 3pt)"
run11_1.font.size = Pt(16)
rPr11_1 = run11_1._r.find('a:rPr', nsmap)
if rPr11_1 is None:
    rPr11_1 = etree.SubElement(run11_1._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run11_1._r.insert(0, rPr11_1)
rPr11_1.set('spc', '300')

# Very wide spacing (spc=1000 = 10pt)
p11_2 = tf11a.add_paragraph()
run11_2 = p11_2.add_run()
run11_2.text = "Very wide (spc=1000, 10pt)"
run11_2.font.size = Pt(16)
rPr11_2 = run11_2._r.find('a:rPr', nsmap)
if rPr11_2 is None:
    rPr11_2 = etree.SubElement(run11_2._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run11_2._r.insert(0, rPr11_2)
rPr11_2.set('spc', '1000')

# Tight spacing (spc=-100 = -1pt)
p11_3 = tf11a.add_paragraph()
run11_3 = p11_3.add_run()
run11_3.text = "Tight spacing (spc=-100, -1pt)"
run11_3.font.size = Pt(16)
rPr11_3 = run11_3._r.find('a:rPr', nsmap)
if rPr11_3 is None:
    rPr11_3 = etree.SubElement(run11_3._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run11_3._r.insert(0, rPr11_3)
rPr11_3.set('spc', '-100')

# Box with lstStyle (shape-specific list style defaults)
box11b = slide11.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
tf11b = box11b.text_frame
tf11b.word_wrap = True
fill11b = box11b.fill
fill11b.solid()
fill11b.fore_color.rgb = RGBColor(255, 250, 235)

# Add a:lstStyle to the txBody with level defaults
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
txBody11b = tf11b._txBody
# Remove existing lstStyle
for old_ls in txBody11b.findall(f'{{{ns_a}}}lstStyle'):
    txBody11b.remove(old_ls)
lstStyle = etree.SubElement(txBody11b, f'{{{ns_a}}}lstStyle')
# Insert lstStyle after bodyPr (should be second child of txBody)
bodyPr11b = txBody11b.find(f'{{{ns_a}}}bodyPr')
if bodyPr11b is not None:
    idx = list(txBody11b).index(bodyPr11b) + 1
    txBody11b.remove(lstStyle)
    txBody11b.insert(idx, lstStyle)

# Level 1: 20pt, bold, dark blue
lvl1pPr_ls = etree.SubElement(lstStyle, f'{{{ns_a}}}lvl1pPr')
defRPr_ls1 = etree.SubElement(lvl1pPr_ls, f'{{{ns_a}}}defRPr')
defRPr_ls1.set('sz', '2000')  # 20pt
defRPr_ls1.set('b', '1')
solidFill_ls1 = etree.SubElement(defRPr_ls1, f'{{{ns_a}}}solidFill')
srgbClr_ls1 = etree.SubElement(solidFill_ls1, f'{{{ns_a}}}srgbClr')
srgbClr_ls1.set('val', '003366')
ea_ls1 = etree.SubElement(defRPr_ls1, f'{{{ns_a}}}ea')
ea_ls1.set('typeface', 'Meiryo')

p11_ls1 = tf11b.paragraphs[0]
p11_ls1.text = "lstStyle lvl1: inherits 20pt bold #003366 Meiryo"
# Don't set font size/color — let lstStyle provide defaults

p11_ls2 = tf11b.add_paragraph()
p11_ls2.text = "Second paragraph: should also inherit lstStyle defaults"

# ── Slide 12: normAutofit (fontScale / lnSpcReduction) ───────────────────
slide12 = prs.slides.add_slide(blank)

title12 = slide12.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title12.text_frame.paragraphs[0].text = "Slide 12: normAutofit (fontScale / lnSpcReduction)"
title12.text_frame.paragraphs[0].font.size = Pt(24)
title12.text_frame.paragraphs[0].font.bold = True

def add_normAutofit_box(slide, left, top, width, height, label, lines,
                         fill_rgb, border_rgb, font_scale=None, ln_spc_reduction=None):
    """Add a textbox with a:normAutofit in bodyPr."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*fill_rgb)
    # Border
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    sp_pr = box._element.find(f'{{{ns_a}}}spPr')
    if sp_pr is None:
        sp_pr = box._element.find(f'{{{ns_p}}}spPr')
    if sp_pr is None:
        # textbox shapes: spPr is direct child of p:sp
        sp_pr = etree.SubElement(box._element, f'{{{ns_a}}}spPr')
    ln_elem = etree.SubElement(sp_pr, f'{{{ns_a}}}ln')
    ln_elem.set('w', '12700')
    ln_fill = etree.SubElement(ln_elem, f'{{{ns_a}}}solidFill')
    ln_clr = etree.SubElement(ln_fill, f'{{{ns_a}}}srgbClr')
    ln_clr.set('val', '%02X%02X%02X' % border_rgb)

    # Title line
    p0 = tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = label
    run0.font.size = Pt(20)
    run0.font.bold = True

    # Content lines
    for line in lines:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(16)

    # Set bodyPr with normAutofit
    bodyPr = tf._txBody.find(f'{{{ns_a}}}bodyPr')
    bodyPr.set('lIns', '91440')
    bodyPr.set('tIns', '45720')
    bodyPr.set('rIns', '91440')
    bodyPr.set('bIns', '45720')

    if font_scale is not None:
        naf = etree.SubElement(bodyPr, f'{{{ns_a}}}normAutofit')
        if font_scale != 100000:
            naf.set('fontScale', str(font_scale))
        if ln_spc_reduction is not None and ln_spc_reduction > 0:
            naf.set('lnSpcReduction', str(ln_spc_reduction))

    return box

# Shape 1: fontScale=80000 (80%)
add_normAutofit_box(slide12, Inches(0.5), Inches(1.2), Inches(4), Inches(1.8),
    "fontScale=80%",
    ["Line 2: scaled down", "Line 3: smaller text", "Line 4: fits in box"],
    (232, 234, 246), (63, 81, 181),
    font_scale=80000)

# Shape 2: fontScale=62500 + lnSpcReduction=20000
add_normAutofit_box(slide12, Inches(5), Inches(1.2), Inches(4), Inches(1.8),
    "62.5% + lnSpc 20%",
    ["Line 2: even smaller", "Line 3: tighter spacing", "Line 4: compact text"],
    (255, 243, 224), (255, 152, 0),
    font_scale=62500, ln_spc_reduction=20000)

# Shape 3: normAutofit default (100%, no fontScale attr)
add_normAutofit_box(slide12, Inches(0.5), Inches(3.5), Inches(4), Inches(1.8),
    "normAutofit default",
    ["Line 2: no scaling", "Line 3: same as original"],
    (232, 245, 233), (76, 175, 80),
    font_scale=100000)

# Shape 4: no autofit (reference)
add_normAutofit_box(slide12, Inches(5), Inches(3.5), Inches(4), Inches(1.8),
    "No autofit (ref)",
    ["Line 2: original size", "Line 3: may overflow"],
    (243, 229, 245), (156, 39, 176),
    font_scale=None)  # no normAutofit

# Info label at bottom
info12 = slide12.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
tf12info = info12.text_frame
tf12info.word_wrap = True
p12info = tf12info.paragraphs[0]
p12info.text = "Compare: top-left (80%) should be visibly smaller than bottom-right (no autofit). Top-right (62.5%+20%) should be smallest and tightest."
p12info.font.size = Pt(12)
p12info.font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 13: Text wrapping ───────────────────────────────────────────────
slide13 = prs.slides.add_slide(blank)

title13 = slide13.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title13.text_frame.paragraphs[0].text = "Slide 13: Text Wrapping"
title13.text_frame.paragraphs[0].font.size = Pt(24)
title13.text_frame.paragraphs[0].font.bold = True

# Box with long Latin text that should wrap
box13a = slide13.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2.5))
tf13a = box13a.text_frame
tf13a.word_wrap = True
fill13a = box13a.fill
fill13a.solid()
fill13a.fore_color.rgb = RGBColor(230, 240, 255)
p13a = tf13a.paragraphs[0]
p13a.text = "This is a long paragraph of Latin text that should automatically wrap at the shape boundary. Word-level wrapping should break at spaces."
p13a.font.size = Pt(16)

# Second paragraph: even longer
p13a2 = tf13a.add_paragraph()
p13a2.text = "Another paragraph with a very long sentence that exercises the wrapping algorithm more thoroughly, including multiple lines of content."
p13a2.font.size = Pt(14)

# Box with CJK text that should wrap character by character
box13b = slide13.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(2.5))
tf13b = box13b.text_frame
tf13b.word_wrap = True
fill13b = box13b.fill
fill13b.solid()
fill13b.fore_color.rgb = RGBColor(255, 245, 230)
p13b = tf13b.paragraphs[0]
run13b = p13b.add_run()
run13b.text = "日本語テキストの折り返しテスト。文字単位で折り返されるべきです。漢字とひらがなが混在しています。"
run13b.font.size = Pt(16)
rPr13b = run13b._r.find('a:rPr', nsmap)
if rPr13b is None:
    rPr13b = etree.SubElement(run13b._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run13b._r.insert(0, rPr13b)
ea13b = etree.SubElement(rPr13b, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
ea13b.set('typeface', 'Yu Gothic')

# Box with mixed Latin + CJK text
p13b2 = tf13b.add_paragraph()
run13b2 = p13b2.add_run()
run13b2.text = "Mixed テキスト: ABC あいう 123 日本語 English words and 漢字が混在"
run13b2.font.size = Pt(14)
rPr13b2 = run13b2._r.find('a:rPr', nsmap)
if rPr13b2 is None:
    rPr13b2 = etree.SubElement(run13b2._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run13b2._r.insert(0, rPr13b2)
ea13b2 = etree.SubElement(rPr13b2, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
ea13b2.set('typeface', 'Yu Gothic')

# Box with wrap="none" — text should NOT wrap
box13c = slide13.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4), Inches(1.5))
tf13c = box13c.text_frame
tf13c.word_wrap = False  # This sets wrap="none"
fill13c = box13c.fill
fill13c.solid()
fill13c.fore_color.rgb = RGBColor(255, 230, 230)
p13c = tf13c.paragraphs[0]
p13c.text = "wrap=none: This text should NOT wrap even though it is long and exceeds the box width"
p13c.font.size = Pt(16)
# Ensure wrap="none" is set in bodyPr
bodyPr13c = tf13c._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
bodyPr13c.set('wrap', 'none')

# Narrow box to test wrapping with large font
box13d = slide13.shapes.add_textbox(Inches(5), Inches(4.2), Inches(2.5), Inches(2.5))
tf13d = box13d.text_frame
tf13d.word_wrap = True
fill13d = box13d.fill
fill13d.solid()
fill13d.fore_color.rgb = RGBColor(230, 255, 230)
p13d = tf13d.paragraphs[0]
p13d.text = "Large font in narrow box"
p13d.font.size = Pt(28)
p13d2 = tf13d.add_paragraph()
p13d2.text = "Should wrap to multiple lines"
p13d2.font.size = Pt(20)

# ── Slide 14: Bullet formatting ──────────────────────────────────────────
slide14 = prs.slides.add_slide(blank)

title14 = slide14.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title14.text_frame.paragraphs[0].text = "Slide 14: Bullet Formatting (buFont / buSzPct / buSzPts / buClr)"
title14.text_frame.paragraphs[0].font.size = Pt(22)
title14.text_frame.paragraphs[0].font.bold = True

# Box with custom bullet font (Wingdings)
box14a = slide14.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2.5))
tf14a = box14a.text_frame
tf14a.word_wrap = True
fill14a = box14a.fill
fill14a.solid()
fill14a.fore_color.rgb = RGBColor(240, 235, 255)

bullet_format_items = [
    ("Wingdings bullet font", "Wingdings", "\u006C"),       # Wingdings 'l' = checkmark-like
    ("Symbol bullet font", "Symbol", "\u00B7"),              # Symbol middle dot
    ("Default font bullet", "", "\u2022"),                    # Regular bullet
]

for i, (text, bu_font, bu_char) in enumerate(bullet_format_items):
    if i == 0:
        p = tf14a.paragraphs[0]
    else:
        p = tf14a.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    pPr.set('marL', str(457200))
    pPr.set('indent', str(-228600))
    if bu_font:
        buFontElem = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
        buFontElem.set('typeface', bu_font)
    buChar14 = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
    buChar14.set('char', bu_char)

# Box with bullet size (percentage + points)
box14b = slide14.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(2.5))
tf14b = box14b.text_frame
tf14b.word_wrap = True
fill14b = box14b.fill
fill14b.solid()
fill14b.fore_color.rgb = RGBColor(230, 255, 240)

bu_size_items = [
    ("buSzPct=150000 (150%)", 150000, None),
    ("buSzPct=75000 (75%)", 75000, None),
    ("buSzPts=3200 (32pt)", None, 3200),
    ("buSzPts=800 (8pt)", None, 800),
]

for i, (text, sz_pct, sz_pts) in enumerate(bu_size_items):
    if i == 0:
        p = tf14b.paragraphs[0]
    else:
        p = tf14b.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    pPr.set('marL', str(457200))
    pPr.set('indent', str(-228600))
    buChar14b = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
    buChar14b.set('char', '\u2022')
    if sz_pct is not None:
        buSzPct = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct')
        buSzPct.set('val', str(sz_pct))
    if sz_pts is not None:
        buSzPts = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPts')
        buSzPts.set('val', str(sz_pts))

# Box with bullet color
box14c = slide14.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4), Inches(2.5))
tf14c = box14c.text_frame
tf14c.word_wrap = True
fill14c = box14c.fill
fill14c.solid()
fill14c.fore_color.rgb = RGBColor(255, 245, 235)

bu_color_items = [
    ("Red bullet", "FF0000"),
    ("Green bullet", "00AA00"),
    ("Blue bullet", "0000FF"),
    ("No bullet color (inherits)", None),
]

for i, (text, bu_clr) in enumerate(bu_color_items):
    if i == 0:
        p = tf14c.paragraphs[0]
    else:
        p = tf14c.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    pPr = p._p.find('a:pPr', nsmap)
    if pPr is None:
        pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        p._p.insert(0, pPr)
    pPr.set('marL', str(457200))
    pPr.set('indent', str(-228600))
    buChar14c = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
    buChar14c.set('char', '\u2022')
    if bu_clr is not None:
        buClrElem = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buClr')
        srgbClr14 = etree.SubElement(buClrElem, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgbClr14.set('val', bu_clr)

# Box with combined: buFont + buSzPct + buClr
box14d = slide14.shapes.add_textbox(Inches(5), Inches(4.2), Inches(4.5), Inches(2.5))
tf14d = box14d.text_frame
tf14d.word_wrap = True
fill14d = box14d.fill
fill14d.solid()
fill14d.fore_color.rgb = RGBColor(245, 240, 255)

p14d = tf14d.paragraphs[0]
p14d.text = "Combined: Wingdings + 120% + red"
p14d.font.size = Pt(16)
pPr14d = p14d._p.find('a:pPr', nsmap)
if pPr14d is None:
    pPr14d = etree.SubElement(p14d._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p14d._p.insert(0, pPr14d)
pPr14d.set('marL', str(457200))
pPr14d.set('indent', str(-228600))
buFont14d = etree.SubElement(pPr14d, '{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
buFont14d.set('typeface', 'Wingdings')
buSzPct14d = etree.SubElement(pPr14d, '{http://schemas.openxmlformats.org/drawingml/2006/main}buSzPct')
buSzPct14d.set('val', '120000')
buClr14d = etree.SubElement(pPr14d, '{http://schemas.openxmlformats.org/drawingml/2006/main}buClr')
srgbClr14d = etree.SubElement(buClr14d, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
srgbClr14d.set('val', 'FF0000')
buChar14d = etree.SubElement(pPr14d, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
buChar14d.set('char', '\u006C')

p14d2 = tf14d.add_paragraph()
p14d2.text = "Normal bullet for comparison"
p14d2.font.size = Pt(16)
pPr14d2 = p14d2._p.find('a:pPr', nsmap)
if pPr14d2 is None:
    pPr14d2 = etree.SubElement(p14d2._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p14d2._p.insert(0, pPr14d2)
pPr14d2.set('marL', str(457200))
pPr14d2.set('indent', str(-228600))
buChar14d2 = etree.SubElement(pPr14d2, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
buChar14d2.set('char', '\u2022')

# ── Slide 15: Capitalization (cap) ──────────────────────────────────────
slide15 = prs.slides.add_slide(blank)

title15 = slide15.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title15.text_frame.paragraphs[0].text = "Slide 15: Capitalization (a:rPr cap)"
title15.text_frame.paragraphs[0].font.size = Pt(24)
title15.text_frame.paragraphs[0].font.bold = True

# Box with cap="all"
box15a = slide15.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2.5))
tf15a = box15a.text_frame
tf15a.word_wrap = True
fill15a = box15a.fill
fill15a.solid()
fill15a.fore_color.rgb = RGBColor(230, 245, 255)

p15a = tf15a.paragraphs[0]
run15a = p15a.add_run()
run15a.text = "cap=all: This Should Be All Caps"
run15a.font.size = Pt(18)
rPr15a = run15a._r.find('a:rPr', nsmap)
if rPr15a is None:
    rPr15a = etree.SubElement(run15a._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run15a._r.insert(0, rPr15a)
rPr15a.set('cap', 'all')

p15a2 = tf15a.add_paragraph()
run15a2 = p15a2.add_run()
run15a2.text = "Mixed: "
run15a2.font.size = Pt(16)
run15a2b = p15a2.add_run()
run15a2b.text = "all caps part"
run15a2b.font.size = Pt(16)
rPr15a2b = run15a2b._r.find('a:rPr', nsmap)
if rPr15a2b is None:
    rPr15a2b = etree.SubElement(run15a2b._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run15a2b._r.insert(0, rPr15a2b)
rPr15a2b.set('cap', 'all')
run15a2c = p15a2.add_run()
run15a2c.text = " and normal part"
run15a2c.font.size = Pt(16)

p15a3 = tf15a.add_paragraph()
run15a3 = p15a3.add_run()
run15a3.text = "No cap: Regular text for comparison"
run15a3.font.size = Pt(16)

# Box with cap="small"
box15b = slide15.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(2.5))
tf15b = box15b.text_frame
tf15b.word_wrap = True
fill15b = box15b.fill
fill15b.solid()
fill15b.fore_color.rgb = RGBColor(255, 245, 230)

p15b = tf15b.paragraphs[0]
run15b = p15b.add_run()
run15b.text = "cap=small: Small Caps Text"
run15b.font.size = Pt(18)
rPr15b = run15b._r.find('a:rPr', nsmap)
if rPr15b is None:
    rPr15b = etree.SubElement(run15b._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run15b._r.insert(0, rPr15b)
rPr15b.set('cap', 'small')

p15b2 = tf15b.add_paragraph()
run15b2 = p15b2.add_run()
run15b2.text = "No cap: Regular text"
run15b2.font.size = Pt(16)

p15b3 = tf15b.add_paragraph()
run15b3 = p15b3.add_run()
run15b3.text = "Bold Small Caps"
run15b3.font.size = Pt(18)
run15b3.font.bold = True
rPr15b3 = run15b3._r.find('a:rPr', nsmap)
if rPr15b3 is None:
    rPr15b3 = etree.SubElement(run15b3._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    run15b3._r.insert(0, rPr15b3)
rPr15b3.set('cap', 'small')

# Info text
info15 = slide15.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
tf15info = info15.text_frame
tf15info.word_wrap = True
p15info = tf15info.paragraphs[0]
p15info.text = "cap=all: lowercase a-z should display as UPPERCASE A-Z. cap=small: should display with SVG font-variant=small-caps."
p15info.font.size = Pt(12)
p15info.font.color.rgb = RGBColor(100, 100, 100)

# ═══════════════════════════════════════════════════════════════════════════════
# Slide 16: Color map override (p:clrMapOvr)
# ═══════════════════════════════════════════════════════════════════════════════
slide16 = prs.slides.add_slide(prs.slide_layouts[0])

nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

# Remove existing p:clrMapOvr (python-pptx adds a:masterClrMapping by default)
existing_ovr = slide16._element.findall(
    '{http://schemas.openxmlformats.org/presentationml/2006/main}clrMapOvr'
)
for ovr in existing_ovr:
    slide16._element.remove(ovr)

# Add p:clrMapOvr with a:overrideClrMapping that swaps bg1↔tx1 (dk1↔lt1)
# This makes the slide use dark background + light text
clrMapOvr = etree.SubElement(
    slide16._element,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}clrMapOvr'
)
overrideMapping = etree.SubElement(
    clrMapOvr,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}overrideClrMapping',
    attrib={
        'bg1': 'dk1',
        'tx1': 'lt1',
        'bg2': 'lt2',
        'tx2': 'dk2',
        'accent1': 'accent1',
        'accent2': 'accent2',
        'accent3': 'accent3',
        'accent4': 'accent4',
        'accent5': 'accent5',
        'accent6': 'accent6',
        'hlink': 'hlink',
        'folHlink': 'folHlink',
    }
)

# Add a textbox with scheme-colored text (should resolve to light color on dark bg)
tb16 = slide16.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
tf16 = tb16.text_frame
tf16.word_wrap = True
p16 = tf16.paragraphs[0]
p16.text = "Color Map Override: bg1=dk1, tx1=lt1"
p16.font.size = Pt(28)
p16.font.bold = True
# Use scheme color tx1 (which should resolve to lt1 = white on this slide)
rPr16 = p16.runs[0]._r.find('a:rPr', nsmap)
if rPr16 is None:
    rPr16 = etree.SubElement(p16.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p16.runs[0]._r.insert(0, rPr16)
solidFill16 = etree.SubElement(rPr16, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
schemeClr16 = etree.SubElement(solidFill16, '{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr', attrib={'val': 'tx1'})

# Add a second textbox with explicit description
tb16b = slide16.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
tf16b = tb16b.text_frame
tf16b.word_wrap = True
p16b = tf16b.paragraphs[0]
p16b.text = "This slide has clrMapOvr swapping bg1/tx1. Text using scheme color tx1 should appear light (white) because tx1 maps to lt1."
p16b.font.size = Pt(16)
# Also use scheme color for this text
rPr16b = p16b.runs[0]._r.find('a:rPr', nsmap)
if rPr16b is None:
    rPr16b = etree.SubElement(p16b.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p16b.runs[0]._r.insert(0, rPr16b)
solidFill16b = etree.SubElement(rPr16b, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
schemeClr16b = etree.SubElement(solidFill16b, '{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr', attrib={'val': 'tx1'})

# Add slide background using scheme color bg1 (which maps to dk1 = dark)
bg16 = slide16._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld', nsmap)
if bg16 is None:
    bg16 = slide16._element.find('p:cSld', nsmap)
# Insert p:bg before p:cSld
pBg = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
bgPr = etree.SubElement(pBg, '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr')
solidFillBg = etree.SubElement(bgPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
# Use a dark solid color for background to make the slide visually dark
srgbClrBg = etree.SubElement(solidFillBg, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr', attrib={'val': '1A1A2E'})
# Insert bg element at the right position in slide XML
cSld = slide16._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
if cSld is not None:
    slide16._element.insert(list(slide16._element).index(cSld), pBg)

# ── Slide 17: CS/Sym fonts + kerning ─────────────────────────────────────────

slide17 = prs.slides.add_slide(blank)

# Textbox with Complex Script font (a:cs)
tb17a = slide17.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
tf17a = tb17a.text_frame
tf17a.word_wrap = True
p17a = tf17a.paragraphs[0]
p17a.text = "Complex Script Font (Arabic style)"
p17a.font.size = Pt(24)
p17a.font.bold = True
# Add a:cs element to rPr
rPr17a = p17a.runs[0]._r.find('a:rPr', nsmap)
if rPr17a is None:
    rPr17a = etree.SubElement(p17a.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p17a.runs[0]._r.insert(0, rPr17a)
etree.SubElement(rPr17a, '{http://schemas.openxmlformats.org/drawingml/2006/main}cs', attrib={'typeface': 'Arial'})

# Textbox with Symbol font (a:sym)
tb17b = slide17.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
tf17b = tb17b.text_frame
p17b = tf17b.paragraphs[0]
p17b.text = "Symbol Font Text"
p17b.font.size = Pt(24)
rPr17b = p17b.runs[0]._r.find('a:rPr', nsmap)
if rPr17b is None:
    rPr17b = etree.SubElement(p17b.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p17b.runs[0]._r.insert(0, rPr17b)
etree.SubElement(rPr17b, '{http://schemas.openxmlformats.org/drawingml/2006/main}sym', attrib={'typeface': 'Wingdings'})

# Textbox with kerning
tb17c = slide17.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
tf17c = tb17c.text_frame
p17c = tf17c.paragraphs[0]
p17c.text = "Kerning enabled at 1200 hundredths-pt"
p17c.font.size = Pt(24)
rPr17c = p17c.runs[0]._r.find('a:rPr', nsmap)
if rPr17c is None:
    rPr17c = etree.SubElement(p17c.runs[0]._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    p17c.runs[0]._r.insert(0, rPr17c)
rPr17c.set('kern', '1200')

# ── Slide 18: Text rotation + tab stops ──────────────────────────────────────

slide18 = prs.slides.add_slide(blank)

# Textbox with rotated text (bodyPr rot)
tb18a = slide18.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(2))
tf18a = tb18a.text_frame
p18a = tf18a.paragraphs[0]
p18a.text = "Rotated text (45 degrees)"
p18a.font.size = Pt(20)
# Set rot on bodyPr (45 degrees = 2700000 in 60000ths)
bodyPr18a = tb18a._element.find('.//a:bodyPr', nsmap)
bodyPr18a.set('rot', '2700000')

# Textbox with tab stops
tb18b = slide18.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(3))
tf18b = tb18b.text_frame
tf18b.word_wrap = True
p18b = tf18b.paragraphs[0]
p18b.text = "Col1\tCol2\tCol3"
p18b.font.size = Pt(18)
# Add tab stops via XML
pPr18b = p18b._pPr
if pPr18b is None:
    pPr18b = etree.SubElement(p18b._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p18b._p.insert(0, pPr18b)
tabLst = etree.SubElement(pPr18b, '{http://schemas.openxmlformats.org/drawingml/2006/main}tabLst')
etree.SubElement(tabLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}tab', attrib={'pos': '2743200', 'algn': 'l'})  # 3 inches
etree.SubElement(tabLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}tab', attrib={'pos': '5486400', 'algn': 'r'})  # 6 inches

# ── Slide 19: Vertical text + text columns ──────────────────────────────────

slide19 = prs.slides.add_slide(blank)

# Textbox with vertical text (vert="vert")
tb19a = slide19.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(5))
tf19a = tb19a.text_frame
p19a = tf19a.paragraphs[0]
p19a.text = "Vertical text"
p19a.font.size = Pt(20)
bodyPr19a = tb19a._element.find('.//a:bodyPr', nsmap)
bodyPr19a.set('vert', 'vert')

# Textbox with eaVert
tb19b = slide19.shapes.add_textbox(Inches(3), Inches(0.5), Inches(2), Inches(5))
tf19b = tb19b.text_frame
p19b = tf19b.paragraphs[0]
p19b.text = "EA Vertical"
p19b.font.size = Pt(20)
bodyPr19b = tb19b._element.find('.//a:bodyPr', nsmap)
bodyPr19b.set('vert', 'eaVert')

# Textbox with text columns (numCol + spcCol)
tb19c = slide19.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(4), Inches(5))
tf19c = tb19c.text_frame
tf19c.word_wrap = True
p19c = tf19c.paragraphs[0]
p19c.text = "This is column text that spans multiple columns. The text should flow from the first column to the second column when it reaches the bottom."
p19c.font.size = Pt(14)
bodyPr19c = tb19c._element.find('.//a:bodyPr', nsmap)
bodyPr19c.set('numCol', '2')
bodyPr19c.set('spcCol', '457200')  # 0.5 inch spacing

# ── Slide 20: Hyperlink + RTL ────────────────────────────────────────────────

slide20 = prs.slides.add_slide(blank)

# Textbox with hyperlink
tb20a = slide20.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
tf20a = tb20a.text_frame
tf20a.word_wrap = True
p20a = tf20a.paragraphs[0]
r20a = p20a.runs[0] if len(p20a.runs) > 0 else p20a.add_run()
r20a.text = "Click here to visit example.com"
r20a.font.size = Pt(20)
r20a.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
r20a.font.underline = True
# Add hlinkClick via XML — need to add a relationship first
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
rel = slide20.part.relate_to('https://example.com', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
rPr20a = r20a._r.find('a:rPr', nsmap)
if rPr20a is None:
    rPr20a = etree.SubElement(r20a._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    r20a._r.insert(0, rPr20a)
etree.SubElement(rPr20a, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rel})

# RTL paragraph
tb20b = slide20.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(2))
tf20b = tb20b.text_frame
tf20b.word_wrap = True
p20b = tf20b.paragraphs[0]
p20b.text = "RTL paragraph text (right-to-left)"
p20b.font.size = Pt(20)
pPr20b = p20b._pPr
if pPr20b is None:
    pPr20b = etree.SubElement(p20b._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p20b._p.insert(0, pPr20b)
pPr20b.set('rtl', '1')

# ── Slide 21: Image bullet (a:buBlip) ───────────────────────────────────────

slide21 = prs.slides.add_slide(blank)

# For image bullets, we need a small embedded image
# We'll create a tiny 1x1 red PNG in memory
def make_tiny_png(r, g, b):
    """Create a minimal 1x1 PNG."""
    # IHDR
    ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
    ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
    # IDAT
    raw = bytes([0, r, g, b])
    compressed = zlib.compress(raw)
    idat_crc = zlib.crc32(b'IDAT' + compressed) & 0xffffffff
    idat = struct.pack('>I', len(compressed)) + b'IDAT' + compressed + struct.pack('>I', idat_crc)
    # IEND
    iend_crc = zlib.crc32(b'IEND') & 0xffffffff
    iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
    return b'\x89PNG\r\n\x1a\n' + ihdr + idat + iend

png_data = make_tiny_png(255, 0, 0)  # red dot
png_stream = io.BytesIO(png_data)

# Add image as a relationship to the slide
from pptx.opc.package import Part
from pptx.opc.constants import CONTENT_TYPE as CT
image_part, rId_img = slide21.part.get_or_add_image_part(png_stream)

# Create textbox with bullet paragraphs
tb21 = slide21.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
tf21 = tb21.text_frame
tf21.word_wrap = True
p21 = tf21.paragraphs[0]
p21.text = "Image bullet paragraph 1"
p21.font.size = Pt(20)

# Add image bullet via XML: <a:buBlip><a:blip r:embed="rIdN"/></a:buBlip>
pPr21 = p21._pPr
if pPr21 is None:
    pPr21 = etree.SubElement(p21._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
    p21._p.insert(0, pPr21)
buBlip = etree.SubElement(pPr21, '{http://schemas.openxmlformats.org/drawingml/2006/main}buBlip')
etree.SubElement(buBlip, '{http://schemas.openxmlformats.org/drawingml/2006/main}blip',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed': rId_img})

# Add second paragraph with same image bullet
p21b_elem = etree.SubElement(tf21._txBody, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
pPr21b = etree.SubElement(p21b_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
buBlip21b = etree.SubElement(pPr21b, '{http://schemas.openxmlformats.org/drawingml/2006/main}buBlip')
etree.SubElement(buBlip21b, '{http://schemas.openxmlformats.org/drawingml/2006/main}blip',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed': rId_img})
r21b = etree.SubElement(p21b_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
rPr21b = etree.SubElement(r21b, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr', attrib={'lang': 'en-US', 'sz': '2000'})
t21b = etree.SubElement(r21b, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
t21b.text = "Image bullet paragraph 2"

# ── Slide 22: Hover link + link color ──────────────────────────────────────────
slide22 = prs.slides.add_slide(blank)
tb22 = slide22.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
tf22 = tb22.text_frame
tf22.word_wrap = True

# Paragraph 1: hlinkClick (normal hyperlink)
p22a = tf22.paragraphs[0]
p22a.text = ""
r22a = p22a.add_run()
r22a.text = "Click link"
r22a.font.size = Pt(24)
# Add hlinkClick via XML
rPr22a = r22a._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
if rPr22a is None:
    rPr22a = etree.SubElement(r22a._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    r22a._r.insert(0, rPr22a)
# Add relationship for click link
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
rId_click = slide22.part.relate_to('https://example.com/click', RT.HYPERLINK, is_external=True)
etree.SubElement(rPr22a, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rId_click})

# Paragraph 2: hlinkHover (mouse-over hyperlink)
p22b = tf22.add_paragraph()
r22b = p22b.add_run()
r22b.text = "Hover link"
r22b.font.size = Pt(24)
rPr22b = r22b._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
if rPr22b is None:
    rPr22b = etree.SubElement(r22b._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    r22b._r.insert(0, rPr22b)
rId_hover = slide22.part.relate_to('https://example.com/hover', RT.HYPERLINK, is_external=True)
etree.SubElement(rPr22b, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkHover',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rId_hover})

# Paragraph 3: both hlinkClick + hlinkHover
p22c = tf22.add_paragraph()
r22c = p22c.add_run()
r22c.text = "Both links"
r22c.font.size = Pt(24)
rPr22c = r22c._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
if rPr22c is None:
    rPr22c = etree.SubElement(r22c._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
    r22c._r.insert(0, rPr22c)
rId_click2 = slide22.part.relate_to('https://example.com/both-click', RT.HYPERLINK, is_external=True)
rId_hover2 = slide22.part.relate_to('https://example.com/both-hover', RT.HYPERLINK, is_external=True)
etree.SubElement(rPr22c, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rId_click2})
etree.SubElement(rPr22c, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkHover',
                 attrib={'{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id': rId_hover2})

# ── Slide 23: Linear gradient fills ──────────────────────────────────────────

slide23 = prs.slides.add_slide(blank)

# Helper to inject gradFill XML into a shape's spPr (after prstGeom, replacing solidFill)
def set_gradient_fill(shape, grad_xml_str):
    """Replace any existing fill with a:gradFill in the shape's spPr."""
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    # spPr can be under a: or p: namespace depending on context
    spPr = shape._element.find(f'{ns_p}spPr')
    if spPr is None:
        spPr = shape._element.find(f'.//{ns_a}spPr')
    if spPr is None:
        print(f"WARNING: spPr not found on shape")
        return
    # Remove existing fills (always a: namespace children)
    for tag in ('solidFill', 'noFill', 'gradFill'):
        for el in spPr.findall(f'{ns_a}{tag}'):
            spPr.remove(el)
    # Insert after prstGeom (a: namespace child)
    prstGeom = spPr.find(f'{ns_a}prstGeom')
    grad_el = etree.fromstring(grad_xml_str)
    if prstGeom is not None:
        prstGeom.addnext(grad_el)
    else:
        spPr.append(grad_el)

# Shape 1: 3-stop linear gradient (left→right, ang=0)
s23a = slide23.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(3), Inches(2))
s23a.text = "Linear 0deg"
s23a.text_frame.paragraphs[0].font.size = Pt(18)
s23a.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s23a, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FF0000"/></a:gs>
    <a:gs pos="50000"><a:srgbClr val="FFFF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="1"/>
</a:gradFill>''')

# Shape 2: 2-stop linear gradient (top→bottom, ang=5400000 = 90°)
s23b = slide23.shapes.add_shape(1, Inches(4), Inches(0.5), Inches(3), Inches(2))
s23b.text = "Linear 90deg"
s23b.text_frame.paragraphs[0].font.size = Pt(18)
s23b.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s23b, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="003366"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="66CCFF"/></a:gs>
  </a:gsLst>
  <a:lin ang="5400000" scaled="1"/>
</a:gradFill>''')

# Shape 3: Diagonal linear gradient (45° = 2700000), rotWithShape="0"
s23c = slide23.shapes.add_shape(1, Inches(7.5), Inches(0.5), Inches(2), Inches(2))
s23c.text = "Linear 45deg"
s23c.text_frame.paragraphs[0].font.size = Pt(14)
set_gradient_fill(s23c, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="0">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="00FF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="006600"/></a:gs>
  </a:gsLst>
  <a:lin ang="2700000" scaled="1"/>
</a:gradFill>''')


# ── Slide 24: Radial/path gradient fills + gradient on ellipse ───────────────

slide24 = prs.slides.add_slide(blank)

# Shape 1: Radial circle gradient
s24a = slide24.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(3), Inches(2.5))
s24a.text = "Radial circle"
s24a.text_frame.paragraphs[0].font.size = Pt(18)
s24a.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s24a, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="FF6600"/></a:gs>
  </a:gsLst>
  <a:path path="circle">
    <a:fillToRect l="50000" t="50000" r="50000" b="50000"/>
  </a:path>
</a:gradFill>''')

# Shape 2: Radial rect gradient with off-center fillToRect
s24b = slide24.shapes.add_shape(1, Inches(4), Inches(0.5), Inches(3), Inches(2.5))
s24b.text = "Radial rect"
s24b.text_frame.paragraphs[0].font.size = Pt(18)
set_gradient_fill(s24b, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FFFF00"/></a:gs>
    <a:gs pos="50000"><a:srgbClr val="FF0000"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="660000"/></a:gs>
  </a:gsLst>
  <a:path path="rect">
    <a:fillToRect l="25000" t="25000" r="75000" b="75000"/>
  </a:path>
</a:gradFill>''')

# Shape 3: Gradient on ellipse (oval)
from pptx.enum.shapes import MSO_SHAPE
s24c = slide24.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(0.5), Inches(2), Inches(2.5))
s24c.text = "Ellipse grad"
s24c.text_frame.paragraphs[0].font.size = Pt(14)
s24c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s24c, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="9933FF"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="330066"/></a:gs>
  </a:gsLst>
  <a:lin ang="16200000" scaled="1"/>
</a:gradFill>''')


# ── Slide 25: Gradient background ────────────────────────────────────────────

slide25 = prs.slides.add_slide(blank)

# Inject gradient background via XML
ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
cSld25 = slide25._element.find(f'{ns_p}cSld')
spTree25 = cSld25.find(f'{ns_p}spTree')
bg_xml = f'''<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:bgPr>
    <a:gradFill>
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="1B2838"/></a:gs>
        <a:gs pos="50000"><a:srgbClr val="2A475E"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="66C0F4"/></a:gs>
      </a:gsLst>
      <a:lin ang="5400000" scaled="0"/>
    </a:gradFill>
    <a:effectLst/>
  </p:bgPr>
</p:bg>'''
bg_el = etree.fromstring(bg_xml)
cSld25.insert(list(cSld25).index(spTree25), bg_el)

# Text shape on gradient background (noFill)
s25a = slide25.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
s25a.text = "Gradient Background"
s25a.text_frame.paragraphs[0].font.size = Pt(36)
s25a.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
s25a.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ── Slide 26: Alpha/transparency ─────────────────────────────────────────────

slide26 = prs.slides.add_slide(blank)

def set_fill_xml(shape, fill_xml_str):
    """Replace any existing fill with custom fill XML in the shape's spPr."""
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    spPr = shape._element.find(f'{ns_p}spPr')
    if spPr is None:
        spPr = shape._element.find(f'.//{ns_a}spPr')
    if spPr is None:
        return
    for tag in ('solidFill', 'noFill', 'gradFill', 'blipFill', 'pattFill'):
        for el in spPr.findall(f'{ns_a}{tag}'):
            spPr.remove(el)
    prstGeom = spPr.find(f'{ns_a}prstGeom')
    fill_el = etree.fromstring(fill_xml_str)
    if prstGeom is not None:
        prstGeom.addnext(fill_el)
    else:
        spPr.append(fill_el)

# Shape 1: Semi-transparent red solid fill (alpha 50%)
s26a = slide26.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(3), Inches(2))
s26a.text = "Alpha 50%"
s26a.text_frame.paragraphs[0].font.size = Pt(18)
set_fill_xml(s26a, '''<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:srgbClr val="FF0000"><a:alpha val="50000"/></a:srgbClr>
</a:solidFill>''')

# Shape 2: Semi-transparent gradient stop
s26b = slide26.shapes.add_shape(1, Inches(4), Inches(0.5), Inches(3), Inches(2))
s26b.text = "Alpha Gradient"
s26b.text_frame.paragraphs[0].font.size = Pt(18)
s26b.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
set_gradient_fill(s26b, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="0000FF"><a:alpha val="80000"/></a:srgbClr></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"><a:alpha val="20000"/></a:srgbClr></a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="1"/>
</a:gradFill>''')

# ── Slide 27: Image fill on AutoShape ────────────────────────────────────────

slide27 = prs.slides.add_slide(blank)

# We need to add an image relationship and then inject a:blipFill into p:spPr
# First, add a simple shape
s27a = slide27.shapes.add_shape(1, Inches(1), Inches(1), Inches(4), Inches(3))
s27a.text = ""

# Add image relationship (use slide's existing rels) — use a small 1px PNG
# Minimal 1x1 red PNG
mini_png = base64.b64decode(
    'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8/5+hHgAHggJ/PchI7wAAAABJRU5ErkJggg=='
)
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Write temp image
tmp_img = os.path.join(tempfile.gettempdir(), '_test_blip.png')
with open(tmp_img, 'wb') as f:
    f.write(mini_png)

# Add image as a relationship on the slide's part
from pptx.opc.package import Part
img_part, img_rid = slide27.part.get_or_add_image_part(tmp_img)

# Inject a:blipFill into p:spPr
set_fill_xml(s27a, f'''<a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <a:blip r:embed="{img_rid}"/>
  <a:stretch><a:fillRect/></a:stretch>
</a:blipFill>''')

os.unlink(tmp_img)

# ── Slide 28: Pattern fill ───────────────────────────────────────────────────

slide28 = prs.slides.add_slide(blank)

# Shape 1: ltDnDiag pattern
s28a = slide28.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(2.5), Inches(2))
s28a.text = "ltDnDiag"
s28a.text_frame.paragraphs[0].font.size = Pt(14)
set_fill_xml(s28a, '''<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="ltDnDiag">
  <a:fgClr><a:srgbClr val="000000"/></a:fgClr>
  <a:bgClr><a:srgbClr val="FFFFFF"/></a:bgClr>
</a:pattFill>''')

# Shape 2: smCheck pattern
s28b = slide28.shapes.add_shape(1, Inches(3.5), Inches(0.5), Inches(2.5), Inches(2))
s28b.text = "smCheck"
s28b.text_frame.paragraphs[0].font.size = Pt(14)
set_fill_xml(s28b, '''<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="smCheck">
  <a:fgClr><a:srgbClr val="FF0000"/></a:fgClr>
  <a:bgClr><a:srgbClr val="FFFF00"/></a:bgClr>
</a:pattFill>''')

# Shape 3: dkHorz pattern
s28c = slide28.shapes.add_shape(1, Inches(6.5), Inches(0.5), Inches(2.5), Inches(2))
s28c.text = "dkHorz"
s28c.text_frame.paragraphs[0].font.size = Pt(14)
set_fill_xml(s28c, '''<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="dkHorz">
  <a:fgClr><a:srgbClr val="003366"/></a:fgClr>
  <a:bgClr><a:srgbClr val="CCCCCC"/></a:bgClr>
</a:pattFill>''')

# ── Slide 29: Gradient tileFlip ──────────────────────────────────────────────

slide29 = prs.slides.add_slide(blank)

# Shape 1: tileFlip="x"
s29a = slide29.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(2.5), Inches(2))
s29a.text = "tileFlip=x"
s29a.text_frame.paragraphs[0].font.size = Pt(14)
set_gradient_fill(s29a, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1" tileFlip="x">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FF0000"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="1"/>
</a:gradFill>''')

# Shape 2: tileFlip="y"
s29b = slide29.shapes.add_shape(1, Inches(3.5), Inches(0.5), Inches(2.5), Inches(2))
s29b.text = "tileFlip=y"
s29b.text_frame.paragraphs[0].font.size = Pt(14)
set_gradient_fill(s29b, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1" tileFlip="y">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="00FF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="FF00FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="5400000" scaled="1"/>
</a:gradFill>''')

# Shape 3: tileFlip="xy"
s29c = slide29.shapes.add_shape(1, Inches(6.5), Inches(0.5), Inches(2.5), Inches(2))
s29c.text = "tileFlip=xy"
s29c.text_frame.paragraphs[0].font.size = Pt(14)
set_gradient_fill(s29c, '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1" tileFlip="xy">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FFFF00"/></a:gs>
    <a:gs pos="50000"><a:srgbClr val="FF6600"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0066FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="2700000" scaled="1"/>
</a:gradFill>''')

# ── Slide 30: Additional pattern fills ──────────────────────────────────────

slide30 = prs.slides.add_slide(blank)

patterns_30 = [
    ("pct50", "333333", "CCCCCC"),
    ("dnDiag", "000080", "FFFFFF"),
    ("cross", "FF0000", "FFFFCC"),
    ("lgCheck", "008000", "FFFFFF"),
    ("solidDmnd", "800080", "FFE0FF"),
    ("trellis", "004040", "E0FFFF"),
]

for idx, (prst, fg_c, bg_c) in enumerate(patterns_30):
    col = idx % 3
    row = idx // 3
    left = Inches(0.5 + col * 3.0)
    top = Inches(0.5 + row * 2.5)
    sh = slide30.shapes.add_shape(1, left, top, Inches(2.5), Inches(2))
    sh.text = prst
    sh.text_frame.paragraphs[0].font.size = Pt(14)
    set_fill_xml(sh, f'''<a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="{prst}">
  <a:fgClr><a:srgbClr val="{fg_c}"/></a:fgClr>
  <a:bgClr><a:srgbClr val="{bg_c}"/></a:bgClr>
</a:pattFill>''')

# ── Slide 31: Image fill tile ───────────────────────────────────────────────

slide31 = prs.slides.add_slide(blank)

# Re-create temp image for tile test
tmp_img2 = os.path.join(tempfile.gettempdir(), '_test_blip_tile.png')
with open(tmp_img2, 'wb') as f:
    f.write(mini_png)

img_part2, img_rid2 = slide31.part.get_or_add_image_part(tmp_img2)

# Shape with tile fill
s31a = slide31.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(4), Inches(3))
s31a.text = "Tile fill"
s31a.text_frame.paragraphs[0].font.size = Pt(14)
set_fill_xml(s31a, f'''<a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <a:blip r:embed="{img_rid2}"/>
  <a:tile tx="0" ty="0" sx="50000" sy="50000" flip="xy" algn="tl"/>
</a:blipFill>''')

os.unlink(tmp_img2)

# ── Helper: set line XML on a shape ────────────────────────────────────────

def set_line_xml(shape, ln_xml_str):
    """Replace any existing a:ln with custom line XML in the shape's spPr."""
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    spPr = shape._element.find(f'{ns_p}spPr')
    if spPr is None:
        spPr = shape._element.find(f'.//{ns_a}spPr')
    if spPr is None:
        return
    for el in spPr.findall(f'{ns_a}ln'):
        spPr.remove(el)
    ln_el = etree.fromstring(ln_xml_str)
    spPr.append(ln_el)

# ── Helper: create a line shape (p:sp with a:prstGeom prst="line") ─────────

def add_line_shape(slide, left, top, width, height, ln_xml_str):
    """Add a p:sp with line geometry + custom a:ln XML."""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    # Create via add_shape(1=rect) then change geometry to line
    s = slide.shapes.add_shape(1, left, top, width, height)
    spPr = s._element.find(f'{{{ns_p}}}spPr')
    if spPr is None:
        spPr = s._element.find(f'.//{{{ns_a}}}spPr')
    # Replace prstGeom with line
    for pg in spPr.findall(f'{{{ns_a}}}prstGeom'):
        spPr.remove(pg)
    pg = etree.SubElement(spPr, f'{{{ns_a}}}prstGeom')
    pg.set('prst', 'line')
    etree.SubElement(pg, f'{{{ns_a}}}avLst')
    # Remove fill (lines don't have fill)
    for tag in ('solidFill', 'noFill', 'gradFill'):
        for el in spPr.findall(f'{{{ns_a}}}{tag}'):
            spPr.remove(el)
    noFill = etree.SubElement(spPr, f'{{{ns_a}}}noFill')
    # Set line XML
    for el in spPr.findall(f'{{{ns_a}}}ln'):
        spPr.remove(el)
    ln_el = etree.fromstring(ln_xml_str)
    spPr.append(ln_el)
    return s

# ── Slide 32: Stroke styles on lines ─────────────────────────────────────

slide32 = prs.slides.add_slide(blank)

title32 = slide32.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
title32.text_frame.paragraphs[0].text = "Slide 32: Line Stroke Styles"
title32.text_frame.paragraphs[0].font.size = Pt(20)
title32.text_frame.paragraphs[0].font.bold = True

# Dash styles on lines
dash_info = [
    ('dash', '0000FF'),
    ('dot', '0066CC'),
    ('dashDot', '009900'),
    ('lgDash', 'CC6600'),
    ('sysDot', '990099'),
    ('sysDash', 'CC0000'),
]
for i, (dash, color) in enumerate(dash_info):
    y = 1.0 + i * 0.9
    add_line_shape(slide32, Inches(1), Inches(y), Inches(8), Emu(0),
        f'''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
          <a:prstDash val="{dash}"/>
        </a:ln>''')
    # Label
    lbl = slide32.shapes.add_textbox(Inches(0.2), Inches(y - 0.25), Inches(0.8), Inches(0.3))
    lbl.text_frame.paragraphs[0].text = dash
    lbl.text_frame.paragraphs[0].font.size = Pt(9)
    lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Custom dash (a:custDash with a:ds elements)
add_line_shape(slide32, Inches(1), Inches(1.0 + 6 * 0.9), Inches(8), Emu(0),
    '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="FF6600"/></a:solidFill>
      <a:custDash>
        <a:ds d="400000" sp="100000"/>
        <a:ds d="100000" sp="100000"/>
      </a:custDash>
    </a:ln>''')
lbl_cd = slide32.shapes.add_textbox(Inches(0.2), Inches(1.0 + 6 * 0.9 - 0.25), Inches(0.8), Inches(0.3))
lbl_cd.text_frame.paragraphs[0].text = "custDash"
lbl_cd.text_frame.paragraphs[0].font.size = Pt(9)
lbl_cd.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide 33: Arrows + Line cap/join ──────────────────────────────────────

slide33 = prs.slides.add_slide(blank)

title33 = slide33.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
title33.text_frame.paragraphs[0].text = "Slide 33: Arrow Lines"
title33.text_frame.paragraphs[0].font.size = Pt(20)
title33.text_frame.paragraphs[0].font.bold = True

# Line with triangle head + stealth tail
add_line_shape(slide33, Inches(0.5), Inches(1.2), Inches(4), Emu(0),
    '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>
      <a:headEnd type="triangle" w="med" len="med"/>
      <a:tailEnd type="stealth" w="lg" len="lg"/>
    </a:ln>''')
lbl33a = slide33.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4), Inches(0.3))
lbl33a.text_frame.paragraphs[0].text = "triangle head + stealth tail"
lbl33a.text_frame.paragraphs[0].font.size = Pt(9)
lbl33a.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Line with diamond head + oval tail
add_line_shape(slide33, Inches(0.5), Inches(2.0), Inches(4), Emu(0),
    '''<a:ln w="19050" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="008000"/></a:solidFill>
      <a:headEnd type="diamond" w="med" len="med"/>
      <a:tailEnd type="oval" w="med" len="med"/>
    </a:ln>''')
lbl33b = slide33.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4), Inches(0.3))
lbl33b.text_frame.paragraphs[0].text = "diamond head + oval tail"
lbl33b.text_frame.paragraphs[0].font.size = Pt(9)
lbl33b.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Line with arrow (open) head + triangle tail
add_line_shape(slide33, Inches(0.5), Inches(2.8), Inches(4), Emu(0),
    '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="0000FF"/></a:solidFill>
      <a:headEnd type="arrow" w="med" len="med"/>
      <a:tailEnd type="triangle" w="sm" len="sm"/>
    </a:ln>''')
lbl33c = slide33.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(4), Inches(0.3))
lbl33c.text_frame.paragraphs[0].text = "arrow head + triangle tail (sm)"
lbl33c.text_frame.paragraphs[0].font.size = Pt(9)
lbl33c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Dashed line with round cap + round join
add_line_shape(slide33, Inches(5), Inches(1.2), Inches(4.5), Emu(0),
    '''<a:ln w="38100" cap="rnd" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="800080"/></a:solidFill>
      <a:prstDash val="dash"/>
      <a:round/>
    </a:ln>''')
lbl33d = slide33.shapes.add_textbox(Inches(5), Inches(1.3), Inches(4.5), Inches(0.3))
lbl33d.text_frame.paragraphs[0].text = "dash + round cap + round join"
lbl33d.text_frame.paragraphs[0].font.size = Pt(9)
lbl33d.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Line with square cap + bevel join + lgDash
add_line_shape(slide33, Inches(5), Inches(2.0), Inches(4.5), Emu(0),
    '''<a:ln w="25400" cap="sq" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="006699"/></a:solidFill>
      <a:prstDash val="lgDash"/>
      <a:bevel/>
    </a:ln>''')
lbl33e = slide33.shapes.add_textbox(Inches(5), Inches(2.1), Inches(4.5), Inches(0.3))
lbl33e.text_frame.paragraphs[0].text = "lgDash + sq cap + bevel join"
lbl33e.text_frame.paragraphs[0].font.size = Pt(9)
lbl33e.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Miter join with limit
add_line_shape(slide33, Inches(5), Inches(2.8), Inches(4.5), Emu(0),
    '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:solidFill><a:srgbClr val="FF6600"/></a:solidFill>
      <a:miter lim="800000"/>
    </a:ln>''')
lbl33f = slide33.shapes.add_textbox(Inches(5), Inches(2.9), Inches(4.5), Inches(0.3))
lbl33f.text_frame.paragraphs[0].text = "miter join (lim=800000)"
lbl33f.text_frame.paragraphs[0].font.size = Pt(9)
lbl33f.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Dashed rect (stroke on rect shape) — keeps one rect for dash+cap test
s33rect = slide33.shapes.add_shape(1, Inches(0.5), Inches(4.0), Inches(4), Inches(2))
s33rect.text = "Rect: dashDot + rnd cap"
s33rect.text_frame.paragraphs[0].font.size = Pt(12)
set_line_xml(s33rect, '''<a:ln w="25400" cap="rnd" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:solidFill><a:srgbClr val="FF00FF"/></a:solidFill>
  <a:prstDash val="dashDot"/>
  <a:round/>
</a:ln>''')

# Compound line: double
s33cmpd = slide33.shapes.add_shape(1, Inches(5), Inches(4.0), Inches(4.5), Inches(0.8))
s33cmpd.text = "Compound: dbl"
s33cmpd.text_frame.paragraphs[0].font.size = Pt(12)
set_line_xml(s33cmpd, '''<a:ln w="38100" cmpd="dbl" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
</a:ln>''')

# noFill line (stroke_no_fill)
s33nf = slide33.shapes.add_shape(1, Inches(5), Inches(5.2), Inches(4.5), Inches(0.8))
s33nf.text = "Line noFill"
s33nf.text_frame.paragraphs[0].font.size = Pt(12)
set_line_xml(s33nf, '''<a:ln w="25400" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:noFill/>
</a:ln>''')

# ── Slide 34: Group shapes ────────────────────────────────────────────────
slide34 = prs.slides.add_slide(blank)

title34 = slide34.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title34.text_frame.paragraphs[0].text = "Slide 34: Group Shapes"
title34.text_frame.paragraphs[0].font.size = Pt(24)
title34.text_frame.paragraphs[0].font.bold = True

# Create a group shape with two rectangles using raw XML
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

spTree34 = slide34.shapes._spTree
grp_xml = f'''<p:grpSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvGrpSpPr>
    <p:cNvPr id="100" name="Group 1"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm>
      <a:off x="457200" y="1371600"/>
      <a:ext cx="3657600" cy="2743200"/>
      <a:chOff x="0" y="0"/>
      <a:chExt cx="3657600" cy="2743200"/>
    </a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="101" name="Rect1"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="0" y="0"/>
        <a:ext cx="1828800" cy="1371600"/>
      </a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="FF6B6B"/></a:solidFill>
      <a:ln w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:lstStyle/>
      <a:p><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Child 1 (red)</a:t></a:r></a:p>
    </p:txBody>
  </p:sp>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="102" name="Rect2"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="1828800" y="1371600"/>
        <a:ext cx="1828800" cy="1371600"/>
      </a:xfrm>
      <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="4ECDC4"/></a:solidFill>
      <a:ln w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:lstStyle/>
      <a:p><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Child 2 (teal)</a:t></a:r></a:p>
    </p:txBody>
  </p:sp>
</p:grpSp>'''
spTree34.append(etree.fromstring(grp_xml))

# Second group: nested group (group inside group)
grp_nested_xml = f'''<p:grpSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvGrpSpPr>
    <p:cNvPr id="200" name="Outer Group"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm>
      <a:off x="4572000" y="1371600"/>
      <a:ext cx="4572000" cy="2743200"/>
      <a:chOff x="0" y="0"/>
      <a:chExt cx="4572000" cy="2743200"/>
    </a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="201" name="Outer Rect"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="0" y="0"/>
        <a:ext cx="4572000" cy="2743200"/>
      </a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="E8E8E8"/></a:solidFill>
      <a:ln w="12700"><a:solidFill><a:srgbClr val="999999"/></a:solidFill></a:ln>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:lstStyle/>
      <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Outer bg</a:t></a:r></a:p>
    </p:txBody>
  </p:sp>
  <p:grpSp>
    <p:nvGrpSpPr>
      <p:cNvPr id="210" name="Inner Group"/>
      <p:cNvGrpSpPr/>
      <p:nvPr/>
    </p:nvGrpSpPr>
    <p:grpSpPr>
      <a:xfrm>
        <a:off x="457200" y="457200"/>
        <a:ext cx="3657600" cy="1828800"/>
        <a:chOff x="0" y="0"/>
        <a:chExt cx="3657600" cy="1828800"/>
      </a:xfrm>
    </p:grpSpPr>
    <p:sp>
      <p:nvSpPr>
        <p:cNvPr id="211" name="Inner Circle"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="1828800" cy="1828800"/>
        </a:xfrm>
        <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="FFD93D"/></a:solidFill>
      </p:spPr>
      <p:txBody>
        <a:bodyPr anchor="ctr"/>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Nested A</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>
    <p:sp>
      <p:nvSpPr>
        <p:cNvPr id="212" name="Inner Rect"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="1828800" y="0"/>
          <a:ext cx="1828800" cy="1828800"/>
        </a:xfrm>
        <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="6C5CE7"/></a:solidFill>
      </p:spPr>
      <p:txBody>
        <a:bodyPr anchor="ctr"/>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>Nested B</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>
  </p:grpSp>
</p:grpSp>'''
spTree34.append(etree.fromstring(grp_nested_xml))

# Label
lbl34 = slide34.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
lbl34.text_frame.paragraphs[0].text = "Left: simple group (red rect + teal ellipse). Right: nested group (outer bg + inner circle + rounded rect)."
lbl34.text_frame.paragraphs[0].font.size = Pt(12)
lbl34.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 35: Connectors ────────────────────────────────────────────────────
slide35 = prs.slides.add_slide(blank)

title35 = slide35.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title35.text_frame.paragraphs[0].text = "Slide 35: Connectors (p:cxnSp)"
title35.text_frame.paragraphs[0].font.size = Pt(24)
title35.text_frame.paragraphs[0].font.bold = True

# Straight connector
spTree35 = slide35.shapes._spTree
cxn1_xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvCxnSpPr>
    <p:cNvPr id="300" name="Straight Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="1371600"/>
      <a:ext cx="3657600" cy="0"/>
    </a:xfrm>
    <a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>
    <a:ln w="25400">
      <a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>
      <a:tailEnd type="triangle"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree35.append(etree.fromstring(cxn1_xml))

# Diagonal connector with arrows on both ends
cxn2_xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvCxnSpPr>
    <p:cNvPr id="301" name="Diagonal Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="2286000"/>
      <a:ext cx="3657600" cy="1828800"/>
    </a:xfrm>
    <a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>
    <a:ln w="19050">
      <a:solidFill><a:srgbClr val="0066CC"/></a:solidFill>
      <a:prstDash val="dash"/>
      <a:headEnd type="diamond"/>
      <a:tailEnd type="stealth"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree35.append(etree.fromstring(cxn2_xml))

# Bent connector (L-shape)
cxn3_xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvCxnSpPr>
    <p:cNvPr id="302" name="Bent Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="5029200" y="1371600"/>
      <a:ext cx="3657600" cy="1371600"/>
    </a:xfrm>
    <a:prstGeom prst="bentConnector3">
      <a:avLst>
        <a:gd name="adj1" fmla="val 50000"/>
      </a:avLst>
    </a:prstGeom>
    <a:ln w="25400">
      <a:solidFill><a:srgbClr val="009900"/></a:solidFill>
      <a:tailEnd type="triangle"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree35.append(etree.fromstring(cxn3_xml))

# Curved connector
cxn4_xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
  <p:nvCxnSpPr>
    <p:cNvPr id="303" name="Curved Connector"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="5029200" y="3200400"/>
      <a:ext cx="3657600" cy="1371600"/>
    </a:xfrm>
    <a:prstGeom prst="curvedConnector3">
      <a:avLst>
        <a:gd name="adj1" fmla="val 50000"/>
      </a:avLst>
    </a:prstGeom>
    <a:ln w="25400">
      <a:solidFill><a:srgbClr val="FF6600"/></a:solidFill>
      <a:headEnd type="oval"/>
      <a:tailEnd type="arrow"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree35.append(etree.fromstring(cxn4_xml))

lbl35 = slide35.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
lbl35.text_frame.paragraphs[0].text = "Left: straight + diagonal. Right: bent (green) + curved (orange)."
lbl35.text_frame.paragraphs[0].font.size = Pt(12)
lbl35.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ────────────────────────────────────────────────────────────────────────────
# Slide 36: Preset geometry shapes (triangle, diamond, arrow, star, heart, etc.)
# ────────────────────────────────────────────────────────────────────────────
slide36 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
spTree36 = slide36.shapes._spTree

# Helper: inject a preset geometry shape via raw XML
def make_prst_shape(name, prst, x_emu, y_emu, cx_emu, cy_emu, fill_hex, avlst_xml=''):
    av_xml = f'<a:avLst>{avlst_xml}</a:avLst>' if avlst_xml else '<a:avLst/>'
    return f'''<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="0" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{x_emu}" y="{y_emu}"/>
      <a:ext cx="{cx_emu}" cy="{cy_emu}"/>
    </a:xfrm>
    <a:prstGeom prst="{prst}">{av_xml}</a:prstGeom>
    <a:solidFill><a:srgbClr val="{fill_hex}"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr lang="en-US" sz="900"/><a:t>{prst}</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''

shapes_36 = [
    ("triangle", "FF6B6B"),
    ("diamond", "4ECDC4"),
    ("pentagon", "FFD93D"),
    ("hexagon", "6C5CE7"),
    ("rightArrow", "74B9FF"),
    ("star5", "FDCB6E"),
    ("heart", "E17055"),
    ("plus", "00CEC9"),
    ("flowChartDecision", "A29BFE"),
    ("chevron", "55EFC4"),
    ("parallelogram", "FF7675"),
    ("octagon", "FFEAA7"),
]

row_y = [Emu(457200), Emu(2286000)]  # 2 rows
col_x = [Emu(457200 + i * 1524000) for i in range(6)]  # 6 columns
sz = Emu(1143000)  # ~1.2 inch
for idx, (prst, fill) in enumerate(shapes_36):
    row = idx // 6
    col = idx % 6
    xml = make_prst_shape(f"prst_{prst}", prst, col_x[col], row_y[row], sz, sz, fill)
    spTree36.append(etree.fromstring(xml))

# Also add shapes with custom adj values
adj_shapes = [
    ("rightArrow_adj", "rightArrow", Emu(457200), Emu(4114800), Emu(2286000), Emu(914400), "74B9FF",
     '<a:gd name="adj1" fmla="val 70000"/><a:gd name="adj2" fmla="val 30000"/>'),
    ("star5_adj", "star5", Emu(3200400), Emu(4114800), Emu(1143000), Emu(1143000), "FDCB6E",
     '<a:gd name="adj" fmla="val 40000"/>'),
]
for name, prst, x, y, cx, cy, fill, avlst in adj_shapes:
    xml = make_prst_shape(name, prst, x, y, cx, cy, fill, avlst)
    spTree36.append(etree.fromstring(xml))

lbl36 = slide36.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl36.text_frame.paragraphs[0].text = "Preset geometry: triangle, diamond, pentagon, hexagon, rightArrow, star5, heart, plus, flowChartDecision, chevron, parallelogram, octagon"
lbl36.text_frame.paragraphs[0].font.size = Pt(10)
lbl36.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 37: Custom geometry (a:custGeom) ──
slide37 = prs.slides.add_slide(blank)
spTree37 = slide37.shapes._spTree

# Custom freeform shape 1: a star-like custom geometry with path coordinates
cust1_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="100" name="CustomStar"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="457200"/><a:ext cx="1828800" cy="1828800"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:pathLst>
        <a:path w="200" h="200">
          <a:moveTo><a:pt x="100" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="130" y="70"/></a:lnTo>
          <a:lnTo><a:pt x="200" y="80"/></a:lnTo>
          <a:lnTo><a:pt x="150" y="130"/></a:lnTo>
          <a:lnTo><a:pt x="160" y="200"/></a:lnTo>
          <a:lnTo><a:pt x="100" y="170"/></a:lnTo>
          <a:lnTo><a:pt x="40" y="200"/></a:lnTo>
          <a:lnTo><a:pt x="50" y="130"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="80"/></a:lnTo>
          <a:lnTo><a:pt x="70" y="70"/></a:lnTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="FFD700"/></a:solidFill>
    <a:ln w="19050"><a:solidFill><a:srgbClr val="B8860B"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree37.append(etree.fromstring(cust1_xml))

# Custom freeform shape 2: curved shape with cubicBezTo
cust2_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="101" name="CustomCurve"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="2743200" y="457200"/><a:ext cx="1828800" cy="1828800"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:pathLst>
        <a:path w="100" h="100">
          <a:moveTo><a:pt x="0" y="100"/></a:moveTo>
          <a:cubicBezTo>
            <a:pt x="0" y="0"/>
            <a:pt x="100" y="0"/>
            <a:pt x="100" y="100"/>
          </a:cubicBezTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="87CEEB"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="4169E1"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree37.append(etree.fromstring(cust2_xml))

# Custom shape 3: with guide formulas
cust3_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="102" name="CustomGuide"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="457200"/><a:ext cx="1828800" cy="1828800"/></a:xfrm>
    <a:custGeom>
      <a:avLst>
        <a:gd name="adj" fmla="val 50000"/>
      </a:avLst>
      <a:gdLst>
        <a:gd name="x1" fmla="*/ w adj 100000"/>
        <a:gd name="y1" fmla="*/ h adj 100000"/>
      </a:gdLst>
      <a:pathLst>
        <a:path>
          <a:moveTo><a:pt x="0" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="x1" y="0"/></a:lnTo>
          <a:lnTo><a:pt x="x1" y="y1"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="y1"/></a:lnTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="98FB98"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="228B22"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree37.append(etree.fromstring(cust3_xml))

lbl37 = slide37.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl37.text_frame.paragraphs[0].text = "Custom geometry (a:custGeom): freeform star, bezier curve, guide-based rect"
lbl37.text_frame.paragraphs[0].font.size = Pt(10)
lbl37.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 38: Gear shapes ────────────────────────────────────────────────────
slide38 = prs.slides.add_slide(blank)
spTree38 = slide38.shapes._spTree

# Gear6
gear6_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="201" name="Gear6"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="gear6"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="2F528F"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree38.append(etree.fromstring(gear6_xml))

# Gear9
gear9_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="202" name="Gear9"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="3429000" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="gear9"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="ED7D31"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="AE5A21"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree38.append(etree.fromstring(gear9_xml))

# Gear6 with custom adj (deeper teeth)
gear6_adj_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="203" name="Gear6Deep"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="6400800" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="gear6"><a:avLst><a:gd name="adj" fmla="val 50000"/></a:avLst></a:prstGeom>
    <a:solidFill><a:srgbClr val="70AD47"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="507E32"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree38.append(etree.fromstring(gear6_adj_xml))

lbl38 = slide38.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl38.text_frame.paragraphs[0].text = "Gear shapes: gear6 (default), gear9 (default), gear6 (adj=50000 deeper teeth)"
lbl38.text_frame.paragraphs[0].font.size = Pt(10)
lbl38.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 39: Text rectangles ────────────────────────────────────────────────
slide39 = prs.slides.add_slide(blank)
spTree39 = slide39.shapes._spTree

# Triangle with text
tri_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="301" name="Triangle"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="triangle"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFD700"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="B8860B"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Text in Triangle</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree39.append(etree.fromstring(tri_xml))

# Diamond with text
diamond_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="302" name="Diamond"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="3429000" y="457200"/><a:ext cx="2286000" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="diamond"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="87CEEB"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="4682B4"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Diamond</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree39.append(etree.fromstring(diamond_xml))

# Right arrow with text
arrow_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="303" name="RightArrow"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="6400800" y="914400"/><a:ext cx="2743200" cy="1371600"/></a:xfrm>
    <a:prstGeom prst="rightArrow"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FF6347"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="8B0000"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Arrow</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree39.append(etree.fromstring(arrow_xml))

lbl39 = slide39.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl39.text_frame.paragraphs[0].text = "Text rectangles: triangle (text in lower half), diamond (text inscribed), right arrow (text in shaft)"
lbl39.text_frame.paragraphs[0].font.size = Pt(10)
lbl39.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 40: Connection points ──────────────────────────────────────────────
slide40 = prs.slides.add_slide(blank)
spTree40 = slide40.shapes._spTree

# Shape 1: Rectangle (target for connectors)
rect1_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="401" name="Rect1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="914400" y="914400"/><a:ext cx="1828800" cy="1371600"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
  </p:spPr>
  <p:txBody><a:bodyPr anchor="ctr"/><a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>Start</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree40.append(etree.fromstring(rect1_xml))

# Shape 2: Rectangle (target for connectors)
rect2_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="402" name="Rect2"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5486400" y="914400"/><a:ext cx="1828800" cy="1371600"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="ED7D31"/></a:solidFill>
  </p:spPr>
  <p:txBody><a:bodyPr anchor="ctr"/><a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>End</a:t></a:r></a:p>
  </p:txBody>
</p:sp>'''
spTree40.append(etree.fromstring(rect2_xml))

# Connector with stCxn/endCxn references
cxn_xml = '''<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvCxnSpPr>
    <p:cNvPr id="403" name="Connector"/>
    <p:cNvCxnSpPr>
      <a:stCxn id="401" idx="3"/>
      <a:endCxn id="402" idx="1"/>
    </p:cNvCxnSpPr>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm><a:off x="2743200" y="1600200"/><a:ext cx="2743200" cy="0"/></a:xfrm>
    <a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>
    <a:ln w="25400">
      <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
      <a:tailEnd type="triangle"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''
spTree40.append(etree.fromstring(cxn_xml))

# Custom geometry with a:cxnLst
cust_cxn_xml = '''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr><p:cNvPr id="404" name="CustomWithCxn"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="3200400" y="3429000"/><a:ext cx="2286000" cy="1828800"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:cxnLst>
        <a:cxn ang="0"><a:pos x="r" y="vc"/></a:cxn>
        <a:cxn ang="5400000"><a:pos x="hc" y="b"/></a:cxn>
        <a:cxn ang="10800000"><a:pos x="l" y="vc"/></a:cxn>
        <a:cxn ang="16200000"><a:pos x="hc" y="t"/></a:cxn>
      </a:cxnLst>
      <a:rect l="l" t="t" r="r" b="b"/>
      <a:pathLst>
        <a:path w="100" h="100">
          <a:moveTo><a:pt x="50" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="100" y="50"/></a:lnTo>
          <a:lnTo><a:pt x="50" y="100"/></a:lnTo>
          <a:lnTo><a:pt x="0" y="50"/></a:lnTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="9370DB"/></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="4B0082"/></a:solidFill></a:ln>
  </p:spPr>
</p:sp>'''
spTree40.append(etree.fromstring(cust_cxn_xml))

lbl40 = slide40.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl40.text_frame.paragraphs[0].text = "Connection points: connector with stCxn/endCxn refs, custom geom with a:cxnLst"
lbl40.text_frame.paragraphs[0].font.size = Pt(10)
lbl40.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 41: Table cell merge + borders + margins + anchor ──────────────────
slide41 = prs.slides.add_slide(blank)
ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

# Build table XML directly for full control over cell merge/border/margin/anchor
table_xml = f'''<p:graphicFrame xmlns:a="{ns_a[1:-1]}" xmlns:p="{ns_p[1:-1]}" xmlns:r="{ns_r[1:-1]}">
  <p:nvGraphicFramePr>
    <p:cNvPr id="100" name="Table41"/>
    <p:cNvGraphicFramePr><a:graphicFrameLocks noGrp="1"/></p:cNvGraphicFramePr>
    <p:nvPr/>
  </p:nvGraphicFramePr>
  <p:xfrm>
    <a:off x="457200" y="457200"/>
    <a:ext cx="8229600" cy="4572000"/>
  </p:xfrm>
  <a:graphic>
    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
      <a:tbl>
        <a:tblPr firstRow="1" bandRow="1"/>
        <a:tblGrid>
          <a:gridCol w="2743200"/>
          <a:gridCol w="2743200"/>
          <a:gridCol w="2743200"/>
        </a:tblGrid>
        <!-- Row 1: horizontal merge (gridSpan=2) + custom borders -->
        <a:tr h="914400">
          <a:tc gridSpan="2">
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Merged 2 cols</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr marL="182880" marR="91440" marT="91440" marB="45720" anchor="ctr">
              <a:lnL w="25400"><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill></a:lnL>
              <a:lnR w="25400"><a:solidFill><a:srgbClr val="0000FF"/></a:solidFill></a:lnR>
              <a:lnT w="38100"><a:solidFill><a:srgbClr val="00AA00"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="FF8800"/></a:solidFill></a:lnB>
              <a:solidFill><a:srgbClr val="DDEEFF"/></a:solidFill>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Normal cell</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr anchor="b">
              <a:lnL w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="333333"/></a:solidFill></a:lnB>
              <a:solidFill><a:srgbClr val="FFFFDD"/></a:solidFill>
            </a:tcPr>
          </a:tc>
        </a:tr>
        <!-- Row 2: vertical merge start (rowSpan=2) + noFill border -->
        <a:tr h="914400">
          <a:tc rowSpan="2">
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Merged 2 rows</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr anchor="ctr">
              <a:lnL w="19050"><a:solidFill><a:srgbClr val="9900CC"/></a:solidFill></a:lnL>
              <a:lnR w="19050"><a:solidFill><a:srgbClr val="9900CC"/></a:solidFill></a:lnR>
              <a:lnT w="19050"><a:solidFill><a:srgbClr val="9900CC"/></a:solidFill></a:lnT>
              <a:lnB w="19050"><a:solidFill><a:srgbClr val="9900CC"/></a:solidFill></a:lnB>
              <a:solidFill><a:srgbClr val="F0E0FF"/></a:solidFill>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Cell B2</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnL w="12700"><a:noFill/></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnB>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Cell C2</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnL w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnB>
            </a:tcPr>
          </a:tc>
        </a:tr>
        <!-- Row 3: vertical merge continuation -->
        <a:tr h="914400">
          <a:tc vMerge="1">
            <a:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Cell B3 (large margin)</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr marL="274320" marT="137160">
              <a:lnL w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnB>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Cell C3 (top anchor)</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr anchor="t">
              <a:lnL w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnL>
              <a:lnR w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnR>
              <a:lnT w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnT>
              <a:lnB w="12700"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:lnB>
              <a:solidFill><a:srgbClr val="EEFFEE"/></a:solidFill>
            </a:tcPr>
          </a:tc>
        </a:tr>
      </a:tbl>
    </a:graphicData>
  </a:graphic>
</p:graphicFrame>'''

spTree41 = slide41._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree41.append(etree.fromstring(table_xml))

lbl41 = slide41.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl41.text_frame.paragraphs[0].text = "Table: cell merge (gridSpan/rowSpan/vMerge) + borders + margins + anchor"
lbl41.text_frame.paragraphs[0].font.size = Pt(10)
lbl41.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 42: Table diagonal borders + tblPr flags + tblStyleId ──────────────
slide42 = prs.slides.add_slide(blank)

table42_xml = f'''<p:graphicFrame xmlns:a="{ns_a[1:-1]}" xmlns:p="{ns_p[1:-1]}" xmlns:r="{ns_r[1:-1]}">
  <p:nvGraphicFramePr>
    <p:cNvPr id="200" name="Table42"/>
    <p:cNvGraphicFramePr><a:graphicFrameLocks noGrp="1"/></p:cNvGraphicFramePr>
    <p:nvPr/>
  </p:nvGraphicFramePr>
  <p:xfrm>
    <a:off x="457200" y="457200"/>
    <a:ext cx="8229600" cy="3657600"/>
  </p:xfrm>
  <a:graphic>
    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
      <a:tbl>
        <a:tblPr firstRow="1" lastRow="1" firstCol="0" lastCol="0" bandRow="1" bandCol="0"
                 tblStyleId="{{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}}"/>
        <a:tblGrid>
          <a:gridCol w="2743200"/>
          <a:gridCol w="2743200"/>
          <a:gridCol w="2743200"/>
        </a:tblGrid>
        <!-- Row 1: diagonal borders -->
        <a:tr h="914400">
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200" b="1"/><a:t>TL-BR diagonal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnTlToBr w="19050"><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill></a:lnTlToBr>
              <a:solidFill><a:srgbClr val="FFF0F0"/></a:solidFill>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200" b="1"/><a:t>BL-TR diagonal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnBlToTr w="19050"><a:solidFill><a:srgbClr val="0000FF"/></a:solidFill></a:lnBlToTr>
              <a:solidFill><a:srgbClr val="F0F0FF"/></a:solidFill>
            </a:tcPr>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200" b="1"/><a:t>Both diagonals</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr>
              <a:lnTlToBr w="12700"><a:solidFill><a:srgbClr val="008800"/></a:solidFill></a:lnTlToBr>
              <a:lnBlToTr w="12700"><a:solidFill><a:srgbClr val="880000"/></a:solidFill></a:lnBlToTr>
              <a:solidFill><a:srgbClr val="FFFFF0"/></a:solidFill>
            </a:tcPr>
          </a:tc>
        </a:tr>
        <!-- Row 2: band row test -->
        <a:tr h="914400">
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Band row 1</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Normal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Normal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
        </a:tr>
        <!-- Row 3: band row test -->
        <a:tr h="914400">
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Band row 2</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Normal</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
          <a:tc>
            <a:txBody><a:bodyPr/><a:lstStyle/>
              <a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Last row cell</a:t></a:r></a:p>
            </a:txBody>
            <a:tcPr/>
          </a:tc>
        </a:tr>
      </a:tbl>
    </a:graphicData>
  </a:graphic>
</p:graphicFrame>'''

spTree42 = slide42._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree42.append(etree.fromstring(table42_xml))

lbl42 = slide42.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9), Inches(0.5))
lbl42.text_frame.paragraphs[0].text = "Table: diagonal borders (lnTlToBr/lnBlToTr) + tblPr flags + tblStyleId"
lbl42.text_frame.paragraphs[0].font.size = Pt(10)
lbl42.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ── Slide 43: Image crop (srcRect) + alpha (alphaModFix) ──────────────────────

slide43 = prs.slides.add_slide(blank)

# Create a test image: 2x2 pixel PNG with 4 colored quadrants (red/green/blue/yellow)
def make_test_png_4color():
    """Create a 4x4 PNG with colored quadrants: TL=red, TR=green, BL=blue, BR=yellow."""
    width, height = 4, 4
    def px(r, g, b):
        return bytes([r, g, b])
    rows = []
    for y in range(height):
        row = b'\x00'  # filter byte
        for x in range(width):
            if y < 2 and x < 2:
                row += px(255, 0, 0)      # TL = red
            elif y < 2 and x >= 2:
                row += px(0, 255, 0)      # TR = green
            elif y >= 2 and x < 2:
                row += px(0, 0, 255)      # BL = blue
            else:
                row += px(255, 255, 0)    # BR = yellow
        rows.append(row)
    raw = b''.join(rows)
    def png_chunk(chunk_type, data):
        c = chunk_type + data
        crc = struct.pack('>I', zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack('>I', len(data)) + c + crc
    ihdr = struct.pack('>IIBBBBB', width, height, 8, 2, 0, 0, 0)
    return (b'\x89PNG\r\n\x1a\n' +
            png_chunk(b'IHDR', ihdr) +
            png_chunk(b'IDAT', zlib.compress(raw)) +
            png_chunk(b'IEND', b''))

test_png = make_test_png_4color()
tmp_img43 = os.path.join(tempfile.gettempdir(), '_test_crop.png')
with open(tmp_img43, 'wb') as f:
    f.write(test_png)

img_part43, img_rid43 = slide43.part.get_or_add_image_part(tmp_img43)
os.unlink(tmp_img43)

# Shape 1: Picture with srcRect crop (crop 25% from each side — shows center)
ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

pic43_crop = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="100" name="CropPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid43}"/>
    <a:srcRect l="25000" t="25000" r="25000" b="25000"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="914400"/>
      <a:ext cx="2743200" cy="2743200"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 2: Picture with alphaModFix (50% opacity)
pic43_alpha = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="101" name="AlphaPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid43}">
      <a:alphaModFix amt="50000"/>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="3657600" y="914400"/>
      <a:ext cx="2743200" cy="2743200"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 3: Picture with crop + alpha combined (crop left 50%, alpha 75%)
pic43_both = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="102" name="CropAlphaPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid43}">
      <a:alphaModFix amt="75000"/>
    </a:blip>
    <a:srcRect l="50000"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="6858000" y="914400"/>
      <a:ext cx="2286000" cy="2743200"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 4: AutoShape with blipFill + srcRect crop
s43_autoshape = slide43.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(3), Inches(2.5))
s43_autoshape.text = ""
set_fill_xml(s43_autoshape, f'''<a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <a:blip r:embed="{img_rid43}"/>
  <a:srcRect l="0" t="50000" r="0" b="0"/>
  <a:stretch><a:fillRect/></a:stretch>
</a:blipFill>''')

spTree43 = slide43._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree43.append(etree.fromstring(pic43_crop))
spTree43.append(etree.fromstring(pic43_alpha))
spTree43.append(etree.fromstring(pic43_both))

lbl43 = slide43.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl43.text_frame.paragraphs[0].text = "Slide 43: Image crop (srcRect) + alpha (alphaModFix)"
lbl43.text_frame.paragraphs[0].font.size = Pt(18)
lbl43.text_frame.paragraphs[0].font.bold = True

# ── Slide 44: External image reference (TargetMode="External") ───────────────

slide44 = prs.slides.add_slide(blank)

# Use stable, CORS-friendly public image URLs
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Image 1: Wikimedia Commons PNG (small, stable, CORS-friendly)
ext_url1 = "https://upload.wikimedia.org/wikipedia/commons/thumb/4/47/PNG_transparency_demonstration_1.png/280px-PNG_transparency_demonstration_1.png"
ext_rid1 = slide44.part.relate_to(ext_url1, RT.IMAGE, is_external=True)

# Image 2: picsum.photos direct image (stable, specific ID = no redirect)
ext_url2 = "https://picsum.photos/id/237/200/300"
ext_rid2 = slide44.part.relate_to(ext_url2, RT.IMAGE, is_external=True)

pic44_ext1 = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="200" name="ExtPic1"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{ext_rid1}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="914400"/>
      <a:ext cx="3657600" cy="3657600"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

pic44_ext2 = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="201" name="ExtPic2"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{ext_rid2}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="4572000" y="914400"/>
      <a:ext cx="2743200" cy="3657600"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

spTree44 = slide44._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree44.append(etree.fromstring(pic44_ext1))
spTree44.append(etree.fromstring(pic44_ext2))

lbl44 = slide44.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl44.text_frame.paragraphs[0].text = "Slide 44: External image reference (TargetMode=External)"
lbl44.text_frame.paragraphs[0].font.size = Pt(18)
lbl44.text_frame.paragraphs[0].font.bold = True

lbl44b = slide44.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.5))
lbl44b.text_frame.paragraphs[0].text = f"Left: Wikimedia Commons PNG\nRight: picsum.photos (id=237)"
lbl44b.text_frame.paragraphs[0].font.size = Pt(12)

# ── Slide 45: Image effects — brightness/contrast (a:lum) + duotone (a:duotone) ──
slide45 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Reuse the test PNG for effects
tmp_img45 = os.path.join(tempfile.gettempdir(), '_test_effects.png')
with open(tmp_img45, 'wb') as f:
    f.write(test_png)
img_part45, img_rid45 = slide45.part.get_or_add_image_part(tmp_img45)
os.unlink(tmp_img45)
# Add second image for duotone
tmp_img45b = os.path.join(tempfile.gettempdir(), '_test_effects2.png')
with open(tmp_img45b, 'wb') as f:
    f.write(test_png)
img_part45b, img_rid45b = slide45.part.get_or_add_image_part(tmp_img45b)
os.unlink(tmp_img45b)
# Third image for combined brightness + contrast
tmp_img45c = os.path.join(tempfile.gettempdir(), '_test_effects3.png')
with open(tmp_img45c, 'wb') as f:
    f.write(test_png)
img_part45c, img_rid45c = slide45.part.get_or_add_image_part(tmp_img45c)
os.unlink(tmp_img45c)

# Shape 1: Picture with brightness +50% (a:lum bright="50000")
pic45_bright = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="300" name="BrightPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid45}">
      <a:lum bright="50000"/>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="914400"/>
      <a:ext cx="2286000" cy="2286000"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 2: Picture with contrast -30% (a:lum contrast="-30000")
pic45_contrast = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="301" name="ContrastPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid45b}">
      <a:lum contrast="-30000"/>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="3429000" y="914400"/>
      <a:ext cx="2286000" cy="2286000"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 3: Picture with bright+contrast combined (a:lum bright="20000" contrast="40000")
pic45_both = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="302" name="BrightContrastPic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid45c}">
      <a:lum bright="20000" contrast="40000"/>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="6400800" y="914400"/>
      <a:ext cx="2286000" cy="2286000"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

spTree45 = slide45._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree45.append(etree.fromstring(pic45_bright))
spTree45.append(etree.fromstring(pic45_contrast))
spTree45.append(etree.fromstring(pic45_both))

lbl45 = slide45.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl45.text_frame.paragraphs[0].text = "Slide 45: Image effects — brightness/contrast (a:lum)"
lbl45.text_frame.paragraphs[0].font.size = Pt(18)
lbl45.text_frame.paragraphs[0].font.bold = True

lbl45b = slide45.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(1))
lbl45b.text_frame.paragraphs[0].text = "Left: bright +50% | Center: contrast -30% | Right: bright +20% contrast +40%"
lbl45b.text_frame.paragraphs[0].font.size = Pt(12)

# ── Slide 46: Duotone (a:duotone) + color change (a:clrChange) ──
slide46 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Image for duotone
tmp_img46 = os.path.join(tempfile.gettempdir(), '_test_duotone.png')
with open(tmp_img46, 'wb') as f:
    f.write(test_png)
img_part46, img_rid46 = slide46.part.get_or_add_image_part(tmp_img46)
os.unlink(tmp_img46)
# Image for clrChange
tmp_img46b = os.path.join(tempfile.gettempdir(), '_test_clrchange.png')
with open(tmp_img46b, 'wb') as f:
    f.write(test_png)
img_part46b, img_rid46b = slide46.part.get_or_add_image_part(tmp_img46b)
os.unlink(tmp_img46b)

# Shape 1: Picture with duotone (dark blue → light yellow)
pic46_duotone = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="400" name="DuotonePic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid46}">
      <a:duotone>
        <a:srgbClr val="000080"/>
        <a:srgbClr val="FFFF00"/>
      </a:duotone>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="457200" y="914400"/>
      <a:ext cx="3657600" cy="3657600"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

# Shape 2: Picture with clrChange (white → transparent)
pic46_clrchange = f'''<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvPicPr>
    <p:cNvPr id="401" name="ClrChangePic"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{img_rid46b}">
      <a:clrChange>
        <a:clrFrom><a:srgbClr val="FF0000"/></a:clrFrom>
        <a:clrTo><a:srgbClr val="00FF00"/></a:clrTo>
      </a:clrChange>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="4572000" y="914400"/>
      <a:ext cx="3657600" cy="3657600"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''

spTree46 = slide46._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree46.append(etree.fromstring(pic46_duotone))
spTree46.append(etree.fromstring(pic46_clrchange))

lbl46 = slide46.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl46.text_frame.paragraphs[0].text = "Slide 46: Duotone (a:duotone) + Color change (a:clrChange)"
lbl46.text_frame.paragraphs[0].font.size = Pt(18)
lbl46.text_frame.paragraphs[0].font.bold = True

lbl46b = slide46.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1))
lbl46b.text_frame.paragraphs[0].text = "Left: Duotone (navy→yellow) | Right: clrChange (red→green, data preserved)"
lbl46b.text_frame.paragraphs[0].font.size = Pt(12)

########################################
# Slide 47: Background pattern fill
########################################
slide47 = prs.slides.add_slide(prs.slide_layouts[6])

# Inject pattern fill background via raw XML
ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
bg47_xml = f"""<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:bgPr>
    <a:pattFill prst="ltDnDiag">
      <a:fgClr><a:srgbClr val="3366CC"/></a:fgClr>
      <a:bgClr><a:srgbClr val="FFFFFF"/></a:bgClr>
    </a:pattFill>
    <a:effectLst/>
  </p:bgPr>
</p:bg>"""
csld47 = slide47._element.find(f'.//{ns_p}cSld')
# Remove existing bg if any
old_bg47 = csld47.find(f'{ns_p}bg')
if old_bg47 is not None:
    csld47.remove(old_bg47)
bg47_elem = etree.fromstring(bg47_xml)
csld47.insert(0, bg47_elem)

lbl47 = slide47.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl47.text_frame.paragraphs[0].text = "Slide 47: Background pattern fill (ltDnDiag)"
lbl47.text_frame.paragraphs[0].font.size = Pt(18)
lbl47.text_frame.paragraphs[0].font.bold = True

########################################
# Slide 48: Line gradient / pattern fill
########################################
slide48 = prs.slides.add_slide(prs.slide_layouts[6])

# Shape 1: Rectangle with gradient stroke
sp48_grad_xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="100" name="GradStrokeRect"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="3657600" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="F0F0F0"/></a:solidFill>
    <a:ln w="76200">
      <a:gradFill>
        <a:gsLst>
          <a:gs pos="0"><a:srgbClr val="FF0000"/></a:gs>
          <a:gs pos="50000"><a:srgbClr val="00FF00"/></a:gs>
          <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
        </a:gsLst>
        <a:lin ang="0" scaled="1"/>
      </a:gradFill>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Gradient stroke</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

# Shape 2: Ellipse with pattern stroke
sp48_patt_xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="101" name="PattStrokeEllipse"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="1371600"/><a:ext cx="3657600" cy="2286000"/></a:xfrm>
    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFCC"/></a:solidFill>
    <a:ln w="57150">
      <a:pattFill prst="smCheck">
        <a:fgClr><a:srgbClr val="990000"/></a:fgClr>
        <a:bgClr><a:srgbClr val="FFCC00"/></a:bgClr>
      </a:pattFill>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:r><a:rPr lang="en-US" sz="1400"/><a:t>Pattern stroke</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

spTree48 = slide48._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree48.append(etree.fromstring(sp48_grad_xml))
spTree48.append(etree.fromstring(sp48_patt_xml))

lbl48 = slide48.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl48.text_frame.paragraphs[0].text = "Slide 48: Line gradient fill + Line pattern fill"
lbl48.text_frame.paragraphs[0].font.size = Pt(18)
lbl48.text_frame.paragraphs[0].font.bold = True

########################################
# Slide 49: Shape hyperlinks + color modifiers
########################################
slide49 = prs.slides.add_slide(prs.slide_layouts[6])

# Shape with hlinkClick on cNvPr
sp49_link_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="200" name="LinkedShape">
      <a:hlinkClick r:id="rId99"/>
    </p:cNvPr>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="3657600" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    <a:ln w="25400"><a:solidFill><a:srgbClr val="2F5597"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1600" b="1"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:rPr><a:t>Click me (shape link)</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

# Shape with color modifier: complement
sp49_comp_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="201" name="CompColor"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="1371600"/><a:ext cx="3657600" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FF6600"><a:comp/></a:srgbClr></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>comp(#FF6600) = complementary hue</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

# Shape with color modifier: inv
sp49_inv_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="202" name="InvColor"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="3657600"/><a:ext cx="3657600" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="3366CC"><a:inv/></a:srgbClr></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>inv(#3366CC) → #CC9933</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

# Shape with hueMod
sp49_huemod_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="203" name="HueModColor"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="3657600"/><a:ext cx="3657600" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FF0000"><a:hueMod val="50000"/></a:srgbClr></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400"/><a:t>hueMod 50% of red</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

spTree49 = slide49._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree49.append(etree.fromstring(sp49_link_xml))
spTree49.append(etree.fromstring(sp49_comp_xml))
spTree49.append(etree.fromstring(sp49_inv_xml))
spTree49.append(etree.fromstring(sp49_huemod_xml))

lbl49 = slide49.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl49.text_frame.paragraphs[0].text = "Slide 49: Shape hyperlinks + Color modifiers (comp/inv/hueMod)"
lbl49.text_frame.paragraphs[0].font.size = Pt(18)
lbl49.text_frame.paragraphs[0].font.bold = True

# Add a fake relationship for the shape link (rId99)
# We need to add this to the slide49 rels
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
slide49_part = slide49.part
slide49_part.rels.get_or_add_ext_rel(RT.HYPERLINK, 'https://example.com/shape-link')

# ── Slide 50: Shape effects (outerShdw, innerShdw, glow, softEdge) ──
slide50 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

sp50_outer_shadow_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="301" name="OuterShadow"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    <a:effectLst>
      <a:outerShdw blurRad="152400" dist="114300" dir="5400000" algn="ctr" rotWithShape="0">
        <a:srgbClr val="000000"><a:alpha val="40000"/></a:srgbClr>
      </a:outerShdw>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Outer Shadow</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp50_inner_shadow_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="302" name="InnerShadow"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="3657600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="70AD47"/></a:solidFill>
    <a:effectLst>
      <a:innerShdw blurRad="114300" dist="76200" dir="2700000">
        <a:srgbClr val="000000"><a:alpha val="50000"/></a:srgbClr>
      </a:innerShdw>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Inner Shadow</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp50_glow_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="303" name="Glow"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="1371600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="ED7D31"/></a:solidFill>
    <a:effectLst>
      <a:glow rad="228600">
        <a:srgbClr val="FFC000"><a:alpha val="60000"/></a:srgbClr>
      </a:glow>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Glow Effect</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp50_softedge_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="304" name="SoftEdge"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="3657600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="A855F7"/></a:solidFill>
    <a:effectLst>
      <a:softEdge rad="317500"/>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>Soft Edge</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

spTree50 = slide50._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree50.append(etree.fromstring(sp50_outer_shadow_xml))
spTree50.append(etree.fromstring(sp50_inner_shadow_xml))
spTree50.append(etree.fromstring(sp50_glow_xml))
spTree50.append(etree.fromstring(sp50_softedge_xml))

lbl50 = slide50.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl50.text_frame.paragraphs[0].text = "Slide 50: Shape effects (outerShdw / innerShdw / glow / softEdge)"
lbl50.text_frame.paragraphs[0].font.size = Pt(18)
lbl50.text_frame.paragraphs[0].font.bold = True

# ── Slide 51: 3D effects (bevel, scene3d, sp3d) + text shadow ──
slide51 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

sp51_bevel_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="401" name="Bevel3D"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="1371600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    <a:scene3d>
      <a:camera prst="orthographicFront"/>
      <a:lightRig rig="threePt" dir="t"/>
    </a:scene3d>
    <a:sp3d prstMaterial="plastic">
      <a:bevelT w="152400" h="50800" prst="relaxedInset"/>
    </a:sp3d>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>3D Bevel</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp51_text_shadow_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="402" name="TextShadow"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="5029200" y="1371600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="2400" b="1">
      <a:solidFill><a:srgbClr val="333333"/></a:solidFill>
      <a:effectLst>
        <a:outerShdw blurRad="38100" dist="38100" dir="5400000">
          <a:srgbClr val="000000"><a:alpha val="40000"/></a:srgbClr>
        </a:outerShdw>
      </a:effectLst>
    </a:rPr><a:t>Text with Shadow</a:t></a:r></a:p>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="2000">
      <a:solidFill><a:srgbClr val="FF4444"/></a:solidFill>
      <a:effectLst>
        <a:glow rad="101600">
          <a:srgbClr val="FF0000"><a:alpha val="50000"/></a:srgbClr>
        </a:glow>
      </a:effectLst>
    </a:rPr><a:t>Text with Glow</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

sp51_extrusion_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="403" name="Extrusion3D"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="3657600"/><a:ext cx="3200400" cy="1828800"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="70AD47"/></a:solidFill>
    <a:scene3d>
      <a:camera prst="isometricLeftDown"/>
      <a:lightRig rig="balanced" dir="t"/>
    </a:scene3d>
    <a:sp3d extrusionH="76200" contourW="12700" prstMaterial="warmMatte">
      <a:bevelT w="101600" h="38100" prst="circle"/>
      <a:bevelB w="50800" h="25400"/>
      <a:extrusionClr><a:srgbClr val="5B8C3C"/></a:extrusionClr>
      <a:contourClr><a:srgbClr val="2E5B14"/></a:contourClr>
    </a:sp3d>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/><a:r><a:rPr lang="en-US" sz="1400" b="1"/><a:t>3D Extrusion</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""

spTree51 = slide51._element.find(f'.//{ns_p}cSld/{ns_p}spTree')
spTree51.append(etree.fromstring(sp51_bevel_xml))
spTree51.append(etree.fromstring(sp51_text_shadow_xml))
spTree51.append(etree.fromstring(sp51_extrusion_xml))

lbl51 = slide51.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl51.text_frame.paragraphs[0].text = "Slide 51: 3D effects (bevel/extrusion/scene3d) + text shadow/glow"
lbl51.text_frame.paragraphs[0].font.size = Pt(18)
lbl51.text_frame.paragraphs[0].font.bold = True

# ── Slide 52: Bar/Column chart ────────────────────────────────────────────────

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

slide52 = prs.slides.add_slide(blank)

chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Sales 2024', (120, 150, 180, 200))
chart_data.add_series('Sales 2025', (100, 130, 160, 190))

chart_frame = slide52.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(5), Inches(4), chart_data
)
chart52 = chart_frame.chart
chart52.has_legend = True

lbl52 = slide52.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl52.text_frame.paragraphs[0].text = "Slide 52: Column chart (clustered)"
lbl52.text_frame.paragraphs[0].font.size = Pt(18)
lbl52.text_frame.paragraphs[0].font.bold = True

# ── Slide 53: Line chart ─────────────────────────────────────────────────────

slide53 = prs.slides.add_slide(blank)

chart_data2 = CategoryChartData()
chart_data2.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May']
chart_data2.add_series('Revenue', (80, 95, 110, 105, 130))
chart_data2.add_series('Cost', (60, 65, 70, 75, 80))

chart_frame2 = slide53.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(1), Inches(1.5),
    Inches(5), Inches(4), chart_data2
)
chart53 = chart_frame2.chart
chart53.has_legend = True

lbl53 = slide53.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl53.text_frame.paragraphs[0].text = "Slide 53: Line chart"
lbl53.text_frame.paragraphs[0].font.size = Pt(18)
lbl53.text_frame.paragraphs[0].font.bold = True

# ── Slide 54: Pie chart ──────────────────────────────────────────────────────

slide54 = prs.slides.add_slide(blank)

chart_data3 = CategoryChartData()
chart_data3.categories = ['Desktop', 'Mobile', 'Tablet', 'Other']
chart_data3.add_series('Market Share', (45, 35, 15, 5))

chart_frame3 = slide54.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(2), Inches(1.5),
    Inches(5), Inches(4), chart_data3
)
chart54 = chart_frame3.chart
chart54.has_legend = True

lbl54 = slide54.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl54.text_frame.paragraphs[0].text = "Slide 54: Pie chart"
lbl54.text_frame.paragraphs[0].font.size = Pt(18)
lbl54.text_frame.paragraphs[0].font.bold = True

# ── Slide 55: Bar chart (horizontal) + Donut chart ──────────────────────────

slide55 = prs.slides.add_slide(blank)

chart_data4 = CategoryChartData()
chart_data4.categories = ['Product A', 'Product B', 'Product C']
chart_data4.add_series('Units Sold', (250, 180, 320))

chart_frame4 = slide55.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.5),
    Inches(4), Inches(4), chart_data4
)

chart_data5 = CategoryChartData()
chart_data5.categories = ['Yes', 'No', 'Maybe']
chart_data5.add_series('Responses', (60, 25, 15))

chart_frame5 = slide55.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, Inches(5), Inches(1.5),
    Inches(4), Inches(4), chart_data5
)

lbl55 = slide55.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl55.text_frame.paragraphs[0].text = "Slide 55: Bar chart (horizontal) + Donut chart"
lbl55.text_frame.paragraphs[0].font.size = Pt(18)
lbl55.text_frame.paragraphs[0].font.bold = True

# ── Slide 56: Column chart with data labels ──────────────────────────────────

slide56 = prs.slides.add_slide(blank)

chart_data6 = CategoryChartData()
chart_data6.categories = ['North', 'South', 'East', 'West']
chart_data6.add_series('Revenue', (320, 280, 190, 410))

chart_frame6 = slide56.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(5), Inches(4), chart_data6
)
chart56 = chart_frame6.chart
chart56.has_legend = True

# Enable data labels with values
from pptx.util import Emu as _Emu
plot = chart56.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.show_value = True
data_labels.show_category_name = False

lbl56 = slide56.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl56.text_frame.paragraphs[0].text = "Slide 56: Column chart with data labels"
lbl56.text_frame.paragraphs[0].font.size = Pt(18)
lbl56.text_frame.paragraphs[0].font.bold = True

# ── Slide 57: Pie chart with data point colors + percentage labels ────────────

slide57 = prs.slides.add_slide(blank)

chart_data7 = CategoryChartData()
chart_data7.categories = ['Chrome', 'Firefox', 'Safari', 'Edge', 'Other']
chart_data7.add_series('Browser Share', (65, 10, 18, 5, 2))

chart_frame7 = slide57.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(2), Inches(1.5),
    Inches(5), Inches(4), chart_data7
)
chart57 = chart_frame7.chart
chart57.has_legend = True

# Set per-slice colors via c:dPt
# Access the chart part XML directly
chart_part57 = chart57.part
chart_xml57 = chart_part57._element
nsmap = {
    'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}
pie_chart = chart_xml57.findall('.//c:pieChart', nsmap)[0]
ser_elem = pie_chart.findall('c:ser', nsmap)[0]

# Add dPt elements with custom colors
dpt_colors = [
    (0, '4285F4'),  # Chrome blue
    (1, 'FF7139'),  # Firefox orange
    (2, '000000'),  # Safari black
    (3, '0078D4'),  # Edge blue
    (4, '888888'),  # Other gray
]
c_ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
for pt_idx, color_hex in dpt_colors:
    dpt = etree.SubElement(ser_elem, f'{{{c_ns}}}dPt')
    idx_elem = etree.SubElement(dpt, f'{{{c_ns}}}idx')
    idx_elem.set('val', str(pt_idx))
    spPr = etree.SubElement(dpt, f'{{{c_ns}}}spPr')
    solidFill = etree.SubElement(spPr, f'{{{a_ns}}}solidFill')
    srgb = etree.SubElement(solidFill, f'{{{a_ns}}}srgbClr')
    srgb.set('val', color_hex)

# Enable percentage labels
plot57 = chart57.plots[0]
plot57.has_data_labels = True
dl57 = plot57.data_labels
dl57.show_percentage = True
dl57.show_value = False
dl57.show_category_name = True

lbl57 = slide57.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl57.text_frame.paragraphs[0].text = "Slide 57: Pie chart with custom colors + % labels"
lbl57.text_frame.paragraphs[0].font.size = Pt(18)
lbl57.text_frame.paragraphs[0].font.bold = True

# ── Slide 58: Line chart with series spPr colors + data labels ────────────────

slide58 = prs.slides.add_slide(blank)

chart_data8 = CategoryChartData()
chart_data8.categories = ['Week 1', 'Week 2', 'Week 3', 'Week 4']
chart_data8.add_series('Actual', (10, 25, 40, 55))
chart_data8.add_series('Target', (15, 30, 45, 60))

chart_frame8 = slide58.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(1), Inches(1.5),
    Inches(5), Inches(4), chart_data8
)
chart58 = chart_frame8.chart
chart58.has_legend = True

# Set explicit series colors
from pptx.chart.series import LineSeries
for i, ser in enumerate(chart58.series):
    ser_format = ser.format
    if i == 0:
        ser_format.line.color.rgb = RGBColor(0x00, 0x88, 0x00)  # Green
    else:
        ser_format.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)  # Red

# Enable data labels showing values
plot58 = chart58.plots[0]
plot58.has_data_labels = True
dl58 = plot58.data_labels
dl58.show_value = True

lbl58 = slide58.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl58.text_frame.paragraphs[0].text = "Slide 58: Line chart with series colors + data labels"
lbl58.text_frame.paragraphs[0].font.size = Pt(18)
lbl58.text_frame.paragraphs[0].font.bold = True

# ── Slide 59: Scatter chart + Area chart ─────────────────────────────────────

from pptx.chart.data import XyChartData

slide59 = prs.slides.add_slide(blank)

xy_data = XyChartData()
s1 = xy_data.add_series('Group A')
s1.add_data_point(10, 20)
s1.add_data_point(30, 50)
s1.add_data_point(50, 40)
s1.add_data_point(70, 80)
s2 = xy_data.add_series('Group B')
s2.add_data_point(15, 60)
s2.add_data_point(35, 30)
s2.add_data_point(55, 70)
s2.add_data_point(75, 45)

chart_frame59a = slide59.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER, Inches(0.3), Inches(1.5),
    Inches(4.5), Inches(4), xy_data
)
chart59a = chart_frame59a.chart
chart59a.has_legend = True

area_data = CategoryChartData()
area_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May']
area_data.add_series('Downloads', (120, 180, 150, 200, 250))
area_data.add_series('Uploads', (80, 100, 130, 110, 170))

chart_frame59b = slide59.shapes.add_chart(
    XL_CHART_TYPE.AREA, Inches(5.2), Inches(1.5),
    Inches(4.5), Inches(4), area_data
)
chart59b = chart_frame59b.chart
chart59b.has_legend = True

lbl59 = slide59.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl59.text_frame.paragraphs[0].text = "Slide 59: Scatter chart + Area chart"
lbl59.text_frame.paragraphs[0].font.size = Pt(18)
lbl59.text_frame.paragraphs[0].font.bold = True

# ── Slide 60: Radar chart ────────────────────────────────────────────────────

slide60 = prs.slides.add_slide(blank)

radar_data = CategoryChartData()
radar_data.categories = ['Str', 'Dex', 'Con', 'Int', 'Wis', 'Cha']
radar_data.add_series('Character A', (15, 12, 14, 10, 8, 16))
radar_data.add_series('Character B', (10, 16, 8, 15, 14, 12))

chart_frame60 = slide60.shapes.add_chart(
    XL_CHART_TYPE.RADAR, Inches(2), Inches(1.5),
    Inches(5), Inches(4.5), radar_data
)
chart60 = chart_frame60.chart
chart60.has_legend = True

lbl60 = slide60.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl60.text_frame.paragraphs[0].text = "Slide 60: Radar chart"
lbl60.text_frame.paragraphs[0].font.size = Pt(18)
lbl60.text_frame.paragraphs[0].font.bold = True

# ── Slide 61: Bubble chart ──────────────────────────────────────────────────

from pptx.chart.data import BubbleChartData

slide61 = prs.slides.add_slide(blank)

bubble_data = BubbleChartData()
bs1 = bubble_data.add_series('Product X')
bs1.add_data_point(10, 40, 15)
bs1.add_data_point(25, 60, 30)
bs1.add_data_point(40, 30, 10)
bs1.add_data_point(55, 70, 25)
bs2 = bubble_data.add_series('Product Y')
bs2.add_data_point(15, 55, 20)
bs2.add_data_point(30, 35, 12)
bs2.add_data_point(50, 50, 35)

chart_frame61 = slide61.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE, Inches(1.5), Inches(1.5),
    Inches(6), Inches(4.5), bubble_data
)
chart61 = chart_frame61.chart
chart61.has_legend = True

lbl61 = slide61.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl61.text_frame.paragraphs[0].text = "Slide 61: Bubble chart"
lbl61.text_frame.paragraphs[0].font.size = Pt(18)
lbl61.text_frame.paragraphs[0].font.bold = True

# ── Slide 62: Stock chart (OHLC) ───────────────────────────────────────────

slide62 = prs.slides.add_slide(blank)

# python-pptx doesn't have direct stock chart support, so inject via lxml
# Create a bar chart as placeholder, then replace with stock chart XML
stock_placeholder = CategoryChartData()
stock_placeholder.categories = ['Day 1', 'Day 2', 'Day 3', 'Day 4', 'Day 5']
stock_placeholder.add_series('Open', (100, 105, 98, 110, 108))

chart_frame62 = slide62.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), stock_placeholder
)
chart62 = chart_frame62.chart

# Replace chart XML with a stock chart
chart_part62 = chart62.part
chart_elem62 = chart_part62._element
ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Find plotArea and remove existing barChart
plot_area = chart_elem62.find(f'.//{{{ns_c}}}plotArea')
for bar in plot_area.findall(f'{{{ns_c}}}barChart'):
    plot_area.remove(bar)

# Build stock chart element
stock_xml = f"""<c:stockChart xmlns:c="{ns_c}" xmlns:a="{ns_a}">
  <c:ser>
    <c:idx val="0"/><c:order val="0"/>
    <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>Open</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$6</c:f><c:strCache><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>Day 1</c:v></c:pt><c:pt idx="1"><c:v>Day 2</c:v></c:pt>
      <c:pt idx="2"><c:v>Day 3</c:v></c:pt><c:pt idx="3"><c:v>Day 4</c:v></c:pt>
      <c:pt idx="4"><c:v>Day 5</c:v></c:pt></c:strCache></c:strRef></c:cat>
    <c:val><c:numRef><c:f>Sheet1!$B$2:$B$6</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="5"/>
      <c:pt idx="0"><c:v>100</c:v></c:pt><c:pt idx="1"><c:v>105</c:v></c:pt>
      <c:pt idx="2"><c:v>98</c:v></c:pt><c:pt idx="3"><c:v>110</c:v></c:pt>
      <c:pt idx="4"><c:v>108</c:v></c:pt></c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="1"/><c:order val="1"/>
    <c:tx><c:strRef><c:f>Sheet1!$C$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>High</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$C$2:$C$6</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="5"/>
      <c:pt idx="0"><c:v>115</c:v></c:pt><c:pt idx="1"><c:v>118</c:v></c:pt>
      <c:pt idx="2"><c:v>112</c:v></c:pt><c:pt idx="3"><c:v>120</c:v></c:pt>
      <c:pt idx="4"><c:v>122</c:v></c:pt></c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="2"/><c:order val="2"/>
    <c:tx><c:strRef><c:f>Sheet1!$D$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>Low</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$D$2:$D$6</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="5"/>
      <c:pt idx="0"><c:v>95</c:v></c:pt><c:pt idx="1"><c:v>100</c:v></c:pt>
      <c:pt idx="2"><c:v>90</c:v></c:pt><c:pt idx="3"><c:v>105</c:v></c:pt>
      <c:pt idx="4"><c:v>102</c:v></c:pt></c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="3"/><c:order val="3"/>
    <c:tx><c:strRef><c:f>Sheet1!$E$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>Close</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$E$2:$E$6</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="5"/>
      <c:pt idx="0"><c:v>108</c:v></c:pt><c:pt idx="1"><c:v>102</c:v></c:pt>
      <c:pt idx="2"><c:v>107</c:v></c:pt><c:pt idx="3"><c:v>115</c:v></c:pt>
      <c:pt idx="4"><c:v>118</c:v></c:pt></c:numCache></c:numRef></c:val>
  </c:ser>
  <c:axId val="10000001"/>
  <c:axId val="10000002"/>
</c:stockChart>"""

stock_elem = etree.fromstring(stock_xml)

# Find existing axes and remove to add stock-compatible ones
for ax in plot_area.findall(f'{{{ns_c}}}catAx') + plot_area.findall(f'{{{ns_c}}}valAx'):
    plot_area.remove(ax)

# Insert stock chart after any layout element
layout_elem = plot_area.find(f'{{{ns_c}}}layout')
if layout_elem is not None:
    idx_pos = list(plot_area).index(layout_elem) + 1
    plot_area.insert(idx_pos, stock_elem)
else:
    plot_area.insert(0, stock_elem)

# Add catAx and valAx for stock chart
cat_ax_xml = f"""<c:catAx xmlns:c="{ns_c}">
  <c:axId val="10000001"/>
  <c:scaling><c:orientation val="minMax"/></c:scaling>
  <c:delete val="0"/>
  <c:axPos val="b"/>
  <c:crossAx val="10000002"/>
</c:catAx>"""
val_ax_xml = f"""<c:valAx xmlns:c="{ns_c}">
  <c:axId val="10000002"/>
  <c:scaling><c:orientation val="minMax"/></c:scaling>
  <c:delete val="0"/>
  <c:axPos val="l"/>
  <c:crossAx val="10000001"/>
  <c:majorGridlines/>
</c:valAx>"""
plot_area.append(etree.fromstring(cat_ax_xml))
plot_area.append(etree.fromstring(val_ax_xml))

# Set legend
chart_node62 = chart_elem62.find(f'.//{{{ns_c}}}chart')
legend_xml = f"""<c:legend xmlns:c="{ns_c}">
  <c:legendPos val="r"/>
  <c:overlay val="0"/>
</c:legend>"""
chart_node62.append(etree.fromstring(legend_xml))

lbl62 = slide62.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl62.text_frame.paragraphs[0].text = "Slide 62: Stock chart (OHLC)"
lbl62.text_frame.paragraphs[0].font.size = Pt(18)
lbl62.text_frame.paragraphs[0].font.bold = True

# ── Slide 63: All chart types overview ──────────────────────────────────────

slide63 = prs.slides.add_slide(blank)

# Mini column chart (top-left)
mini_col = CategoryChartData()
mini_col.categories = ['A', 'B', 'C']
mini_col.add_series('S1', (30, 50, 40))

slide63.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.3), Inches(1.2),
    Inches(3), Inches(2.5), mini_col
)

# Mini line chart (top-center)
mini_line = CategoryChartData()
mini_line.categories = ['A', 'B', 'C', 'D']
mini_line.add_series('S1', (10, 30, 20, 40))

slide63.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(3.5), Inches(1.2),
    Inches(3), Inches(2.5), mini_line
)

# Mini pie chart (top-right)
mini_pie = CategoryChartData()
mini_pie.categories = ['X', 'Y', 'Z']
mini_pie.add_series('S1', (40, 35, 25))

slide63.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(6.7), Inches(1.2),
    Inches(3), Inches(2.5), mini_pie
)

# Mini bar chart (bottom-left)
mini_bar = CategoryChartData()
mini_bar.categories = ['P', 'Q', 'R']
mini_bar.add_series('S1', (25, 40, 35))

slide63.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.3), Inches(4.0),
    Inches(3), Inches(2.5), mini_bar
)

# Mini scatter chart (bottom-center)
mini_scatter = XyChartData()
ms1 = mini_scatter.add_series('S1')
ms1.add_data_point(5, 10)
ms1.add_data_point(15, 30)
ms1.add_data_point(25, 20)
ms1.add_data_point(35, 40)

slide63.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER, Inches(3.5), Inches(4.0),
    Inches(3), Inches(2.5), mini_scatter
)

# Mini donut chart (bottom-right)
mini_donut = CategoryChartData()
mini_donut.categories = ['A', 'B', 'C']
mini_donut.add_series('S1', (50, 30, 20))

slide63.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, Inches(6.7), Inches(4.0),
    Inches(3), Inches(2.5), mini_donut
)

lbl63 = slide63.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl63.text_frame.paragraphs[0].text = "Slide 63: All chart types overview"
lbl63.text_frame.paragraphs[0].font.size = Pt(18)
lbl63.text_frame.paragraphs[0].font.bold = True

# ── Slide 64: Line chart with linear trendline ──────────────────────────────

slide64 = prs.slides.add_slide(blank)

trend_data = CategoryChartData()
trend_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
trend_data.add_series('Sales', (20, 35, 28, 45, 52, 60))

chart_frame64 = slide64.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), trend_data
)
chart64 = chart_frame64.chart
chart64.has_legend = True

# Add trendline via lxml (python-pptx has limited trendline support)
chart_part64 = chart64.part
chart_elem64 = chart_part64._element
ns_c64 = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
ns_a64 = 'http://schemas.openxmlformats.org/drawingml/2006/main'

line_chart64 = chart_elem64.findall(f'.//{{{ns_c64}}}lineChart')[0]
ser64 = line_chart64.findall(f'{{{ns_c64}}}ser')[0]

trendline_xml = f"""<c:trendline xmlns:c="{ns_c64}" xmlns:a="{ns_a64}">
  <c:trendlineType val="linear"/>
  <c:spPr>
    <a:ln w="12700">
      <a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>
      <a:prstDash val="dash"/>
    </a:ln>
  </c:spPr>
</c:trendline>"""
ser64.append(etree.fromstring(trendline_xml))

lbl64 = slide64.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl64.text_frame.paragraphs[0].text = "Slide 64: Line chart with linear trendline"
lbl64.text_frame.paragraphs[0].font.size = Pt(18)
lbl64.text_frame.paragraphs[0].font.bold = True

# ── Slide 65: Column chart with error bars ──────────────────────────────────

slide65 = prs.slides.add_slide(blank)

err_data = CategoryChartData()
err_data.categories = ['Group A', 'Group B', 'Group C', 'Group D']
err_data.add_series('Measurement', (85, 92, 78, 95))

chart_frame65 = slide65.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), err_data
)
chart65 = chart_frame65.chart
chart65.has_legend = True

# Add error bars via lxml
chart_part65 = chart65.part
chart_elem65 = chart_part65._element
ns_c65 = ns_c64

bar_chart65 = chart_elem65.findall(f'.//{{{ns_c65}}}barChart')[0]
ser65 = bar_chart65.findall(f'{{{ns_c65}}}ser')[0]

errbar_xml = f"""<c:errBars xmlns:c="{ns_c65}">
  <c:errDir val="y"/>
  <c:errBarType val="both"/>
  <c:errValType val="fixedVal"/>
  <c:val val="8"/>
</c:errBars>"""
ser65.append(etree.fromstring(errbar_xml))

lbl65 = slide65.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl65.text_frame.paragraphs[0].text = "Slide 65: Column chart with error bars (±8)"
lbl65.text_frame.paragraphs[0].font.size = Pt(18)
lbl65.text_frame.paragraphs[0].font.bold = True

# ── Slide 66: Composite chart (column + line on same plot) ───────────────────

slide66 = prs.slides.add_slide(blank)

# Create a column chart first
combo_data = CategoryChartData()
combo_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
combo_data.add_series('Revenue', (300, 350, 400, 450))

chart_frame66 = slide66.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), combo_data
)
chart66 = chart_frame66.chart
chart66.has_legend = True

# Add a lineChart group to the plotArea via lxml
chart_part66 = chart66.part
chart_elem66 = chart_part66._element
ns_c66 = ns_c64

plot_area66 = chart_elem66.find(f'.//{{{ns_c66}}}plotArea')

# Find existing barChart's axId values
bar_chart66 = plot_area66.find(f'{{{ns_c66}}}barChart')
ax_ids66 = bar_chart66.findall(f'{{{ns_c66}}}axId')
ax_id_vals = [a.get('val') for a in ax_ids66]

line_chart_xml = f"""<c:lineChart xmlns:c="{ns_c66}" xmlns:a="{ns_a64}">
  <c:grouping val="standard"/>
  <c:ser>
    <c:idx val="1"/><c:order val="1"/>
    <c:tx><c:strRef><c:f>Sheet1!$C$1</c:f><c:strCache><c:ptCount val="1"/>
      <c:pt idx="0"><c:v>Profit</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:spPr><a:ln w="25400"><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill></a:ln></c:spPr>
    <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$5</c:f><c:strCache><c:ptCount val="4"/>
      <c:pt idx="0"><c:v>Q1</c:v></c:pt><c:pt idx="1"><c:v>Q2</c:v></c:pt>
      <c:pt idx="2"><c:v>Q3</c:v></c:pt><c:pt idx="3"><c:v>Q4</c:v></c:pt>
    </c:strCache></c:strRef></c:cat>
    <c:val><c:numRef><c:f>Sheet1!$C$2:$C$5</c:f><c:numCache><c:formatCode>General</c:formatCode>
      <c:ptCount val="4"/>
      <c:pt idx="0"><c:v>80</c:v></c:pt><c:pt idx="1"><c:v>120</c:v></c:pt>
      <c:pt idx="2"><c:v>150</c:v></c:pt><c:pt idx="3"><c:v>200</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:axId val="{ax_id_vals[0] if ax_id_vals else '10000001'}"/>
  <c:axId val="{ax_id_vals[1] if len(ax_id_vals) > 1 else '10000002'}"/>
</c:lineChart>"""

# Insert lineChart after barChart
bar_idx = list(plot_area66).index(bar_chart66)
plot_area66.insert(bar_idx + 1, etree.fromstring(line_chart_xml))

lbl66 = slide66.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl66.text_frame.paragraphs[0].text = "Slide 66: Composite chart (column + line)"
lbl66.text_frame.paragraphs[0].font.size = Pt(18)
lbl66.text_frame.paragraphs[0].font.bold = True

# ── Slide 67: Surface chart (2D heatmap approximation) ─────────────────────
slide67 = prs.slides.add_slide(blank)

# Create a bar chart as placeholder, then replace with surface chart XML
surf_placeholder = CategoryChartData()
surf_placeholder.categories = ['X1', 'X2', 'X3', 'X4', 'X5']
surf_placeholder.add_series('Y1', (10, 20, 30, 40, 50))

chart_frame67 = slide67.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), surf_placeholder
)
chart67 = chart_frame67.chart
chart_part67 = chart67.part
chart_elem67 = chart_part67._element
ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Find plotArea and remove existing barChart
plot_area67 = chart_elem67.find(f'.//{{{ns_c}}}plotArea')
for bar in plot_area67.findall(f'{{{ns_c}}}barChart'):
    plot_area67.remove(bar)

# Build surface chart with 4 series × 5 data points (grid)
surface_xml = f"""<c:surfaceChart xmlns:c="{ns_c}" xmlns:a="{ns_a}">
  <c:wireframe val="0"/>
  <c:ser>
    <c:idx val="0"/><c:order val="0"/>
    <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Row 1</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$6</c:f><c:strCache><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>X1</c:v></c:pt><c:pt idx="1"><c:v>X2</c:v></c:pt><c:pt idx="2"><c:v>X3</c:v></c:pt><c:pt idx="3"><c:v>X4</c:v></c:pt><c:pt idx="4"><c:v>X5</c:v></c:pt>
    </c:strCache></c:strRef></c:cat>
    <c:val><c:numRef><c:f>Sheet1!$B$2:$B$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>10</c:v></c:pt><c:pt idx="1"><c:v>25</c:v></c:pt><c:pt idx="2"><c:v>40</c:v></c:pt><c:pt idx="3"><c:v>30</c:v></c:pt><c:pt idx="4"><c:v>15</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="1"/><c:order val="1"/>
    <c:tx><c:strRef><c:f>Sheet1!$C$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Row 2</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$C$2:$C$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>20</c:v></c:pt><c:pt idx="1"><c:v>35</c:v></c:pt><c:pt idx="2"><c:v>50</c:v></c:pt><c:pt idx="3"><c:v>45</c:v></c:pt><c:pt idx="4"><c:v>30</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="2"/><c:order val="2"/>
    <c:tx><c:strRef><c:f>Sheet1!$D$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Row 3</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$D$2:$D$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>35</c:v></c:pt><c:pt idx="1"><c:v>50</c:v></c:pt><c:pt idx="2"><c:v>60</c:v></c:pt><c:pt idx="3"><c:v>55</c:v></c:pt><c:pt idx="4"><c:v>40</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:ser>
    <c:idx val="3"/><c:order val="3"/>
    <c:tx><c:strRef><c:f>Sheet1!$E$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Row 4</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:val><c:numRef><c:f>Sheet1!$E$2:$E$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>25</c:v></c:pt><c:pt idx="1"><c:v>40</c:v></c:pt><c:pt idx="2"><c:v>45</c:v></c:pt><c:pt idx="3"><c:v>50</c:v></c:pt><c:pt idx="4"><c:v>35</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:bandFmts/>
  <c:axId val="111111"/>
  <c:axId val="222222"/>
  <c:axId val="333333"/>
</c:surfaceChart>"""

# Remove existing axes
for ax_tag in ['catAx', 'valAx', 'serAx']:
    for ax in plot_area67.findall(f'{{{ns_c}}}{ax_tag}'):
        plot_area67.remove(ax)

# Insert surface chart
plot_area67.append(etree.fromstring(surface_xml))

# Add axes for surface chart
surf_axes = f"""<c:catAx xmlns:c="{ns_c}"><c:axId val="111111"/><c:scaling><c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="b"/><c:crossAx val="222222"/></c:catAx>"""
plot_area67.append(etree.fromstring(surf_axes))
surf_val_ax = f"""<c:valAx xmlns:c="{ns_c}"><c:axId val="222222"/><c:scaling><c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="l"/><c:crossAx val="111111"/><c:majorGridlines/></c:valAx>"""
plot_area67.append(etree.fromstring(surf_val_ax))
# serAx for surface chart
surf_ser_ax = f"""<c:serAx xmlns:c="{ns_c}"><c:axId val="333333"/><c:scaling><c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="b"/><c:crossAx val="222222"/></c:serAx>"""
plot_area67.append(etree.fromstring(surf_ser_ax))

# Add view3D to chart element
chart_node67 = chart_elem67.find(f'{{{ns_c}}}chart')
view3d_xml = f"""<c:view3D xmlns:c="{ns_c}"><c:rotX val="15"/><c:rotY val="20"/><c:depthPercent val="100"/><c:rAngAx val="1"/><c:perspective val="30"/></c:view3D>"""
chart_node67.insert(0, etree.fromstring(view3d_xml))

lbl67 = slide67.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl67.text_frame.paragraphs[0].text = "Slide 67: Surface chart (2D heatmap)"
lbl67.text_frame.paragraphs[0].font.size = Pt(18)
lbl67.text_frame.paragraphs[0].font.bold = True

# ── Slide 68: Pie of pie chart (ofPieChart) ─────────────────────────────────
slide68 = prs.slides.add_slide(blank)

# Create a pie chart as placeholder, then replace with ofPieChart XML
ofpie_placeholder = CategoryChartData()
ofpie_placeholder.categories = ['Product A', 'Product B', 'Product C', 'Product D', 'Product E']
ofpie_placeholder.add_series('Sales', (40, 25, 15, 12, 8))

chart_frame68 = slide68.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), ofpie_placeholder
)
chart68 = chart_frame68.chart
chart_part68 = chart68.part
chart_elem68 = chart_part68._element

# Find plotArea and remove existing pieChart
plot_area68 = chart_elem68.find(f'.//{{{ns_c}}}plotArea')
for pie in plot_area68.findall(f'{{{ns_c}}}pieChart'):
    plot_area68.remove(pie)

# Build ofPieChart element
ofpie_xml = f"""<c:ofPieChart xmlns:c="{ns_c}" xmlns:a="{ns_a}">
  <c:ofPieType val="pie"/>
  <c:varyColors val="1"/>
  <c:ser>
    <c:idx val="0"/><c:order val="0"/>
    <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Sales</c:v></c:pt></c:strCache></c:strRef></c:tx>
    <c:cat><c:strRef><c:f>Sheet1!$A$2:$A$6</c:f><c:strCache><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>Product A</c:v></c:pt><c:pt idx="1"><c:v>Product B</c:v></c:pt>
      <c:pt idx="2"><c:v>Product C</c:v></c:pt><c:pt idx="3"><c:v>Product D</c:v></c:pt>
      <c:pt idx="4"><c:v>Product E</c:v></c:pt>
    </c:strCache></c:strRef></c:cat>
    <c:val><c:numRef><c:f>Sheet1!$B$2:$B$6</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="5"/>
      <c:pt idx="0"><c:v>40</c:v></c:pt><c:pt idx="1"><c:v>25</c:v></c:pt>
      <c:pt idx="2"><c:v>15</c:v></c:pt><c:pt idx="3"><c:v>12</c:v></c:pt>
      <c:pt idx="4"><c:v>8</c:v></c:pt>
    </c:numCache></c:numRef></c:val>
  </c:ser>
  <c:gapWidth val="150"/>
  <c:splitPos val="2"/>
</c:ofPieChart>"""
plot_area68.append(etree.fromstring(ofpie_xml))

lbl68 = slide68.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl68.text_frame.paragraphs[0].text = "Slide 68: Pie of pie chart (ofPieChart)"
lbl68.text_frame.paragraphs[0].font.size = Pt(18)
lbl68.text_frame.paragraphs[0].font.bold = True

# ── Slide 69: 3D bar chart with view3D ────────────────────────────────────
slide69 = prs.slides.add_slide(blank)

chart_data69 = CategoryChartData()
chart_data69.categories = ['East', 'West', 'North', 'South']
chart_data69.add_series('Revenue', (320, 280, 190, 250))
chart_data69.add_series('Cost', (200, 180, 140, 160))

chart_frame69 = slide69.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5),
    Inches(7), Inches(4.5), chart_data69
)

# Replace barChart with bar3DChart and add view3D via lxml
chart69 = chart_frame69.chart
chart_part69 = chart69.part
chart_elem69 = chart_part69._element
plot_area69 = chart_elem69.find(f'.//{{{ns_c}}}plotArea')
for bar in plot_area69.findall(f'{{{ns_c}}}barChart'):
    # Convert barChart to bar3DChart by changing tag
    bar.tag = f'{{{ns_c}}}bar3DChart'

# Add view3D element
chart_node69 = chart_elem69.find(f'{{{ns_c}}}chart')
view3d_69 = f"""<c:view3D xmlns:c="{ns_c}"><c:rotX val="20"/><c:rotY val="30"/><c:depthPercent val="150"/><c:rAngAx val="1"/><c:perspective val="40"/></c:view3D>"""
chart_node69.insert(0, etree.fromstring(view3d_69))

lbl69 = slide69.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl69.text_frame.paragraphs[0].text = "Slide 69: 3D bar chart (view3D preserved)"
lbl69.text_frame.paragraphs[0].font.size = Pt(18)
lbl69.text_frame.paragraphs[0].font.bold = True

# ── Slide 70: Text outline (a:rPr/a:ln) ──────────────────────────────────────
slide70 = prs.slides.add_slide(blank)
lbl70 = slide70.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl70.text_frame.paragraphs[0].text = "Slide 70: Text outline (a:rPr/a:ln)"
lbl70.text_frame.paragraphs[0].font.size = Pt(18)
lbl70.text_frame.paragraphs[0].font.bold = True

ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Text box with outlined text
tb70a = slide70.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1.5))
tf70a = tb70a.text_frame
p70a = tf70a.paragraphs[0]
r70a = p70a.add_run()
r70a.text = "Red Outlined Text"
r70a.font.size = Pt(48)
r70a.font.bold = True
# Inject a:ln into a:rPr via lxml
rpr70a = r70a._r.find(f'{{{ns_a}}}rPr')
ln70a = etree.SubElement(rpr70a, f'{{{ns_a}}}ln', attrib={'w': '25400'})
sf70a = etree.SubElement(ln70a, f'{{{ns_a}}}solidFill')
etree.SubElement(sf70a, f'{{{ns_a}}}srgbClr', attrib={'val': 'FF0000'})

# Thinner outline with different color
tb70b = slide70.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(8), Inches(1.5))
tf70b = tb70b.text_frame
p70b = tf70b.paragraphs[0]
r70b = p70b.add_run()
r70b.text = "Blue Outlined (thin)"
r70b.font.size = Pt(36)
r70b.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
rpr70b = r70b._r.find(f'{{{ns_a}}}rPr')
ln70b = etree.SubElement(rpr70b, f'{{{ns_a}}}ln', attrib={'w': '12700'})
sf70b = etree.SubElement(ln70b, f'{{{ns_a}}}solidFill')
etree.SubElement(sf70b, f'{{{ns_a}}}srgbClr', attrib={'val': '0000FF'})

# ── Slide 71: Text gradient fill (a:rPr/a:gradFill) ─────────────────────────
slide71 = prs.slides.add_slide(blank)
lbl71 = slide71.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl71.text_frame.paragraphs[0].text = "Slide 71: Text gradient fill (a:rPr/a:gradFill)"
lbl71.text_frame.paragraphs[0].font.size = Pt(18)
lbl71.text_frame.paragraphs[0].font.bold = True

# Text with gradient fill
tb71a = slide71.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(2))
tf71a = tb71a.text_frame
p71a = tf71a.paragraphs[0]
r71a = p71a.add_run()
r71a.text = "Gradient Text"
r71a.font.size = Pt(60)
r71a.font.bold = True
# Inject a:gradFill into a:rPr
rpr71a = r71a._r.find(f'{{{ns_a}}}rPr')
grad_xml71 = f"""<a:gradFill xmlns:a="{ns_a}">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="FF0000"/></a:gs>
    <a:gs pos="50000"><a:srgbClr val="FFFF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="1"/>
</a:gradFill>"""
rpr71a.append(etree.fromstring(grad_xml71))

# Second text with solid fill + gradient (gradient takes priority in display)
tb71b = slide71.shapes.add_textbox(Inches(0.5), Inches(4), Inches(8), Inches(2))
tf71b = tb71b.text_frame
p71b = tf71b.paragraphs[0]
r71b = p71b.add_run()
r71b.text = "Green-Blue Grad"
r71b.font.size = Pt(48)
rpr71b = r71b._r.find(f'{{{ns_a}}}rPr')
grad_xml71b = f"""<a:gradFill xmlns:a="{ns_a}">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="00FF00"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="0000FF"/></a:gs>
  </a:gsLst>
  <a:lin ang="5400000" scaled="1"/>
</a:gradFill>"""
rpr71b.append(etree.fromstring(grad_xml71b))

# ── Slide 72: Text warp (a:prstTxWarp) ──────────────────────────────────────
slide72 = prs.slides.add_slide(blank)
lbl72 = slide72.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl72.text_frame.paragraphs[0].text = "Slide 72: Text warp (a:prstTxWarp)"
lbl72.text_frame.paragraphs[0].font.size = Pt(18)
lbl72.text_frame.paragraphs[0].font.bold = True

ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# Text box with wave warp
tb72a = slide72.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(2))
tf72a = tb72a.text_frame
p72a = tf72a.paragraphs[0]
r72a = p72a.add_run()
r72a.text = "Wave Text"
r72a.font.size = Pt(36)
r72a.font.bold = True
# Inject prstTxWarp into bodyPr
sp72a = tb72a._element
body_pr72a = sp72a.find(f'.//{{{ns_a}}}bodyPr')
warp72a = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textWave1">
  <a:avLst>
    <a:gd name="adj" fmla="val 19773"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72a.insert(0, etree.fromstring(warp72a))

# Text box with arch up warp
tb72b = slide72.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(2))
tf72b = tb72b.text_frame
p72b = tf72b.paragraphs[0]
r72b = p72b.add_run()
r72b.text = "Arch Up"
r72b.font.size = Pt(36)
r72b.font.bold = True
sp72b = tb72b._element
body_pr72b = sp72b.find(f'.//{{{ns_a}}}bodyPr')
warp72b = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textArchUp">
  <a:avLst>
    <a:gd name="adj" fmla="val 10800000"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72b.insert(0, etree.fromstring(warp72b))

# Text box with no adjust values (textDeflate)
tb72c = slide72.shapes.add_textbox(Inches(0.5), Inches(4), Inches(8), Inches(2))
tf72c = tb72c.text_frame
p72c = tf72c.paragraphs[0]
r72c = p72c.add_run()
r72c.text = "Deflate Text"
r72c.font.size = Pt(36)
sp72c = tb72c._element
body_pr72c = sp72c.find(f'.//{{{ns_a}}}bodyPr')
warp72c = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textDeflate">
  <a:avLst/>
</a:prstTxWarp>"""
body_pr72c.insert(0, etree.fromstring(warp72c))

# ── Slide 72b: Additional text warp presets ──────────────────────────────────
slide72b = prs.slides.add_slide(blank)
lbl72b = slide72b.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl72b.text_frame.paragraphs[0].text = "Slide 72b: Additional text warp presets"
lbl72b.text_frame.paragraphs[0].font.size = Pt(18)
lbl72b.text_frame.paragraphs[0].font.bold = True

# textSlantUp
tb72d = slide72b.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(1.5))
tf72d = tb72d.text_frame
r72d = tf72d.paragraphs[0].add_run()
r72d.text = "Slant Up"
r72d.font.size = Pt(36)
r72d.font.bold = True
sp72d = tb72d._element
body_pr72d = sp72d.find(f'.//{{{ns_a}}}bodyPr')
warp72d = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textSlantUp">
  <a:avLst>
    <a:gd name="adj" fmla="val 55556"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72d.insert(0, etree.fromstring(warp72d))

# textCurveUp
tb72e = slide72b.shapes.add_textbox(Inches(5), Inches(1), Inches(4), Inches(1.5))
tf72e = tb72e.text_frame
r72e = tf72e.paragraphs[0].add_run()
r72e.text = "Curve Up"
r72e.font.size = Pt(36)
r72e.font.bold = True
r72e.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
sp72e = tb72e._element
body_pr72e = sp72e.find(f'.//{{{ns_a}}}bodyPr')
warp72e = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textCurveUp">
  <a:avLst>
    <a:gd name="adj" fmla="val 45977"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72e.insert(0, etree.fromstring(warp72e))

# textInflate
tb72f = slide72b.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4), Inches(1.5))
tf72f = tb72f.text_frame
r72f = tf72f.paragraphs[0].add_run()
r72f.text = "Inflate"
r72f.font.size = Pt(36)
r72f.font.bold = True
r72f.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
sp72f = tb72f._element
body_pr72f = sp72f.find(f'.//{{{ns_a}}}bodyPr')
warp72f = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textInflate">
  <a:avLst/>
</a:prstTxWarp>"""
body_pr72f.insert(0, etree.fromstring(warp72f))

# textChevron
tb72g = slide72b.shapes.add_textbox(Inches(5), Inches(3), Inches(4), Inches(1.5))
tf72g = tb72g.text_frame
r72g = tf72g.paragraphs[0].add_run()
r72g.text = "Chevron"
r72g.font.size = Pt(36)
r72g.font.bold = True
r72g.font.color.rgb = RGBColor(0x00, 0xB0, 0x50)
sp72g = tb72g._element
body_pr72g = sp72g.find(f'.//{{{ns_a}}}bodyPr')
warp72g = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textChevron">
  <a:avLst>
    <a:gd name="adj" fmla="val 25000"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72g.insert(0, etree.fromstring(warp72g))

# textCircle
tb72h = slide72b.shapes.add_textbox(Inches(2.5), Inches(5), Inches(4), Inches(2.5))
tf72h = tb72h.text_frame
r72h = tf72h.paragraphs[0].add_run()
r72h.text = "Circle Text Example"
r72h.font.size = Pt(28)
r72h.font.bold = True
sp72h = tb72h._element
body_pr72h = sp72h.find(f'.//{{{ns_a}}}bodyPr')
warp72h = f"""<a:prstTxWarp xmlns:a="{ns_a}" prst="textCircle">
  <a:avLst>
    <a:gd name="adj" fmla="val 10800000"/>
  </a:avLst>
</a:prstTxWarp>"""
body_pr72h.insert(0, etree.fromstring(warp72h))

# ── Slide 73: Stacked / Percent-stacked bar charts ──────────────────────────

slide73 = prs.slides.add_slide(blank)

# Percent-stacked horizontal bar (like gender breakdown: 58.1% / 41.9%)
chart_data73a = CategoryChartData()
chart_data73a.categories = ['Gender']
chart_data73a.add_series('Male', (58.1,))
chart_data73a.add_series('Female', (41.9,))

chart_frame73a = slide73.shapes.add_chart(
    XL_CHART_TYPE.BAR_STACKED_100, Inches(0.5), Inches(1.5),
    Inches(4), Inches(3), chart_data73a
)

# Stacked vertical column chart (multiple categories)
chart_data73b = CategoryChartData()
chart_data73b.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data73b.add_series('Product A', (100, 120, 90, 150))
chart_data73b.add_series('Product B', (80, 70, 110, 60))
chart_data73b.add_series('Product C', (50, 40, 30, 80))

chart_frame73b = slide73.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED, Inches(5), Inches(1.5),
    Inches(4), Inches(3), chart_data73b
)

# Percent-stacked vertical column chart
chart_data73c = CategoryChartData()
chart_data73c.categories = ['East', 'West', 'North']
chart_data73c.add_series('2024', (300, 200, 150))
chart_data73c.add_series('2025', (250, 350, 200))

chart_frame73c = slide73.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED_100, Inches(2.5), Inches(5),
    Inches(4), Inches(2), chart_data73c
)

lbl73 = slide73.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl73.text_frame.paragraphs[0].text = "Slide 73: Stacked / Percent-stacked bar charts"
lbl73.text_frame.paragraphs[0].font.size = Pt(18)
lbl73.text_frame.paragraphs[0].font.bold = True

# ── Slide 74: Speaker notes + comments ──────────────────────────────────────

slide74 = prs.slides.add_slide(blank)

lbl74 = slide74.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
lbl74.text_frame.paragraphs[0].text = "Slide 74: Speaker notes + comments"
lbl74.text_frame.paragraphs[0].font.size = Pt(18)
lbl74.text_frame.paragraphs[0].font.bold = True

body74 = slide74.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(3))
body74.text_frame.paragraphs[0].text = (
    "This slide has speaker notes and comments attached.\n"
    "Notes and comments are metadata not rendered on slide, "
    "but preserved in round-trip export."
)
body74.text_frame.paragraphs[0].font.size = Pt(14)
body74.text_frame.word_wrap = True

# Add speaker notes
notes_slide74 = slide74.notes_slide
notes_tf = notes_slide74.notes_text_frame
notes_tf.text = ""
p1 = notes_tf.paragraphs[0]
p1.text = "These are the speaker notes for slide 74."
p2 = notes_tf.add_paragraph()
p2.text = "Second paragraph of notes with key points."
p3 = notes_tf.add_paragraph()
p3.text = "Remember to mention the round-trip preservation."

# Add comments via raw XML (python-pptx has no public comment API)
from pptx.opc.package import Part as OpcPart
from pptx.opc.packuri import PackURI

authors_xml = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:cmAuthorLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
    '<p:cmAuthor id="1" name="Test User" initials="TU" lastIdx="2" clrIdx="0"/>'
    '</p:cmAuthorLst>'
)

# Add commentAuthors part
authors_part = OpcPart(
    PackURI('/ppt/commentAuthors.xml'),
    'application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml',
    prs.part.package,
    authors_xml.encode('utf-8'),
)
prs.part.relate_to(authors_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors')

# Add comments part for slide 74
comment_xml = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<p:cmLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
    '<p:cm authorId="1" dt="2025-01-15T10:30:00.000" idx="1">'
    '<p:pos x="100" y="200"/>'
    '<p:text>This is a test comment on slide 74.</p:text>'
    '</p:cm>'
    '<p:cm authorId="1" dt="2025-01-15T11:00:00.000" idx="2">'
    '<p:pos x="300" y="400"/>'
    '<p:text>Second comment with review feedback.</p:text>'
    '</p:cm>'
    '</p:cmLst>'
)

slide74_idx = len(prs.slides)
comment_part = OpcPart(
    PackURI(f'/ppt/comments/comment{slide74_idx}.xml'),
    'application/vnd.openxmlformats-officedocument.presentationml.comments+xml',
    prs.part.package,
    comment_xml.encode('utf-8'),
)
slide74.part.relate_to(comment_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments')

# ── Slide 75: SmartArt fallback (mc:AlternateContent) ────────────────────────

slide75 = prs.slides.add_slide(blank)

# We inject raw XML with mc:AlternateContent containing:
#   mc:Choice  → dummy p:graphicFrame referencing diagram (preserved for round-trip)
#   mc:Fallback → p:grpSp with 3 rectangles representing a simple process flow
slide75_sp_tree = slide75.shapes._spTree
# Register mc: namespace prefix on the root <p:sld> element so lxml uses mc: prefix
sld_root = slide75_sp_tree.getparent()  # p:cSld -> p:sld
if sld_root is not None and sld_root.getparent() is not None:
    sld_root = sld_root.getparent()
MC_NS = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
DGM_NS = 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
etree.register_namespace('mc', MC_NS)
etree.register_namespace('dgm', DGM_NS)

mc_ac = etree.SubElement(slide75_sp_tree, f'{{{MC_NS}}}AlternateContent')

# mc:Choice — dummy SmartArt reference (for round-trip preservation)
mc_choice = etree.SubElement(mc_ac, f'{{{MC_NS}}}Choice', attrib={'Requires': 'dgm'})
choice_gf = etree.SubElement(mc_choice,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}graphicFrame')
choice_nvgf = etree.SubElement(choice_gf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvGraphicFramePr')
choice_cnvpr = etree.SubElement(choice_nvgf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr',
    attrib={'id': '100', 'name': 'SmartArt Diagram'})
choice_cnvgf = etree.SubElement(choice_nvgf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvGraphicFramePr')
choice_nvpr = etree.SubElement(choice_nvgf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
choice_xfrm = etree.SubElement(choice_gf,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}xfrm')
etree.SubElement(choice_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}off',
    attrib={'x': '457200', 'y': '1371600'})
etree.SubElement(choice_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}ext',
    attrib={'cx': '8229600', 'y': '4572000'})
# Dummy graphic with dgm:relIds
choice_graphic = etree.SubElement(choice_gf,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}graphic')
choice_gd = etree.SubElement(choice_graphic,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData',
    attrib={'uri': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'})
etree.SubElement(choice_gd,
    f'{{{DGM_NS}}}relIds',
    attrib={
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}dm': 'rId10',
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}lo': 'rId11',
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}qs': 'rId12',
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}cs': 'rId13',
    })

# mc:Fallback — pre-rendered shapes (3 rectangles as a simple process flow)
mc_fallback = etree.SubElement(mc_ac,
    f'{{{MC_NS}}}Fallback')

# Create a group shape containing 3 process boxes
grp_sp = etree.SubElement(mc_fallback,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}grpSp')
grp_nv = etree.SubElement(grp_sp,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvGrpSpPr')
etree.SubElement(grp_nv,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr',
    attrib={'id': '101', 'name': 'SmartArt Fallback Group'})
etree.SubElement(grp_nv,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvGrpSpPr')
etree.SubElement(grp_nv,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')

grp_sp_pr = etree.SubElement(grp_sp,
    '{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr')
grp_xfrm = etree.SubElement(grp_sp_pr,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
etree.SubElement(grp_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}off',
    attrib={'x': '457200', 'y': '1371600'})
etree.SubElement(grp_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}ext',
    attrib={'cx': '8229600', 'cy': '2743200'})
etree.SubElement(grp_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}chOff',
    attrib={'x': '0', 'y': '0'})
etree.SubElement(grp_xfrm,
    '{http://schemas.openxmlformats.org/drawingml/2006/main}chExt',
    attrib={'cx': '8229600', 'cy': '2743200'})

# Helper to add a rectangle shape inside the group
def add_smartart_box(parent, sp_id, name, x, y, cx, cy, text, fill_hex):
    sp = etree.SubElement(parent,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sp')
    nv = etree.SubElement(sp,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
    etree.SubElement(nv,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr',
        attrib={'id': str(sp_id), 'name': name})
    etree.SubElement(nv,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvSpPr')
    etree.SubElement(nv,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
    spPr = etree.SubElement(sp,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
    xfrm = etree.SubElement(spPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    etree.SubElement(xfrm,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}off',
        attrib={'x': str(x), 'y': str(y)})
    etree.SubElement(xfrm,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}ext',
        attrib={'cx': str(cx), 'cy': str(cy)})
    etree.SubElement(spPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom',
        attrib={'prst': 'roundRect'}).append(
        etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst'))
    solid = etree.SubElement(spPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    etree.SubElement(solid,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
        attrib={'val': fill_hex})
    # Text body
    txBody = etree.SubElement(sp,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}txBody')
    bodyPr = etree.SubElement(txBody,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr',
        attrib={'anchor': 'ctr', 'wrap': 'square'})
    etree.SubElement(txBody,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
    p = etree.SubElement(txBody,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
    pPr = etree.SubElement(p,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr',
        attrib={'algn': 'ctr'})
    r = etree.SubElement(p,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
    rPr = etree.SubElement(r,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr',
        attrib={'lang': 'en-US', 'sz': '1800', 'b': '1'})
    solid_font = etree.SubElement(rPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    etree.SubElement(solid_font,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
        attrib={'val': 'FFFFFF'})
    t = etree.SubElement(r,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
    t.text = text

# Three process boxes: Plan → Build → Ship
add_smartart_box(grp_sp, 102, 'Step 1', 0, 457200, 2286000, 1828800, 'Plan', '4472C4')
add_smartart_box(grp_sp, 103, 'Step 2', 2971800, 457200, 2286000, 1828800, 'Build', 'ED7D31')
add_smartart_box(grp_sp, 104, 'Step 3', 5943600, 457200, 2286000, 1828800, 'Ship', '70AD47')

# Also add a title textbox
title75 = slide75.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title75.text_frame.paragraphs[0].text = "Slide 75: SmartArt Fallback (mc:AlternateContent)"
title75.text_frame.paragraphs[0].font.size = Pt(24)
title75.text_frame.paragraphs[0].font.bold = True

# ── Slide 76: OLE embedded object (p:oleObj with fallback image) ──────────────

slide76 = prs.slides.add_slide(blank)

# Create a small 2x2 red PNG as the OLE fallback image
def make_tiny_png(r, g, b, w=2, h=2):
    """Create a minimal RGBA PNG."""
    raw = b''
    for _ in range(h):
        raw += b'\x00'  # filter: none
        for _ in range(w):
            raw += struct.pack('BBBB', r, g, b, 255)
    compressed = zlib.compress(raw)
    def chunk(ctype, data):
        c = ctype + data
        return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
    ihdr = struct.pack('>IIBBBBB', w, h, 8, 6, 0, 0, 0)  # 8-bit RGBA
    return b'\x89PNG\r\n\x1a\n' + chunk(b'IHDR', ihdr) + chunk(b'IDAT', compressed) + chunk(b'IEND', b'')

ole_fallback_png = make_tiny_png(200, 50, 50)

# Add the PNG as an image part via the slide's relationships
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part as OpcPart
from pptx.opc.packuri import PackURI

slide76_idx = len(prs.slides)
ole_img_partname = PackURI(f'/ppt/media/oleImage{slide76_idx}.png')
ole_img_part = OpcPart(
    ole_img_partname,
    'image/png',
    prs.part.package,
    ole_fallback_png,
)
ole_img_rel = slide76.part.relate_to(ole_img_part, RT.IMAGE)

# Also create a dummy OLE binary part (empty, just for structure)
ole_bin_partname = PackURI(f'/ppt/embeddings/oleObject{slide76_idx}.bin')
ole_bin_part = OpcPart(
    ole_bin_partname,
    'application/vnd.openxmlformats-officedocument.oleObject',
    prs.part.package,
    b'\x00' * 16,  # Minimal dummy data
)
ole_bin_rel = slide76.part.relate_to(
    ole_bin_part,
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject',
)

# Inject p:graphicFrame with p:oleObj into spTree
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

slide76_sp_tree = slide76.shapes._spTree
gf = etree.SubElement(slide76_sp_tree, f'{{{P_NS}}}graphicFrame')

# nvGraphicFramePr
nv = etree.SubElement(gf, f'{{{P_NS}}}nvGraphicFramePr')
etree.SubElement(nv, f'{{{P_NS}}}cNvPr', attrib={'id': '200', 'name': 'OLE Object'})
etree.SubElement(nv, f'{{{P_NS}}}cNvGraphicFramePr')
etree.SubElement(nv, f'{{{P_NS}}}nvPr')

# xfrm
xfrm = etree.SubElement(gf, f'{{{P_NS}}}xfrm')
etree.SubElement(xfrm, f'{{{A_NS}}}off', attrib={'x': '914400', 'y': '1828800'})
etree.SubElement(xfrm, f'{{{A_NS}}}ext', attrib={'cx': '4572000', 'cy': '3429000'})

# a:graphic > a:graphicData (OLE URI)
graphic = etree.SubElement(gf, f'{{{A_NS}}}graphic')
gdata = etree.SubElement(graphic, f'{{{A_NS}}}graphicData',
    attrib={'uri': 'http://schemas.openxmlformats.org/presentationml/2006/ole'})

# p:oleObj with r:id and fallback p:pic
ole_obj = etree.SubElement(gdata, f'{{{P_NS}}}oleObj', attrib={
    f'{{{R_NS}}}id': ole_bin_rel,
    'imgW': '4572000',
    'imgH': '3429000',
    'progId': 'Excel.Sheet.12',
    'name': 'Embedded Spreadsheet',
})
etree.SubElement(ole_obj, f'{{{P_NS}}}embed')

# p:pic inside oleObj (fallback image)
ole_pic = etree.SubElement(ole_obj, f'{{{P_NS}}}pic')
ole_pic_nv = etree.SubElement(ole_pic, f'{{{P_NS}}}nvPicPr')
etree.SubElement(ole_pic_nv, f'{{{P_NS}}}cNvPr', attrib={'id': '201', 'name': 'OLE Fallback'})
etree.SubElement(ole_pic_nv, f'{{{P_NS}}}cNvPicPr')
etree.SubElement(ole_pic_nv, f'{{{P_NS}}}nvPr')
ole_pic_bf = etree.SubElement(ole_pic, f'{{{P_NS}}}blipFill')
etree.SubElement(ole_pic_bf, f'{{{A_NS}}}blip', attrib={f'{{{R_NS}}}embed': ole_img_rel})
etree.SubElement(ole_pic_bf, f'{{{A_NS}}}stretch').append(
    etree.Element(f'{{{A_NS}}}fillRect'))
ole_pic_sp = etree.SubElement(ole_pic, f'{{{P_NS}}}spPr')
ole_pic_xfrm = etree.SubElement(ole_pic_sp, f'{{{A_NS}}}xfrm')
etree.SubElement(ole_pic_xfrm, f'{{{A_NS}}}off', attrib={'x': '0', 'y': '0'})
etree.SubElement(ole_pic_xfrm, f'{{{A_NS}}}ext', attrib={'cx': '4572000', 'cy': '3429000'})
etree.SubElement(ole_pic_sp, f'{{{A_NS}}}prstGeom', attrib={'prst': 'rect'}).append(
    etree.Element(f'{{{A_NS}}}avLst'))

# Title
title76 = slide76.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title76.text_frame.paragraphs[0].text = "Slide 76: OLE Embedded Object (p:oleObj)"
title76.text_frame.paragraphs[0].font.size = Pt(24)
title76.text_frame.paragraphs[0].font.bold = True

# ── Slide 77: Media (video with poster frame) ────────────────────────────────

slide77 = prs.slides.add_slide(blank)

# Create a poster frame image (blue rectangle) for the video
poster_png = make_tiny_png(30, 100, 200, w=4, h=3)

slide77_idx = len(prs.slides)
poster_partname = PackURI(f'/ppt/media/posterFrame{slide77_idx}.png')
poster_part = OpcPart(
    poster_partname,
    'image/png',
    prs.part.package,
    poster_png,
)
poster_rel = slide77.part.relate_to(poster_part, RT.IMAGE)

# Create a dummy video part
video_partname = PackURI(f'/ppt/media/video{slide77_idx}.mp4')
video_part = OpcPart(
    video_partname,
    'video/mp4',
    prs.part.package,
    b'\x00' * 32,  # Minimal dummy data
)
video_rel = slide77.part.relate_to(
    video_part,
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/video',
)

# Inject p:pic with a:videoFile into spTree
slide77_sp_tree = slide77.shapes._spTree

vid_pic = etree.SubElement(slide77_sp_tree, f'{{{P_NS}}}pic')

# nvPicPr
vid_nv = etree.SubElement(vid_pic, f'{{{P_NS}}}nvPicPr')
etree.SubElement(vid_nv, f'{{{P_NS}}}cNvPr', attrib={'id': '300', 'name': 'Video Placeholder'})
etree.SubElement(vid_nv, f'{{{P_NS}}}cNvPicPr')
vid_nvpr = etree.SubElement(vid_nv, f'{{{P_NS}}}nvPr')
etree.SubElement(vid_nvpr, f'{{{A_NS}}}videoFile', attrib={f'{{{R_NS}}}link': video_rel})

# blipFill (poster frame)
vid_bf = etree.SubElement(vid_pic, f'{{{P_NS}}}blipFill')
etree.SubElement(vid_bf, f'{{{A_NS}}}blip', attrib={f'{{{R_NS}}}embed': poster_rel})
etree.SubElement(vid_bf, f'{{{A_NS}}}stretch').append(
    etree.Element(f'{{{A_NS}}}fillRect'))

# spPr
vid_sp = etree.SubElement(vid_pic, f'{{{P_NS}}}spPr')
vid_xfrm = etree.SubElement(vid_sp, f'{{{A_NS}}}xfrm')
etree.SubElement(vid_xfrm, f'{{{A_NS}}}off', attrib={'x': '1371600', 'y': '1828800'})
etree.SubElement(vid_xfrm, f'{{{A_NS}}}ext', attrib={'cx': '6858000', 'cy': '3886200'})
etree.SubElement(vid_sp, f'{{{A_NS}}}prstGeom', attrib={'prst': 'rect'}).append(
    etree.Element(f'{{{A_NS}}}avLst'))

# Title
title77 = slide77.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title77.text_frame.paragraphs[0].text = "Slide 77: Media (Video with Poster Frame)"
title77.text_frame.paragraphs[0].font.size = Pt(24)
title77.text_frame.paragraphs[0].font.bold = True

# ── Slide 78: Math equation (OMML) ────────────────────────────────────────────

slide78 = prs.slides.add_slide(blank)

# Title
title78 = slide78.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
title78.text_frame.paragraphs[0].text = "Slide 78: Math Equation (OMML)"
title78.text_frame.paragraphs[0].font.size = Pt(28)
title78.text_frame.paragraphs[0].font.bold = True

# Add a textbox, then inject OMML XML directly into the slide XML
math_tb = slide78.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
math_tb.text_frame.paragraphs[0].text = "MATH_PLACEHOLDER"
math_tb.text_frame.paragraphs[0].font.size = Pt(24)

# ── Slide 79: Transition + Timing (round-trip preservation) ───────────────────

slide79 = prs.slides.add_slide(blank)

title79 = slide79.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title79.text_frame.paragraphs[0].text = "Slide 79: Transition + Timing (round-trip)"
title79.text_frame.paragraphs[0].font.size = Pt(24)
title79.text_frame.paragraphs[0].font.bold = True

desc79 = slide79.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
desc79.text_frame.paragraphs[0].text = "This slide has a fade transition and timing data for round-trip."
desc79.text_frame.paragraphs[0].font.size = Pt(16)

# ── Slide 80: Hidden slide ───────────────────────────────────────────────────

slide80 = prs.slides.add_slide(blank)

title80 = slide80.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title80.text_frame.paragraphs[0].text = "Slide 80: Hidden Slide"
title80.text_frame.paragraphs[0].font.size = Pt(24)
title80.text_frame.paragraphs[0].font.bold = True

desc80 = slide80.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
desc80.text_frame.paragraphs[0].text = "This slide is marked as hidden (show='0')."
desc80.text_frame.paragraphs[0].font.size = Pt(16)

# ── Slide 81: WMF → SVG conversion ───────────────────────────────────────────

slide81 = prs.slides.add_slide(blank)

title81 = slide81.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title81.text_frame.paragraphs[0].text = "Slide 81: WMF Image (WMF → SVG conversion)"
title81.text_frame.paragraphs[0].font.size = Pt(24)
title81.text_frame.paragraphs[0].font.bold = True

desc81 = slide81.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
desc81.text_frame.paragraphs[0].text = "WMF images below are converted to SVG at runtime."
desc81.text_frame.paragraphs[0].font.size = Pt(14)

# Placeholder — will be replaced with p:pic referencing a WMF via ZIP patching
wmf_placeholder = slide81.shapes.add_textbox(Inches(1), Inches(2.5), Inches(3), Inches(3))
wmf_placeholder.text_frame.paragraphs[0].text = "WMF_PLACEHOLDER"
wmf_placeholder.text_frame.paragraphs[0].font.size = Pt(12)

# Save first, then patch the OMML into the slide XML
output_path = 'test_fixtures/test_features.pptx'
prs.save(output_path)

# Patch slide78 XML to inject actual OMML
OMML_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
etree.register_namespace('m', OMML_NS)

total_slides = len(prs.slides)  # 82
# python-pptx numbers slide files sequentially
slide78_path = 'ppt/slides/slide79.xml'

# Read entire ZIP into memory, then close it before writing
zin = zipfile.ZipFile(output_path, 'r')
slide78_xml = zin.read(slide78_path).decode('utf-8')
all_entries = {}
for item in zin.infolist():
    all_entries[item.filename] = (item, zin.read(item.filename))
zin.close()

# Replace the placeholder paragraph with OMML
omml_xml = (
    '<m:oMathPara xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">'
    '<m:oMath>'
    '<m:r><m:t>x</m:t></m:r>'
    '<m:r><m:t>=</m:t></m:r>'
    '<m:f>'  # fraction
    '<m:num><m:r><m:t>-b±</m:t></m:r>'
    '<m:rad><m:radPr><m:degHide m:val="1"/></m:radPr>'
    '<m:deg/><m:e><m:sSup><m:e><m:r><m:t>b</m:t></m:r></m:e>'
    '<m:sup><m:r><m:t>2</m:t></m:r></m:sup></m:sSup>'
    '<m:r><m:t>-4ac</m:t></m:r></m:e></m:rad></m:num>'
    '<m:den><m:r><m:t>2a</m:t></m:r></m:den>'
    '</m:f>'
    '</m:oMath>'
    '</m:oMathPara>'
)
slide78_xml = re.sub(
    r'<a:r>(?:<a:rPr[^/]*/>)?<a:t>MATH_PLACEHOLDER</a:t></a:r>',
    omml_xml,
    slide78_xml
)
# Add m: namespace to root if not present
if 'xmlns:m=' not in slide78_xml:
    slide78_xml = slide78_xml.replace(
        'xmlns:a=',
        f'xmlns:m="{OMML_NS}" xmlns:a='
    )

# ── Patch slide 79: inject transition + timing XML ──
slide79_path = 'ppt/slides/slide80.xml'
slide79_xml = all_entries[slide79_path][1].decode('utf-8') if isinstance(all_entries[slide79_path][1], bytes) else all_entries[slide79_path][1]
# Inject p:transition and p:timing before </p:sld>
transition_xml = '<p:transition spd="med"><p:fade/></p:transition>'
timing_xml = (
    '<p:timing>'
    '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"/></p:par></p:tnLst>'
    '</p:timing>'
)
slide79_xml = slide79_xml.replace('</p:sld>', transition_xml + timing_xml + '</p:sld>')

# ── Patch slide 80: inject show="0" for hidden slide ──
slide80_path = 'ppt/slides/slide81.xml'
slide80_xml = all_entries[slide80_path][1].decode('utf-8') if isinstance(all_entries[slide80_path][1], bytes) else all_entries[slide80_path][1]
# Add show="0" to p:sld element
slide80_xml = slide80_xml.replace('<p:sld ', '<p:sld show="0" ', 1)

# ── Patch slide 81: inject WMF image + picture shape ──
slide81_path = 'ppt/slides/slide82.xml'
slide81_xml = all_entries[slide81_path][1].decode('utf-8') if isinstance(all_entries[slide81_path][1], bytes) else all_entries[slide81_path][1]

# Build a Placeable WMF binary with colored shapes (rectangle + ellipse + polygon)
def build_wmf_binary():
    """Generate a Placeable WMF with a blue rectangle, red ellipse, and green triangle."""
    records = bytearray()

    def add_record(func_num, params=b''):
        # WMF record: size (uint32 words) + function (uint16) + params
        param_words = len(params) // 2
        size = 3 + param_words  # 3 words for header
        records.extend(struct.pack('<IH', size, func_num) + params)

    # SetMapMode (MM_ANISOTROPIC = 8)
    add_record(0x0103, struct.pack('<h', 8))
    # SetWindowOrg (0, 0)
    add_record(0x020B, struct.pack('<hh', 0, 0))
    # SetWindowExt (1000, 800)
    add_record(0x020C, struct.pack('<hh', 800, 1000))

    # ── Blue filled rectangle ──
    # CreateBrushIndirect: style=0(solid), colorR=0x0000FF, hatch=0
    add_record(0x02FC, struct.pack('<h', 0) + struct.pack('<BBBx', 0, 0, 255) + struct.pack('<h', 0))
    # SelectObject (index 0)
    add_record(0x012D, struct.pack('<h', 0))
    # Rectangle (bottom=400, right=450, top=50, left=50) — Y,X order
    add_record(0x041B, struct.pack('<hhhh', 400, 450, 50, 50))

    # ── Red filled ellipse ──
    # CreateBrushIndirect: style=0(solid), color=0xFF0000, hatch=0
    add_record(0x02FC, struct.pack('<h', 0) + struct.pack('<BBBx', 255, 0, 0) + struct.pack('<h', 0))
    # SelectObject (index 1)
    add_record(0x012D, struct.pack('<h', 1))
    # Ellipse (bottom=400, right=950, top=50, left=550)
    add_record(0x0418, struct.pack('<hhhh', 400, 950, 50, 550))

    # ── Green filled triangle (Polygon) ──
    # CreateBrushIndirect: style=0(solid), color=0x00FF00, hatch=0
    add_record(0x02FC, struct.pack('<h', 0) + struct.pack('<BBBx', 0, 255, 0) + struct.pack('<h', 0))
    # SelectObject (index 2)
    add_record(0x012D, struct.pack('<h', 2))
    # Polygon: 3 points — (500,750), (250,500), (750,500) — X,Y order in params
    add_record(0x0324, struct.pack('<h', 3) + struct.pack('<hh', 500, 750) + struct.pack('<hh', 250, 500) + struct.pack('<hh', 750, 500))

    # DeleteObject (0, 1, 2)
    add_record(0x01F0, struct.pack('<h', 0))
    add_record(0x01F0, struct.pack('<h', 1))
    add_record(0x01F0, struct.pack('<h', 2))

    # EOF record
    records.extend(struct.pack('<IH', 3, 0x0000))

    # Build standard WMF header (18 bytes)
    file_type = 1  # memory metafile
    header_size = 9  # words
    version = 0x0300
    file_size = (18 + len(records)) // 2  # in words
    num_objects = 3
    max_record = max(struct.unpack_from('<I', records, i)[0] for i in range(0, len(records), 2) if i + 4 <= len(records) and i % 2 == 0)  # approximate
    # Recalculate max_record properly
    max_rec_size = 0
    pos = 0
    while pos < len(records):
        rec_size = struct.unpack_from('<I', records, pos)[0]
        if rec_size > max_rec_size:
            max_rec_size = rec_size
        pos += rec_size * 2
    std_header = struct.pack('<HHIHIHH', file_type, header_size, version, file_size, num_objects, max_rec_size, 0)

    # Placeable WMF header (22 bytes)
    magic = 0x9AC6CDD7
    hmf = 0
    bbox_left, bbox_top, bbox_right, bbox_bottom = 0, 0, 1000, 800
    inch = 96  # units per inch
    reserved = 0
    placeable = struct.pack('<IHhhhhHI', magic, hmf, bbox_left, bbox_top, bbox_right, bbox_bottom, inch, reserved)
    # Checksum: XOR of first 10 uint16 values
    chk = 0
    for i in range(0, 20, 2):
        chk ^= struct.unpack_from('<H', placeable, i)[0]
    placeable += struct.pack('<H', chk & 0xFFFF)

    return bytes(placeable) + bytes(std_header) + bytes(records)

wmf_data = build_wmf_binary()

# Replace WMF_PLACEHOLDER textbox with a p:pic shape
# Find and remove the placeholder <p:sp> containing "WMF_PLACEHOLDER"
slide81_xml = re.sub(
    r'<p:sp>.*?WMF_PLACEHOLDER.*?</p:sp>',
    '',
    slide81_xml,
    flags=re.DOTALL
)

# Insert p:pic shape referencing the WMF image before </p:spTree>
pic_xml = (
    '<p:pic>'
    '<p:nvPicPr>'
    '<p:cNvPr id="100" name="WMF Picture"/>'
    '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
    '<p:nvPr/>'
    '</p:nvPicPr>'
    '<p:blipFill>'
    '<a:blip r:embed="rWmf1"/>'
    '<a:stretch><a:fillRect/></a:stretch>'
    '</p:blipFill>'
    '<p:spPr>'
    '<a:xfrm>'
    '<a:off x="914400" y="2286000"/>'
    '<a:ext cx="4572000" cy="3657600"/>'
    '</a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
    '</p:spPr>'
    '</p:pic>'
)
slide81_xml = slide81_xml.replace('</p:spTree>', pic_xml + '</p:spTree>')

# Add relationship for the WMF image in slide81.xml.rels
slide81_rels_path = 'ppt/slides/_rels/slide82.xml.rels'
if slide81_rels_path in all_entries:
    slide81_rels = all_entries[slide81_rels_path][1].decode('utf-8')
else:
    slide81_rels = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"></Relationships>'
wmf_rel = '<Relationship Id="rWmf1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/image_wmf1.wmf"/>'
slide81_rels = slide81_rels.replace('</Relationships>', wmf_rel + '</Relationships>')

# Add [Content_Types].xml entry for .wmf
content_types_path = '[Content_Types].xml'
content_types_xml = all_entries[content_types_path][1].decode('utf-8')
if 'Extension="wmf"' not in content_types_xml:
    content_types_xml = content_types_xml.replace(
        '</Types>',
        '<Default Extension="wmf" ContentType="image/x-wmf"/></Types>'
    )

# Rewrite ZIP with patched slides + WMF image
with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
    for fname, (item, data) in all_entries.items():
        if fname == slide78_path:
            zout.writestr(item, slide78_xml.encode('utf-8'))
        elif fname == slide79_path:
            zout.writestr(item, slide79_xml.encode('utf-8'))
        elif fname == slide80_path:
            zout.writestr(item, slide80_xml.encode('utf-8'))
        elif fname == slide81_path:
            zout.writestr(item, slide81_xml.encode('utf-8'))
        elif fname == slide81_rels_path:
            zout.writestr(item, slide81_rels.encode('utf-8'))
        elif fname == content_types_path:
            zout.writestr(item, content_types_xml.encode('utf-8'))
        else:
            zout.writestr(item, data)
    # Add WMF binary as new entry
    zout.writestr('ppt/media/image_wmf1.wmf', wmf_data)

print(f"Saved {output_path} with {total_slides} slides")
